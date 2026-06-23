from __future__ import annotations

import re
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation

PLACEHOLDER_RE = re.compile(r"\{\{[^{}]+\}\}")
PENDING_FIELD_RE = re.compile(r"\[待补充[:：][^\]]+\]")


def _checked_at() -> str:
    return datetime.now(UTC).isoformat()


def _is_selected_case(value: Any) -> bool:
    return value is not None and str(value).strip() not in ("", "null", "None", "__none__")


def _check(name: str, passed: bool, severity: str, message: str) -> dict[str, Any]:
    return {"name": name, "passed": passed, "severity": severity, "message": message}


def _sample(matches: list[str], limit: int = 5) -> str:
    unique = []
    for item in matches:
        if item not in unique:
            unique.append(item)
        if len(unique) >= limit:
            break
    return "、".join(unique)


def _extract_text(prs: Presentation) -> str:
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def _finalize_report(errors: list[str], warnings: list[str], checks: list[dict[str, Any]]) -> dict[str, Any]:
    return {
        "passed": not errors,
        "severity": "error" if errors else ("warning" if warnings else "pass"),
        "errors": errors,
        "warnings": warnings,
        "checks": checks,
        "checked_at": _checked_at(),
    }


def review_project_quality(project: dict[str, Any], final_ppt_path: str | Path) -> dict[str, Any]:
    """Deterministic post-generation QA. It reports findings only and never mutates project."""
    errors: list[str] = []
    warnings: list[str] = []
    checks: list[dict[str, Any]] = []
    final_path = Path(final_ppt_path)

    exists = final_path.exists()
    checks.append(_check("final_file_exists", exists, "error", f"最终 PPTX 文件{'存在' if exists else '不存在'}：{final_path}"))
    if not exists:
        errors.append(f"最终 PPTX 文件不存在：{final_path}")
        return _finalize_report(errors, warnings, checks)

    try:
        prs = Presentation(str(final_path))
    except Exception as exc:
        checks.append(_check("pptx_openable", False, "error", f"最终 PPTX 无法打开：{exc}"))
        errors.append(f"最终 PPTX 无法被 python-pptx 打开：{exc}")
        return _finalize_report(errors, warnings, checks)

    checks.append(_check("pptx_openable", True, "error", "最终 PPTX 可被 python-pptx 打开"))
    final_text = _extract_text(prs)

    placeholders = PLACEHOLDER_RE.findall(final_text)
    no_placeholders = not placeholders
    checks.append(_check("no_template_placeholders", no_placeholders, "error", "未发现 {{...}} 占位符残留" if no_placeholders else f"发现 {len(placeholders)} 处 {{...}} 占位符残留"))
    if placeholders:
        errors.append("最终 PPTX 仍残留 {{...}} 占位符：" + _sample(placeholders))

    pending_fields = PENDING_FIELD_RE.findall(final_text)
    no_pending_fields = not pending_fields
    checks.append(_check("no_pending_fields", no_pending_fields, "warning", "未发现 [待补充：...] 字段" if no_pending_fields else f"发现 {len(pending_fields)} 处 [待补充：...] 字段"))
    if pending_fields:
        warnings.append(f"最终 PPTX 仍残留 [待补充：...] 字段：{_sample(pending_fields)}")

    confirmed_case_id = project.get("case_selection", {}).get("confirmed_case_id")
    has_case = _is_selected_case(confirmed_case_id)
    include_m3 = project.get("m3_selection", "m3_template") == "m3_template"
    required_modules = {"M1", "M2", "M6"}
    if include_m3:
        required_modules.add("M3")
    if has_case:
        required_modules.add("M5")

    modules = {module.get("module_id"): module for module in project.get("modules", [])}
    for module_id in ["M1", "M2", "M3", "M5", "M6"]:
        module = modules.get(module_id)
        if not module:
            errors.append(f"{module_id} 模块状态缺失")
            checks.append(_check(f"{module_id}_module_present", False, "error", f"{module_id} 模块状态缺失"))
            continue

        chapter_path = str(module.get("chapter_ppt_path") or "")
        status = module.get("status")
        if module_id in required_modules:
            ok = status == "rendered" and bool(chapter_path)
            checks.append(_check(f"{module_id}_chapter_rendered", ok, "error", f"{module_id} 应已渲染并写入章节路径"))
            if not ok:
                errors.append(f"{module_id} 已选择但模块状态或 chapter_ppt_path 不完整")
        else:
            ok = status == "skipped" and not chapter_path
            checks.append(_check(f"{module_id}_chapter_skipped", ok, "error", f"{module_id} 未选择时应跳过且章节路径为空"))
            if not ok:
                errors.append(f"{module_id} 未选择但模块状态或 chapter_ppt_path 与跳过逻辑不一致")

    final_name = final_path.name
    if not include_m3 and "M3" in final_name:
        checks.append(_check("filename_m3_consistency", False, "error", "跳过 M3 时最终文件名不应包含 M3"))
        errors.append("已选择跳过 M3，但最终合并文件名仍包含 M3")
    else:
        checks.append(_check("filename_m3_consistency", True, "error", "最终文件名与 M3 选择一致"))

    if not has_case and "M5" in final_name:
        checks.append(_check("filename_m5_consistency", False, "error", "未选择 M5 案例时最终文件名不应包含 M5"))
        errors.append("未选择 M5 案例，但最终合并文件名仍包含 M5")
    else:
        checks.append(_check("filename_m5_consistency", True, "error", "最终文件名与 M5 案例选择一致"))

    return _finalize_report(errors, warnings, checks)


def quality_review_failed_report(exc: Exception) -> dict[str, Any]:
    message = f"QAReviewAgent 执行异常：{exc}"
    return _finalize_report(
        [message],
        [],
        [_check("qa_review_agent", False, "error", message)],
    )
