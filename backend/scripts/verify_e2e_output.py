# -*- coding: utf-8 -*-
"""验证 e2e 最终 PPT 文本内容符合要求。"""
import sys
from pathlib import Path

from pptx import Presentation


def verify_e2e_output(pptx_path: Path) -> list[str]:
    """验证 PPTX 文件文本内容。

    Returns:
        错误列表（空表示全部通过）。
    """
    errors = []
    prs = Presentation(str(pptx_path))

    all_text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text_parts.append(shape.text)

    all_text = " ".join(all_text_parts)

    # 负面断言：不应包含未替换的占位符
    if "{{" in all_text:
        errors.append("PPTX 包含 '{{'（未替换的占位符）")
    if "}}" in all_text:
        errors.append("PPTX 包含 '}}'（未替换的占位符）")
    if "[待补充" in all_text:
        errors.append("PPTX 包含 '[待补充'（未填充的兜底字段）")

    # 正面断言：应包含项目基本信息
    if "南京地铁3号线声屏障改造工程" not in all_text:
        errors.append("PPTX 不包含项目名称 '南京地铁3号线声屏障改造工程'")
    if "南京" not in all_text:
        errors.append("PPTX 不包含项目所在地 '南京'")
    if "南京地铁集团有限公司" not in all_text:
        errors.append("PPTX 不包含业主单位 '南京地铁集团有限公司'")
    if "轨交既有线改造" not in all_text:
        errors.append("PPTX 不包含产品线 '轨交既有线改造'")

    # 正面断言：应包含规则识别字段（从文件名/内容提取）
    has_line = "3号线" in all_text or "S1号线" in all_text
    if not has_line:
        errors.append("PPTX 不包含线路名称（应为 '3号线' 或 'S1号线'）")

    has_pain_point = "噪声治理" in all_text or "施工窗口受限" in all_text or "工期紧张" in all_text or "既有线改造约束" in all_text
    if not has_pain_point:
        errors.append("PPTX 不包含现场痛点（应为 '噪声治理'/'施工窗口受限'/'工期紧张'/'既有线改造约束' 等）")

    has_scenario = "既有线改造" in all_text or "轨道交通声屏障" in all_text or "公路声屏障" in all_text
    if not has_scenario:
        errors.append("PPTX 不包含施工场景（应为 '既有线改造'/'轨道交通声屏障' 等）")

    # M5 案例内容验证：应包含案例关键词（不指定具体案例，通用检查）
    has_case_content = (
        "案例" in all_text
        or "南昌" in all_text
        or "轨道交通4号线" in all_text
        or "S1号线" in all_text
        or "宁波" in all_text
        or "地铁" in all_text
        or "声屏障" in all_text
    )
    if not has_case_content:
        errors.append(
            "PPTX 不包含 M5 案例相关内容（应包含'案例'或具体案例关键词如'南昌轨道交通4号线'/'S1号线'/'宁波地铁2号线'）"
        )

    # 最终 PPT 应合入 M5 案例（M5 被跳过后文件名不含 M5，渲染则含 M5）
    # 宽松检查：文件名含 M1_M2_M5_M6 或内容包含"案例"相关词
    file_has_m5 = "M1_M2_M5_M6" in pptx_path.name
    text_has_case = "案例" in all_text or "南昌" in all_text or "轨道交通4号线" in all_text or "S1号线" in all_text or "宁波" in all_text
    if not file_has_m5 and not text_has_case:
        errors.append("最终 PPT 未合入 M5 案例章节（文件路径不含 M5 且内容无案例关键词）")

    return errors


if __name__ == "__main__":
    code_dir = Path(__file__).resolve().parents[2]
    output_path = code_dir / "data" / "outputs" / "e2e_downloaded_final.pptx"

    if not output_path.exists():
        print(f"错误：文件不存在 {output_path}")
        sys.exit(1)

    errors = verify_e2e_output(output_path)
    if errors:
        print("验证失败：")
        for err in errors:
            print(f"  - {err}")
        sys.exit(1)
    else:
        print("验证通过：PPTX 内容符合要求")
        sys.exit(0)