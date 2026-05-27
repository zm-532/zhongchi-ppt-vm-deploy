import importlib
import os
import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation


def _png_bytes(color=(20, 120, 200)) -> bytes:
    buffer = BytesIO()
    Image.new("RGB", (640, 360), color).save(buffer, format="PNG")
    return buffer.getvalue()


class ProjectGenerationM3MaterialsTest(unittest.TestCase):
    def setUp(self):
        test_tmp_root = Path(__file__).resolve().parents[1] / "test_tmp"
        test_tmp_root.mkdir(exist_ok=True)
        self.temp_dir = tempfile.TemporaryDirectory(dir=test_tmp_root)
        os.environ["ZHONGCHI_DATA_DIR"] = self.temp_dir.name
        os.environ["ZHONGCHI_PPT_MERGE_ENGINE"] = "python-pptx"

        from fastapi.testclient import TestClient

        app_module = importlib.import_module("app.main")
        self.client = TestClient(app_module.app)

    def tearDown(self):
        self.temp_dir.cleanup()
        os.environ.pop("ZHONGCHI_DATA_DIR", None)
        os.environ.pop("ZHONGCHI_PPT_MERGE_ENGINE", None)

    def _review_project(self, project_id: int) -> None:
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": None,
                "m3_selection": "m3_template",
                "notes": "确认正式 M3 materials",
            },
        )
        self.assertEqual(response.status_code, 200, response.text)

    def test_generate_uses_saved_m3_materials_with_multi_image_expansion(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "M3正式生成项目", "project_location": "南京"},
        ).json()["project_id"]
        self._review_project(project_id)

        save_response = self.client.post(
            f"/api/projects/{project_id}/m3-materials",
            files=[
                ("descriptions", (None, "项目基本情况-1：正式第一张\n项目基本情况-2：正式第二张\n项目线路图-1：正式线路图")),
                ("files", ("项目基本情况-1.png", _png_bytes(), "image/png")),
                ("files", ("项目基本情况-2.png", _png_bytes((200, 80, 20)), "image/png")),
                ("files", ("项目线路图-1.png", _png_bytes((80, 20, 200)), "image/png")),
            ],
        )
        self.assertEqual(save_response.status_code, 200, save_response.text)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202, generate_response.text)

        detail = self.client.get(f"/api/projects/{project_id}").json()
        m3_module = next(module for module in detail["modules"] if module["module_id"] == "M3")
        m3_path = Path(m3_module["chapter_ppt_path"])
        self.assertTrue(m3_path.exists())

        m3_prs = Presentation(str(m3_path))
        self.assertEqual(len(m3_prs.slides), 10)
        m3_text = "\n".join(
            shape.text
            for slide in m3_prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        self.assertIn("正式第一张", m3_text)
        self.assertIn("正式第二张", m3_text)
        self.assertIn("正式线路图", m3_text)
        self.assertNotIn("{{m3_", m3_text)
        self.assertNotIn("{{image:m3_basic}}", m3_text)

    def test_generate_without_m3_images_still_succeeds(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "M3无图正式生成项目", "project_location": "南京"},
        ).json()["project_id"]
        self._review_project(project_id)
        save_response = self.client.post(
            f"/api/projects/{project_id}/m3-materials",
            files=[],
        )
        self.assertEqual(save_response.status_code, 200, save_response.text)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202, generate_response.text)
        detail = self.client.get(f"/api/projects/{project_id}").json()
        self.assertTrue(Path(detail["final_ppt_path"]).exists())

    def test_generate_adds_fixed_tail_by_default_without_print_tail(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "默认固定尾页项目", "project_location": "南京"},
        ).json()["project_id"]
        self._review_project(project_id)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202, generate_response.text)

        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        chapters_dir = final_path.parent / "chapters"
        self.assertFalse(detail["include_print_tail_page"])
        self.assertTrue((chapters_dir / "TAIL_FIXED_固定尾页.pptx").exists())
        self.assertFalse((chapters_dir / "TAIL_PRINT_尾页打印版.pptx").exists())
        self.assertEqual(
            len(Presentation(str(final_path)).slides),
            sum(len(Presentation(str(path)).slides) for path in chapters_dir.glob("*.pptx")),
        )

    def test_generate_adds_print_tail_after_fixed_tail_when_selected(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "打印版尾页项目", "project_location": "南京"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": None,
                "m3_selection": "m3_template",
                "include_print_tail_page": True,
                "notes": "确认添加打印版尾页",
            },
        )
        self.assertEqual(response.status_code, 200, response.text)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202, generate_response.text)

        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        chapters_dir = final_path.parent / "chapters"
        fixed_tail = chapters_dir / "TAIL_FIXED_固定尾页.pptx"
        print_tail = chapters_dir / "TAIL_PRINT_尾页打印版.pptx"
        self.assertTrue(detail["include_print_tail_page"])
        self.assertTrue(fixed_tail.exists())
        self.assertTrue(print_tail.exists())
        self.assertGreater(len(Presentation(str(print_tail)).slides), 0)
        self.assertEqual(
            len(Presentation(str(final_path)).slides),
            sum(len(Presentation(str(path)).slides) for path in chapters_dir.glob("*.pptx")),
        )


if __name__ == "__main__":
    unittest.main()
