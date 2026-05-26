import importlib
import os
import tempfile
import unittest
import zipfile
from io import BytesIO
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from app.quality_review import review_project_quality


def _zip_bytes(entries: dict[str, str]) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, text in entries.items():
            archive.writestr(name, text)
    return buffer.getvalue()


def _docx_bytes(text: str) -> bytes:
    return _zip_bytes({"word/document.xml": f"<w:document><w:body><w:t>{text}</w:t></w:body></w:document>"})


def _pptx_bytes(text: str) -> bytes:
    return _zip_bytes({"ppt/slides/slide1.xml": f"<p:sld><a:t>{text}</a:t></p:sld>"})


def _xlsx_bytes(text: str) -> bytes:
    return _zip_bytes({"xl/sharedStrings.xml": f"<sst><si><t>{text}</t></si></sst>"})


class BackendApiTest(unittest.TestCase):
    def setUp(self):
        test_tmp_root = Path(__file__).resolve().parents[1] / "test_tmp"
        test_tmp_root.mkdir(exist_ok=True)
        self.temp_dir = tempfile.TemporaryDirectory(dir=test_tmp_root)
        os.environ["ZHONGCHI_DATA_DIR"] = self.temp_dir.name
        self._old_merge_engine = os.environ.get("ZHONGCHI_PPT_MERGE_ENGINE")
        os.environ["ZHONGCHI_PPT_MERGE_ENGINE"] = "python-pptx"
        self._old_llm_env = {
            "ZHONGCHI_LLM_BASE_URL": os.environ.get("ZHONGCHI_LLM_BASE_URL"),
            "ZHONGCHI_LLM_API_KEY": os.environ.get("ZHONGCHI_LLM_API_KEY"),
            "ZHONGCHI_LLM_MODEL": os.environ.get("ZHONGCHI_LLM_MODEL"),
            "ZHONGCHI_M1M2_LLM_ENABLED": os.environ.get("ZHONGCHI_M1M2_LLM_ENABLED"),
        }
        for name in self._old_llm_env:
            os.environ.pop(name, None)
        os.environ["ZHONGCHI_M1M2_LLM_ENABLED"] = "0"

        from fastapi.testclient import TestClient

        app_module = importlib.import_module("app.main")
        self.client = TestClient(app_module.app)

    def tearDown(self):
        self.temp_dir.cleanup()
        os.environ.pop("ZHONGCHI_DATA_DIR", None)
        if self._old_merge_engine is None:
            os.environ.pop("ZHONGCHI_PPT_MERGE_ENGINE", None)
        else:
            os.environ["ZHONGCHI_PPT_MERGE_ENGINE"] = self._old_merge_engine
        for name, old_value in self._old_llm_env.items():
            if old_value is None:
                os.environ.pop(name, None)
            else:
                os.environ[name] = old_value

    def test_project_create_list_and_detail(self):
        response = self.client.post(
            "/api/projects",
            json={
                "project_name": "某城市轨道交通声屏障改造项目",
                "project_location": "南京",
                "owner_unit": "某建设单位",
                "product_line": "轨交既有线改造",
            },
        )

        self.assertEqual(response.status_code, 201)
        project = response.json()
        self.assertEqual(project["project_name"], "某城市轨道交通声屏障改造项目")
        self.assertEqual(project["task_status"], "待生成")
        self.assertEqual([module["module_id"] for module in project["modules"]], ["M1", "M2", "M3", "M5", "M6"])

        list_response = self.client.get("/api/projects")
        self.assertEqual(list_response.status_code, 200)
        self.assertEqual(len(list_response.json()), 1)

        detail_response = self.client.get(f"/api/projects/{project['project_id']}")
        self.assertEqual(detail_response.status_code, 200)
        self.assertEqual(detail_response.json()["project_id"], project["project_id"])

    def test_delete_project_removes_it_from_list_and_detail(self):
        first_id = self.client.post("/api/projects", json={"project_name": "待删除项目"}).json()["project_id"]
        second_id = self.client.post("/api/projects", json={"project_name": "保留项目"}).json()["project_id"]

        delete_response = self.client.delete(f"/api/projects/{first_id}")

        self.assertEqual(delete_response.status_code, 204)
        list_response = self.client.get("/api/projects")
        self.assertEqual([item["project_id"] for item in list_response.json()], [second_id])
        detail_response = self.client.get(f"/api/projects/{first_id}")
        self.assertEqual(detail_response.status_code, 404)

    def test_dev_llm_test_uses_configured_url_without_appending_path(self):
        class FakeResponse:
            status_code = 200

            def raise_for_status(self):
                return None

            def json(self):
                return {
                    "model": "fake-model",
                    "choices": [{"message": {"content": "LLM连接成功"}}],
                }

        calls = []

        class FakeClient:
            def __init__(self, *args, **kwargs):
                self.kwargs = kwargs

            def __enter__(self):
                return self

            def __exit__(self, exc_type, exc, tb):
                return False

            def post(self, url, headers=None, json=None):
                calls.append({"url": url, "headers": headers, "json": json, "kwargs": self.kwargs})
                return FakeResponse()

        old_base = os.environ.get("ZHONGCHI_LLM_BASE_URL")
        old_key = os.environ.get("ZHONGCHI_LLM_API_KEY")
        old_model = os.environ.get("ZHONGCHI_LLM_MODEL")
        os.environ["ZHONGCHI_LLM_BASE_URL"] = "https://example.test/custom/chat"
        os.environ["ZHONGCHI_LLM_API_KEY"] = "test-key"
        os.environ["ZHONGCHI_LLM_MODEL"] = "test-model"
        try:
            importlib.import_module("app.llm")
            with patch("app.llm.httpx.Client", FakeClient):
                response = self.client.post("/api/dev/llm-test", json={"prompt": "测试提示词"})
        finally:
            for name, old_value in (
                ("ZHONGCHI_LLM_BASE_URL", old_base),
                ("ZHONGCHI_LLM_API_KEY", old_key),
                ("ZHONGCHI_LLM_MODEL", old_model),
            ):
                if old_value is None:
                    os.environ.pop(name, None)
                else:
                    os.environ[name] = old_value

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["reply"], "LLM连接成功")
        self.assertEqual(calls[0]["url"], "https://example.test/custom/chat")
        self.assertEqual(calls[0]["json"]["model"], "test-model")
        self.assertEqual(calls[0]["json"]["messages"][-1]["content"], "测试提示词")

    def test_upload_records_file_against_allowed_module(self):
        project_id = self.client.post("/api/projects", json={"project_name": "上传测试项目"}).json()["project_id"]

        response = self.client.post(
            f"/api/projects/{project_id}/modules/M1/files",
            files={"file": ("m1.pdf", b"M1 material", "application/pdf")},
        )

        self.assertEqual(response.status_code, 201)
        payload = response.json()
        self.assertEqual(payload["module_id"], "M1")
        self.assertEqual(payload["filename"], "m1.pdf")

        detail = self.client.get(f"/api/projects/{project_id}").json()
        m1 = next(module for module in detail["modules"] if module["module_id"] == "M1")
        self.assertEqual(m1["status"], "uploaded")
        self.assertEqual(m1["uploaded_file_ids"], [payload["file_id"]])

    def test_unified_upload_accepts_multiple_files_without_module_id(self):
        project_id = self.client.post("/api/projects", json={"project_name": "统一上传测试项目"}).json()["project_id"]

        response = self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("项目简介.pdf", b"metro project brief", "application/pdf")),
                ("files", ("案例材料.docx", b"case material", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
            ],
        )

        self.assertEqual(response.status_code, 201)
        payload = response.json()
        self.assertEqual(len(payload), 2)
        self.assertTrue(all(item["module_id"] is None for item in payload))
        self.assertTrue(all(item["assigned_modules"] == [] for item in payload))
        self.assertTrue(all(item["parse_status"] == "pending" for item in payload))

    def test_upload_rejects_disallowed_module(self):
        project_id = self.client.post("/api/projects", json={"project_name": "模块校验项目"}).json()["project_id"]

        response = self.client.post(
            f"/api/projects/{project_id}/modules/M4/files",
            files={"file": ("m4.pdf", b"M4 material", "application/pdf")},
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("module_id", response.json()["detail"])

    def test_analyze_classification_review_and_generate_use_confirmed_selection(self):
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "南京地铁声屏障项目",
                "project_location": "南京",
                "owner_unit": "某建设单位",
                "product_line": "轨道交通声屏障",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf")),
                ("files", ("南昌轨道交通案例.docx", b"case rail transit", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
            ],
        )

        analyze_response = self.client.post(f"/api/projects/{project_id}/analyze")
        self.assertEqual(analyze_response.status_code, 200)
        classification = analyze_response.json()
        self.assertEqual(classification["detected_project_type"], "metro")
        self.assertEqual(classification["classification_status"], "analyzed")
        self.assertEqual(classification["template_selection"]["M1_M2"]["template_key"], "metro")
        self.assertEqual(classification["template_selection"]["M1_M2"]["template_filename"], "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx")
        self.assertNotIn("template_path", classification["template_selection"]["M1_M2"])
        self.assertEqual(classification["template_selection"]["M5"]["template_filename"], "M5示例.pptx")
        self.assertEqual(classification["template_selection"]["M6"]["template_key"], "enterprise")
        self.assertIn("recommended_cases", classification["case_selection"])

        get_response = self.client.get(f"/api/projects/{project_id}/classification")
        self.assertEqual(get_response.status_code, 200)
        self.assertEqual(get_response.json()["detected_project_type"], "metro")

        review_response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": 1,
                "notes": "确认地铁模板和推荐案例",
            },
        )
        self.assertEqual(review_response.status_code, 200)
        reviewed = review_response.json()
        self.assertEqual(reviewed["classification_status"], "reviewed")
        self.assertEqual(reviewed["confirmed_project_type"], "metro")
        self.assertEqual(reviewed["case_selection"]["confirmed_case_id"], 1)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202)
        generated = generate_response.json()
        self.assertEqual(generated["task_status"], "完成")
        self.assertEqual([module["module_id"] for module in generated["modules"]], ["M1", "M2", "M3", "M5", "M6"])
        self.assertIn("quality_report", generated)
        self.assertTrue(generated["quality_report"]["passed"])

        detail = self.client.get(f"/api/projects/{project_id}").json()
        self.assertEqual(detail["confirmed_project_type"], "metro")
        self.assertTrue(Path(detail["final_ppt_path"]).exists())
        self.assertTrue(detail["quality_report"]["passed"])
        self.assertIn(detail["quality_report"]["severity"], ["pass", "warning"])
        task_detail = self.client.get(f"/api/projects/{project_id}/task").json()
        self.assertEqual(task_detail["quality_report"]["checked_at"], detail["quality_report"]["checked_at"])
        final_text = " ".join(
            shape.text for slide in Presentation(detail["final_ppt_path"]).slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        self.assertIn("南京地铁声屏障项目", final_text)
        self.assertNotIn("{{m3_", final_text)
        self.assertNotIn("工程量与施工周期测算", final_text)

    def test_quality_review_reports_placeholders_and_pending_fields(self):
        with tempfile.TemporaryDirectory(dir=Path(__file__).resolve().parents[1] / "test_tmp") as tmp:
            pptx_path = Path(tmp) / "demo_M1_M2_M3_M6.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            textbox = slide.shapes.add_textbox(0, 0, 5000000, 1000000)
            textbox.text = "项目 {{project_name}} [待补充：项目关键信息]"
            prs.save(pptx_path)

            project = {
                "case_selection": {"confirmed_case_id": None},
                "m3_selection": "m3_template",
                "modules": [
                    {"module_id": "M1", "status": "rendered", "chapter_ppt_path": str(pptx_path)},
                    {"module_id": "M2", "status": "rendered", "chapter_ppt_path": str(pptx_path)},
                    {"module_id": "M3", "status": "rendered", "chapter_ppt_path": str(pptx_path)},
                    {"module_id": "M5", "status": "skipped", "chapter_ppt_path": ""},
                    {"module_id": "M6", "status": "rendered", "chapter_ppt_path": str(pptx_path)},
                ],
            }

            report = review_project_quality(project, pptx_path)

        self.assertFalse(report["passed"])
        self.assertEqual(report["severity"], "error")
        self.assertTrue(any("{{...}}" in item for item in report["errors"]))
        self.assertTrue(any("[待补充" in item for item in report["warnings"]))

    def test_generate_still_completes_when_quality_review_fails(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "高速公路声屏障项目", "product_line": "公路声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("高速公路项目简介.pdf", b"highway noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": classification["detected_project_type"],
                "template_selection": classification["template_selection"],
                "confirmed_case_id": None,
                "notes": "确认无案例",
            },
        )

        with patch("app.storage.review_project_quality", side_effect=RuntimeError("qa down")):
            generate_response = self.client.post(f"/api/projects/{project_id}/generate")

        self.assertEqual(generate_response.status_code, 202)
        generated = generate_response.json()
        self.assertEqual(generated["task_status"], "完成")
        self.assertFalse(generated["quality_report"]["passed"])
        self.assertEqual(generated["quality_report"]["severity"], "error")
        self.assertIn("qa down", " ".join(generated["quality_report"]["errors"]))

    def test_analyze_project_uses_llm_classifier_when_available(self):
        os.environ["ZHONGCHI_M1M2_LLM_ENABLED"] = "1"
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "普通高速公路项目", "product_line": "公路声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("高速项目简介.pdf", b"highway noise barrier", "application/pdf"))],
        )

        llm_result = {
            "project_type": "railway",
            "confidence": 0.88,
            "matched_keywords": ["铁路", "声屏障"],
            "evidence": [
                {
                    "project_type": "railway",
                    "keyword": "铁路",
                    "source": "LLM",
                    "snippet": "资料描述为铁路声屏障项目",
                }
            ],
            "reasoning_summary": "模型根据模板画像判断该项目更接近铁路声屏障模板。",
        }

        with patch("app.m1m2_classifier.call_llm_project_classifier", return_value=llm_result):
            response = self.client.post(f"/api/projects/{project_id}/analyze")

        self.assertEqual(response.status_code, 200)
        classification = response.json()
        self.assertEqual(classification["detected_project_type"], "railway")
        self.assertEqual(classification["template_selection"]["M1_M2"]["template_key"], "railway")
        self.assertEqual(classification["classification_method"], "llm")
        self.assertEqual(classification["llm_reasoning_summary"], "模型根据模板画像判断该项目更接近铁路声屏障模板。")
        self.assertEqual(classification["fallback_reason"], "")

    def test_analyze_project_falls_back_to_rules_when_llm_fails(self):
        os.environ["ZHONGCHI_M1M2_LLM_ENABLED"] = "1"
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "南京地铁声屏障项目", "product_line": "轨道交通声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )

        with patch("app.m1m2_classifier.call_llm_project_classifier", side_effect=RuntimeError("llm down")):
            response = self.client.post(f"/api/projects/{project_id}/analyze")

        self.assertEqual(response.status_code, 200)
        classification = response.json()
        self.assertEqual(classification["detected_project_type"], "metro")
        self.assertEqual(classification["template_selection"]["M1_M2"]["template_key"], "metro")
        self.assertEqual(classification["classification_method"], "rule_fallback")
        self.assertIn("llm down", classification["fallback_reason"])

    def test_analyze_project_falls_back_to_rules_when_llm_returns_invalid_project_type(self):
        os.environ["ZHONGCHI_M1M2_LLM_ENABLED"] = "1"
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "高速公路声屏障项目", "product_line": "公路声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("高速公路项目简介.pdf", b"highway noise barrier", "application/pdf"))],
        )

        with patch(
            "app.m1m2_classifier.call_llm_project_classifier",
            return_value={"project_type": "airport", "confidence": 0.92, "matched_keywords": ["机场"]},
        ):
            response = self.client.post(f"/api/projects/{project_id}/analyze")

        self.assertEqual(response.status_code, 200)
        classification = response.json()
        self.assertEqual(classification["detected_project_type"], "highway")
        self.assertEqual(classification["template_selection"]["M1_M2"]["template_key"], "highway")
        self.assertEqual(classification["classification_method"], "rule_fallback")
        self.assertIn("非法 project_type", classification["fallback_reason"])

    def test_generate_omits_m5_when_review_confirms_no_case(self):
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "不选择案例项目",
                "project_location": "南京",
                "product_line": "轨道交通声屏障",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        review_response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": None,
                "notes": "本次不选择 M5 案例",
            },
        )
        self.assertEqual(review_response.status_code, 200)
        self.assertIsNone(review_response.json()["case_selection"]["confirmed_case_id"])

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202)
        generated = generate_response.json()
        m5_module = next(module for module in generated["modules"] if module["module_id"] == "M5")
        self.assertEqual(m5_module["status"], "skipped")
        self.assertEqual(m5_module["chapter_ppt_path"], "")

        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        self.assertTrue(final_path.exists())
        self.assertIn("M1_M2_M3_M6", final_path.name)
        self.assertNotIn("M5", final_path.name)
        self.assertFalse((final_path.parent / "chapters" / "M5_同类型案例匹配.pptx").exists())
        self.assertTrue((final_path.parent / "chapters" / "M3_项目深化方案.pptx").exists())

    def test_generate_skips_m3_when_user_selects_m3_skip(self):
        """验证用户选择'暂不选择'M3时，最终PPT不包含M3，文件名为M1_M2_M6或M1_M2_M5_M6。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "SkipM3Project",
                "project_location": "Nanjing",
                "product_line": "Rail Transit",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("NanjingMetro.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        review_response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": None,
                "m3_selection": "m3_skip",
                "notes": "Skip M3 this time",
            },
        )
        self.assertEqual(review_response.status_code, 200)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202)
        generated = generate_response.json()
        # M3 should be skipped
        m3_module = next(module for module in generated["modules"] if module["module_id"] == "M3")
        self.assertEqual(m3_module["status"], "skipped")
        self.assertEqual(m3_module["chapter_ppt_path"], "")
        # M1/M2/M6 rendered normally
        for mid in ("M1", "M2", "M6"):
            module = next(m for m in generated["modules"] if m["module_id"] == mid)
            self.assertEqual(module["status"], "rendered")
            self.assertTrue(Path(module["chapter_ppt_path"]).exists())

        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        self.assertTrue(final_path.exists())
        # Final filename should not contain M3 in merge order
        self.assertNotIn("M3_M6", final_path.name)
        self.assertIn("M1_M2_M6", final_path.name)
        # M3 chapter file should not exist
        self.assertFalse((final_path.parent / "chapters" / "M3_项目深化方案.pptx").exists())
        # M1/M2 and M6 chapters should exist
        self.assertTrue((final_path.parent / "chapters" / "M1_M2_行业背景与技术标准.pptx").exists())
        self.assertTrue((final_path.parent / "chapters" / "M6_企业背书与荣誉.pptx").exists())

    def test_generate_skips_m3_with_case_selection(self):
        """验证用户选择'暂不选择'M3且同时选择M5案例时，最终文件名为M1_M2_M5_M6。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "SkipM3WithCaseProject",
                "project_location": "Nanjing",
                "product_line": "Rail Transit",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("NanjingMetro.pdf", b"metro line noise barrier metro", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        cases = classification.get("case_selection", {}).get("recommended_cases", [])
        first_case_id = cases[0]["case_id"] if cases else None

        review_response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": first_case_id,
                "m3_selection": "m3_skip",
                "notes": "No M3, but use M5 case",
            },
        )
        self.assertEqual(review_response.status_code, 200)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202)
        generated = generate_response.json()
        # M3 skipped
        m3_module = next(module for module in generated["modules"] if module["module_id"] == "M3")
        self.assertEqual(m3_module["status"], "skipped")
        # M5 rendered normally
        m5_module = next(module for module in generated["modules"] if module["module_id"] == "M5")
        self.assertEqual(m5_module["status"], "rendered")
        self.assertTrue(Path(m5_module["chapter_ppt_path"]).exists())

        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        self.assertTrue(final_path.exists())
        self.assertNotIn("M3_M6", final_path.name)
        self.assertIn("M1_M2_M5_M6", final_path.name)
        self.assertFalse((final_path.parent / "chapters" / "M3_项目深化方案.pptx").exists())

    def test_generate_keeps_string_case_id_and_renders_m5(self):
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "字符串案例项目",
                "project_location": "南京",
                "product_line": "轨道交通声屏障",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        # Use a real recommended case ID (string format from case_matcher)
        recommended = classification.get("case_selection", {}).get("recommended_cases", [])
        if not recommended:
            self.skipTest("无推荐案例，跳过字符串 case_id 测试")
        string_case_id = recommended[0]["case_id"]
        self.assertIsInstance(string_case_id, str)

        review_response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": string_case_id,
                "notes": "选择真实案例库字符串 ID",
            },
        )
        self.assertEqual(review_response.status_code, 200)
        self.assertEqual(review_response.json()["case_selection"]["confirmed_case_id"], string_case_id)

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202)
        generated = generate_response.json()
        m5_module = next(module for module in generated["modules"] if module["module_id"] == "M5")
        self.assertEqual(m5_module["status"], "rendered")
        self.assertTrue(m5_module["chapter_ppt_path"])

        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        self.assertTrue(final_path.exists())
        self.assertIn("M1_M2_M3_M5_M6", final_path.name)

    def test_project_type_detection_distinguishes_four_fixed_m1_m2_templates(self):
        cases = [
            (
                {
                    "project_name": "沪宁高速公路声屏障改造项目",
                    "product_line": "公路声屏障",
                },
                [("files", ("高速公路降噪需求.pdf", "高速 公路 全封闭 声屏障 改造".encode(), "application/pdf"))],
                "highway",
                "公路全封闭声屏障（M1_&_M2）.pptx",
            ),
            (
                {
                    "project_name": "京沪铁路声屏障新建工程",
                    "product_line": "铁路声屏障",
                },
                [("files", ("铁路技术标准.pdf", "铁路 干线 声屏障 技术标准".encode(), "application/pdf"))],
                "railway",
                "铁路声屏障行业背景与技术发展（M1_&_M2）.pptx",
            ),
            (
                {
                    "project_name": "南京地铁三号线声屏障项目",
                    "product_line": "地铁声屏障",
                },
                [("files", ("地铁项目简介.pdf", "地铁 车站 区间 声屏障".encode(), "application/pdf"))],
                "metro",
                "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx",
            ),
            (
                {
                    "project_name": "既有轨道交通线路声屏障改造工程",
                    "product_line": "轨交既有线改造",
                },
                [("files", ("既有线夜间施工窗口.pdf", "既有线 轨道交通 改造 夜间 施工窗口".encode(), "application/pdf"))],
                "existing_rail_transit",
                "铁路_&_轨道交通既有线声屏障_（M1_&_M2）.pptx",
            ),
        ]

        for project_payload, files, expected_type, expected_template in cases:
            with self.subTest(expected_type=expected_type):
                project_id = self.client.post("/api/projects", json=project_payload).json()["project_id"]
                self.client.post(f"/api/projects/{project_id}/files", files=files)

                classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

                self.assertEqual(classification["detected_project_type"], expected_type)
                self.assertEqual(classification["template_selection"]["M1_M2"]["template_key"], expected_type)
                self.assertEqual(classification["template_selection"]["M1_M2"]["template_filename"], expected_template)
                self.assertGreater(classification["confidence"], 0.5)
                self.assertTrue(classification["matched_keywords"])

    def test_real_pdf_extracts_rail_transit_keyword_and_returns_evidence(self):
        source_pdf = Path(r"D:\中驰股份\06中海寰宇\20.噪声污染意见.pdf")
        if not source_pdf.exists():
            self.skipTest("真实 PDF 样例不存在")
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "M1/M2真实PDF识别测试", "product_line": ""},
        ).json()["project_id"]

        with source_pdf.open("rb") as file:
            upload_response = self.client.post(
                f"/api/projects/{project_id}/files",
                files=[("files", (source_pdf.name, file.read(), "application/pdf"))],
            )
        self.assertEqual(upload_response.status_code, 201)

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        self.assertEqual(classification["detected_project_type"], "metro")
        self.assertIn("轨道交通", classification["matched_keywords"])
        self.assertGreater(classification["confidence"], 0.5)
        self.assertTrue(any("轨道交通" in item["snippet"] for item in classification["detection_evidence"]))
        self.assertEqual(classification["template_selection"]["M1_M2"]["template_key"], "metro")

    def test_review_rejects_project_type_outside_fixed_enum(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "南京地铁声屏障项目", "product_line": "地铁声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        self.client.post(f"/api/projects/{project_id}/analyze")

        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={"confirmed_project_type": "airport"},
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("project_type", response.json()["detail"])

    def test_review_override_reselects_m1_m2_template_when_selection_is_not_supplied(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "南京地铁声屏障项目", "product_line": "地铁声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        self.client.post(f"/api/projects/{project_id}/analyze")

        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={"confirmed_project_type": "railway"},
        )

        self.assertEqual(response.status_code, 200)
        reviewed = response.json()
        self.assertEqual(reviewed["confirmed_project_type"], "railway")
        self.assertEqual(reviewed["template_selection"]["M1_M2"]["template_key"], "railway")
        self.assertEqual(reviewed["template_selection"]["M1_M2"]["template_filename"], "铁路声屏障行业背景与技术发展（M1_&_M2）.pptx")

    def test_review_override_keeps_m1_m2_template_consistent_with_confirmed_type(self):
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "南京地铁声屏障项目", "product_line": "地铁声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "highway",
                "template_selection": classification["template_selection"],
            },
        )

        self.assertEqual(response.status_code, 200)
        reviewed = response.json()
        self.assertEqual(reviewed["confirmed_project_type"], "highway")
        self.assertEqual(reviewed["template_selection"]["M1_M2"]["template_key"], "highway")
        self.assertEqual(reviewed["template_selection"]["M1_M2"]["template_filename"], "公路全封闭声屏障（M1_&_M2）.pptx")

    def test_analyze_extracts_office_text_and_classifies_document_roles(self):
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "北京高速声屏障项目",
                "project_location": "北京",
                "product_line": "公路声屏障",
            },
        ).json()["project_id"]

        upload_response = self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("技术标书.docx", _docx_bytes("招标文件 技术标准 高速 公路 声屏障"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
                ("files", ("现场调研表.xlsx", _xlsx_bytes("现场调研 痛点 施工窗口 降噪需求"), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")),
                ("files", ("企业资质介绍.pptx", _pptx_bytes("企业介绍 CNAS 专利 荣誉 产能"), "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
                ("files", ("总平面图.png", b"\x89PNG\r\n\x1a\n", "image/png")),
            ],
        )
        self.assertEqual(upload_response.status_code, 201)

        analyze_response = self.client.post(f"/api/projects/{project_id}/analyze")
        self.assertEqual(analyze_response.status_code, 200)
        files_by_name = {item["filename"]: item for item in analyze_response.json()["files"]}

        self.assertEqual(files_by_name["技术标书.docx"]["document_role"], "tender")
        self.assertEqual(files_by_name["技术标书.docx"]["assigned_modules"], ["M1", "M2", "M5"])
        self.assertEqual(files_by_name["技术标书.docx"]["parse_status"], "parsed")
        self.assertTrue(Path(files_by_name["技术标书.docx"]["parsed_text_path"]).exists())

        self.assertEqual(files_by_name["现场调研表.xlsx"]["document_role"], "survey")
        self.assertEqual(files_by_name["现场调研表.xlsx"]["assigned_modules"], ["M2", "M5"])

        self.assertEqual(files_by_name["企业资质介绍.pptx"]["document_role"], "enterprise_material")
        self.assertEqual(files_by_name["企业资质介绍.pptx"]["assigned_modules"], ["M6"])

        self.assertEqual(files_by_name["总平面图.png"]["document_role"], "drawing")
        self.assertEqual(files_by_name["总平面图.png"]["parse_status"], "pending_enhancement")

    def test_task_generate_and_review_flow_uses_mock_statuses(self):
        project_id = self.client.post("/api/projects", json={"project_name": "生成状态项目"}).json()["project_id"]

        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202)
        self.assertEqual(generate_response.json()["task_status"], "待确认")

        task_response = self.client.get(f"/api/projects/{project_id}/task")
        self.assertEqual(task_response.status_code, 200)
        task = task_response.json()
        self.assertEqual(task["task_status"], "待确认")
        self.assertEqual({module["status"] for module in task["modules"]}, {"outlined"})
        self.assertEqual(
            task["status_history"],
            ["待生成", "模块解析中", "素材匹配中", "章节生成中", "待确认"],
        )
        for module in task["modules"]:
            self.assertEqual(module["status_history"], ["pending", "uploaded", "parsed", "matched", "outlined"])

        review_response = self.client.post(
            f"/api/projects/{project_id}/review",
            json={"approved": True, "notes": "确认使用 Mock 大纲"},
        )
        self.assertEqual(review_response.status_code, 200)
        self.assertEqual(review_response.json()["task_status"], "完成")
        self.assertEqual(
            review_response.json()["status_history"],
            ["待生成", "模块解析中", "素材匹配中", "章节生成中", "待确认", "章节渲染中", "合并中", "完成"],
        )

    def test_m5_case_match_and_render_in_final_ppt(self):
        """验证 M5 案例能稳定匹配并在最终 PPT 中合入。

        流程：创建项目 → 上传含轨道交通/声屏障/既有线改造关键词的资料
        → analyze 应返回推荐案例 → review 传入推荐案例
        → generate 后 M5 状态为 rendered → 最终 PPT 包含案例内容
        """
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "南京地铁3号线声屏障改造工程",
                "project_location": "南京",
                "owner_unit": "南京地铁集团有限公司",
                "product_line": "轨交既有线改造",
            },
        ).json()["project_id"]

        # 上传包含案例匹配关键词的资料
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                (
                    "files",
                    (
                        "南京地铁3号线声屏障改造工程施工组织设计.txt",
                        (
                            "南京地铁3号线声屏障改造工程施工组织设计\n"
                            "既有线运营期间施工，施工窗口受限，夜间天窗点作业。\n"
                            "沿线噪声问题突出，需要采取降噪措施。\n"
                            "地铁轨道交通声屏障加装工程，全线约15公里既有线改造。"
                        ).encode("utf-8"),
                        "text/plain",
                    ),
                ),
            ],
        )

        # 分析项目，应返回推荐案例
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        self.assertEqual(classification["classification_status"], "analyzed")
        self.assertIn(classification["detected_project_type"], {"metro", "existing_rail_transit"})

        recommended_cases = classification.get("case_selection", {}).get("recommended_cases", [])
        self.assertGreater(
            len(recommended_cases), 0,
            "analyze 应返回至少一个推荐案例（演示样例含轨道交通/声屏障/既有线改造关键词）",
        )

        # 取第一个推荐案例
        selected_case = recommended_cases[0]
        confirmed_case_id = selected_case["case_id"]
        self.assertIsNotNone(confirmed_case_id)
        self.assertIsInstance(confirmed_case_id, str)
        self.assertTrue(len(confirmed_case_id) > 0)

        # 人工确认，传入推荐案例
        review_response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": classification["detected_project_type"],
                "template_selection": classification["template_selection"],
                "confirmed_case_id": confirmed_case_id,
                "notes": "M5 案例匹配演示",
            },
        )
        self.assertEqual(review_response.status_code, 200)
        reviewed = review_response.json()
        self.assertEqual(reviewed["classification_status"], "reviewed")
        self.assertEqual(reviewed["case_selection"]["confirmed_case_id"], confirmed_case_id)

        # 生成 PPT
        generate_response = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_response.status_code, 202)
        generated = generate_response.json()
        self.assertEqual(generated["task_status"], "完成")

        # M5 应为 rendered 状态
        m5_module = next(m for m in generated["modules"] if m["module_id"] == "M5")
        self.assertEqual(m5_module["status"], "rendered", "M5 案例章节应被渲染（不是 skipped）")
        self.assertTrue(m5_module["chapter_ppt_path"], "M5 章节路径不应为空")

        # 最终 PPT 可下载，且文件名含 M5
        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        self.assertTrue(final_path.exists(), "最终 PPT 应已生成")
        self.assertIn("M1_M2_M3_M5_M6", final_path.name, "文件名应包含 M3/M5 表示章节已合入")

        # 下载并验证内容
        download_response = self.client.get(f"/api/projects/{project_id}/download")
        self.assertEqual(download_response.status_code, 200)
        self.assertTrue(download_response.content.startswith(b"PK"))

        # 验证 PPT 包含案例相关内容（关键词来自推荐案例）
        from io import BytesIO
        from pptx import Presentation
        prs = Presentation(BytesIO(download_response.content))
        all_text = " ".join(
            shape.text for slide in prs.slides for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        # 至少应包含声屏障或地铁等案例领域关键词
        self.assertTrue(
            any(kw in all_text for kw in ["声屏障", "地铁", "案例", "轨道交通", "既有线"]),
            "最终 PPT 应包含 M5 案例相关内容",
        )

    def test_m1_m2_placeholder_replacement_in_final_ppt(self):
        """验证主流程最终 PPT 不含 {{...}} 占位符，所有字段均被正确替换。

        回归测试：防止 template copy + replace 链路断裂导致占位符泄漏。
        """
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "南京地铁3号线既有线声屏障改造工程",
                "project_location": "南京",
                "owner_unit": "南京地铁集团有限公司",
                "product_line": "轨交既有线改造",
            },
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                (
                    "files",
                    (
                        "南京地铁3号线声屏障改造工程施工组织设计.txt",
                        (
                            "南京地铁3号线声屏障改造工程施工组织设计\n"
                            "既有线运营期间施工，施工窗口受限，夜间天窗点作业。\n"
                            "沿线噪声问题突出，需要采取降噪措施。\n"
                            "地铁轨道交通声屏障加装工程，全线约15公里既有线改造。\n"
                            "工期紧张，需在短窗口内完成噪声治理施工。"
                        ).encode("utf-8"),
                        "text/plain",
                    ),
                ),
            ],
        )

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        recommended_cases = classification.get("case_selection", {}).get("recommended_cases", [])
        confirmed_case_id = recommended_cases[0]["case_id"] if recommended_cases else None

        self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": classification["detected_project_type"],
                "template_selection": classification["template_selection"],
                "confirmed_case_id": confirmed_case_id,
                "notes": "占位符回归测试",
            },
        )

        self.client.post(f"/api/projects/{project_id}/generate")

        detail = self.client.get(f"/api/projects/{project_id}").json()
        download = self.client.get(f"/api/projects/{project_id}/download")
        self.assertEqual(download.status_code, 200)

        from io import BytesIO
        from pptx import Presentation
        prs = Presentation(BytesIO(download.content))
        all_text = " ".join(
            shape.text for slide in prs.slides for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )

        # 负面断言：不应包含任何 {{ 或 }}
        self.assertNotIn("{{", all_text, "PPTX 不应包含 '{{' 占位符")
        self.assertNotIn("}}", all_text, "PPTX 不应包含 '}}' 占位符")
        self.assertNotIn("{{project_name}}", all_text)

        # 正面断言：项目基础字段应被替换
        self.assertIn("南京地铁3号线既有线声屏障改造工程", all_text)
        self.assertIn("南京", all_text)
        self.assertIn("南京地铁集团有限公司", all_text)
        self.assertIn("轨交既有线改造", all_text)

        # 规则识别字段应被替换或兜底
        self.assertIn("3号线", all_text)
        self.assertTrue(
            any(kw in all_text for kw in ["噪声治理", "施工窗口受限", "工期紧张", "既有线改造约束"]),
            "现场痛点字段应被提取或兜底",
        )
        self.assertTrue(
            any(kw in all_text for kw in ["既有线改造", "轨道交通声屏障", "公路声屏障"]),
            "施工场景字段应被提取或兜底",
        )

    def test_update_project_basic_info(self):
        """验证 PATCH /api/projects/{id} 可部分更新项目基础信息。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "原始项目名",
                "project_location": "原始城市",
                "owner_unit": "原始单位",
                "product_line": "原始产品线",
            },
        ).json()["project_id"]

        # 更新部分字段
        updated = self.client.patch(
            f"/api/projects/{project_id}",
            json={
                "project_name": "更新后的项目名",
                "project_location": "新城市",
            },
        ).json()
        self.assertEqual(updated["project_name"], "更新后的项目名")
        self.assertEqual(updated["project_location"], "新城市")
        self.assertEqual(updated["owner_unit"], "原始单位")  # 未传值，保留原值
        self.assertEqual(updated["product_line"], "原始产品线")  # 未传值，保留原值

        # 再次更新，仅改 product_line
        updated2 = self.client.patch(
            f"/api/projects/{project_id}",
            json={"product_line": "新公路声屏障"},
        ).json()
        self.assertEqual(updated2["project_name"], "更新后的项目名")  # 保持不变
        self.assertEqual(updated2["product_line"], "新公路声屏障")

        # project_name 传空字符串应被拒绝（保留原值）
        updated3 = self.client.patch(
            f"/api/projects/{project_id}",
            json={"project_name": "  "},
        ).json()
        self.assertEqual(updated3["project_name"], "更新后的项目名")  # 空字符串不覆盖

        # 不存在的项目返回 404
        not_found = self.client.patch("/api/projects/99999", json={"project_location": "某地"})
        self.assertEqual(not_found.status_code, 404)

        # 验证项目列表中的数据也已更新
        listed = self.client.get("/api/projects").json()
        found = next(p for p in listed if p["project_id"] == project_id)
        self.assertEqual(found["project_name"], "更新后的项目名")
        self.assertEqual(found["product_line"], "新公路声屏障")

    def test_end_to_end_generate_review_and_download_final_pptx(self):
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "端到端验收项目",
                "project_location": "南京",
                "owner_unit": "某建设单位",
                "product_line": "轨交既有线改造",
            },
        ).json()["project_id"]

        for module_id in ["M1", "M2", "M3", "M5", "M6"]:
            upload_response = self.client.post(
                f"/api/projects/{project_id}/modules/{module_id}/files",
                files={"file": (f"{module_id}.pdf", f"{module_id} material".encode(), "application/pdf")},
            )
            self.assertEqual(upload_response.status_code, 201)

        self.assertEqual(self.client.post(f"/api/projects/{project_id}/generate").status_code, 202)
        review_response = self.client.post(
            f"/api/projects/{project_id}/review",
            json={"approved": True, "notes": "端到端验收确认"},
        )

        self.assertEqual(review_response.status_code, 200)
        reviewed = review_response.json()
        self.assertEqual(reviewed["task_status"], "完成")
        self.assertEqual(reviewed["status_history"][-3:], ["章节渲染中", "合并中", "完成"])
        modules_by_id = {module["module_id"]: module for module in reviewed["modules"]}
        self.assertEqual(modules_by_id["M5"]["status"], "skipped")
        self.assertEqual(modules_by_id["M5"]["chapter_ppt_path"], "")
        self.assertTrue(all(module["status"] == "rendered" for module_id, module in modules_by_id.items() if module_id != "M5"))

        detail = self.client.get(f"/api/projects/{project_id}").json()
        final_path = Path(detail["final_ppt_path"])
        self.assertTrue(final_path.exists())
        self.assertIn("M1_M2_M3_M6", final_path.name)
        final_text = " ".join(
            shape.text for slide in Presentation(str(final_path)).slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        self.assertIn("端到端验收项目", final_text)
        self.assertNotIn("{{m3_", final_text)

        download_response = self.client.get(f"/api/projects/{project_id}/download")
        self.assertEqual(download_response.status_code, 200)
        self.assertTrue(download_response.content.startswith(b"PK"))

    def test_review_failure_marks_project_and_modules_failed(self):
        project_id = self.client.post("/api/projects", json={"project_name": "失败路径项目"}).json()["project_id"]
        self.assertEqual(self.client.post(f"/api/projects/{project_id}/generate").status_code, 202)

        with patch("app.storage.render_project_ppt", side_effect=RuntimeError("render failed")):
            review_response = self.client.post(
                f"/api/projects/{project_id}/review",
                json={"approved": True, "notes": "触发失败路径"},
            )

        self.assertEqual(review_response.status_code, 200)
        payload = review_response.json()
        self.assertEqual(payload["task_status"], "失败")
        self.assertTrue(all(module["status"] == "failed" for module in payload["modules"]))
        self.assertTrue(all(module["status_history"][-1] == "failed" for module in payload["modules"]))

    def test_assets_are_filtered_by_module(self):
        response = self.client.get("/api/assets", params={"module_id": "M5"})

        self.assertEqual(response.status_code, 200)
        assets = response.json()
        self.assertGreaterEqual(len(assets), 1)
        self.assertTrue(all(asset["module_id"] == "M5" for asset in assets))

    def test_cors_allows_frontend_dev_server(self):
        response = self.client.options(
            "/api/projects",
            headers={
                "Origin": "http://127.0.0.1:3001",
                "Access-Control-Request-Method": "POST",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.headers["access-control-allow-origin"], "http://127.0.0.1:3001")

    def test_document_parse_test_accepts_txt_and_returns_parsed_result(self):
        response = self.client.post(
            "/api/document-parse-test",
            files=[("files", ("test.txt", b"Metro Project Brief\nProject: Nanjing Metro Line 3\nType: Railway", "text/plain"))],
        )

        self.assertEqual(response.status_code, 200)
        results = response.json()
        self.assertEqual(len(results), 1)
        result = results[0]
        self.assertEqual(result["filename"], "test.txt")
        self.assertEqual(result["suffix"], ".txt")
        self.assertEqual(result["parse_status"], "parsed")
        self.assertIn("document_role", result)
        self.assertIn("text", result)
        self.assertTrue(len(result["text"]) > 0)
        self.assertEqual(result["error_message"], "")

    def test_document_parse_test_accepts_multiple_files_including_office_formats(self):
        response = self.client.post(
            "/api/document-parse-test",
            files=[
                ("files", ("project.txt", b"Highway noise barrier project", "text/plain")),
                ("files", ("docx_test.docx", _docx_bytes("highway tender project"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
                ("files", ("xlsx_test.xlsx", _xlsx_bytes("survey pain point"), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")),
                ("files", ("pptx_test.pptx", _pptx_bytes("enterprise CNAS"), "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ],
        )

        self.assertEqual(response.status_code, 200)
        results = response.json()
        self.assertEqual(len(results), 4)
        filenames = {r["filename"] for r in results}
        self.assertEqual(filenames, {"project.txt", "docx_test.docx", "xlsx_test.xlsx", "pptx_test.pptx"})
        for result in results:
            self.assertIn("parse_status", result)
            self.assertIn("document_role", result)
            self.assertIn("assigned_modules", result)
            self.assertIn("text_preview", result)

    def test_document_parse_test_returns_pending_enhancement_for_images_and_cad(self):
        response = self.client.post(
            "/api/document-parse-test",
            files=[
                ("files", ("plan.png", b"\x89PNG\r\n\x1a\n", "image/png")),
                ("files", ("drawing.dwg", b"CAD content", "application/vnd.dwg")),
            ],
        )

        self.assertEqual(response.status_code, 200)
        results = response.json()
        self.assertEqual(len(results), 2)
        for result in results:
            self.assertEqual(result["parse_status"], "pending_enhancement")
            self.assertIn("error_message", result)
            self.assertTrue("pending_enhancement" in result["error_message"] or "图片" in result["error_message"] or "CAD" in result["error_message"])
            self.assertEqual(result["text"], "")

    def test_document_parse_test_returns_pending_ocr_for_empty_pdf(self):
        # PDF with no extractable text - simulated by passing text bytes as PDF
        response = self.client.post(
            "/api/document-parse-test",
            files=[("files", ("empty.pdf", b"%PDF-1.4\n%%EOF\n", "application/pdf"))],
        )

        self.assertEqual(response.status_code, 200)
        results = response.json()
        self.assertEqual(len(results), 1)
        result = results[0]
        self.assertIn(result["parse_status"], ["pending_ocr", "parsed"])
        self.assertIn(result["filename"], "empty.pdf")

    def test_document_parse_test_rejects_unsupported_extension(self):
        response = self.client.post(
            "/api/document-parse-test",
            files=[("files", ("malware.exe", b"MZ\x90\x00", "application/x-msdownload"))],
        )

        self.assertEqual(response.status_code, 200)
        results = response.json()
        self.assertEqual(len(results), 1)
        self.assertEqual(results[0]["parse_status"], "failed")
        self.assertIn("不支持", results[0]["error_message"])

    def test_document_parse_test_does_not_create_project_or_write_to_main_flow(self):
        import json
        list_before = self.client.get("/api/projects").json()

        response = self.client.post(
            "/api/document-parse-test",
            files=[("files", ("test.txt", b"test content", "text/plain"))],
        )

        self.assertEqual(response.status_code, 200)
        list_after = self.client.get("/api/projects").json()
        self.assertEqual(len(list_before), len(list_after))

    def test_document_parse_test_includes_required_fields_in_response(self):
        response = self.client.post(
            "/api/document-parse-test",
            files=[("files", ("doc.docx", _docx_bytes("tender project highway railway"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))],
        )

        self.assertEqual(response.status_code, 200)
        result = response.json()[0]
        required_fields = [
            "filename", "suffix", "content_type", "parse_status",
            "document_role", "assigned_modules", "text", "text_preview",
            "sections", "tables", "slides", "metadata", "error_message",
        ]
        for field in required_fields:
            self.assertIn(field, result, f"Missing field: {field}")

    def test_analyze_project_does_not_write_to_vector_store(self):
        """Verify that analyze_project does NOT write to vector store."""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "向量库测试项目",
                "project_location": "南京",
                "product_line": "轨道交通声屏障",
            },
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf")),
            ],
        )

        # Make sure vector DSN is NOT set
        old_dsn = os.environ.pop("ZHONGCHI_VECTOR_DSN", None)
        try:
            analyze_response = self.client.post(f"/api/projects/{project_id}/analyze")
            self.assertEqual(analyze_response.status_code, 200)
            classification = analyze_response.json()

            # analyze should NOT have written to vector store
            # vector_status should still be "not_indexed" for all files
            for file_record in classification.get("files", []):
                self.assertEqual(file_record.get("vector_status", "not_indexed"), "not_indexed")
                self.assertEqual(file_record.get("vector_chunk_count", 0), 0)
        finally:
            if old_dsn:
                os.environ["ZHONGCHI_VECTOR_DSN"] = old_dsn

    def test_vector_index_returns_proper_response_structure(self):
        """Verify the vector-index endpoint returns proper response structure."""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "向量库索引测试",
                "project_location": "南京",
                "product_line": "轨道交通声屏障",
            },
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf")),
            ],
        )
        self.client.post(f"/api/projects/{project_id}/analyze")

        # Call vector-index
        vector_response = self.client.post(
            f"/api/projects/{project_id}/vector-index",
            json={},
        )
        self.assertEqual(vector_response.status_code, 200)
        result = vector_response.json()

        # Check response structure
        self.assertEqual(result["project_id"], project_id)
        self.assertIn("status", result)
        self.assertIn("indexed_files", result)
        self.assertIn("indexed_chunks", result)
        self.assertIn("skipped_files", result)
        self.assertIn("message", result)

        # Without ZHONGCHI_VECTOR_DSN, should return not_configured
        self.assertEqual(result["status"], "not_configured")
        self.assertIn("向量库未配置", result["message"])

    def test_vector_index_with_file_ids_filtering(self):
        """Test vector-index with specific file_ids in request body."""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "选择性索引测试",
                "project_location": "南京",
                "product_line": "轨道交通声屏障",
            },
        ).json()["project_id"]

        upload_response = self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("南京地铁项目简介.pdf", b"metro line noise barrier", "application/pdf")),
                ("files", ("企业资质介绍.pptx", b"enterprise CNAS", "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ],
        )
        self.assertEqual(upload_response.status_code, 201)
        files = upload_response.json()
        self.assertEqual(len(files), 2)

        self.client.post(f"/api/projects/{project_id}/analyze")

        # Index only the first file
        first_file_id = files[0]["file_id"]
        vector_response = self.client.post(
            f"/api/projects/{project_id}/vector-index",
            json={"file_ids": [first_file_id]},
        )
        self.assertEqual(vector_response.status_code, 200)
        result = vector_response.json()
        self.assertEqual(result["project_id"], project_id)
        self.assertIn("indexed_files", result)
        self.assertIn("indexed_chunks", result)

    def test_vector_index_skips_pending_enhancement_files(self):
        """Test that vector-index skips files with parse_status pending_enhancement."""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "跳过待增强文件测试", "product_line": "轨道交通声屏障"},
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("总平面图.png", b"\x89PNG\r\n\x1a\n", "image/png"))],
        )
        self.client.post(f"/api/projects/{project_id}/analyze")

        vector_response = self.client.post(
            f"/api/projects/{project_id}/vector-index",
            json={},
        )
        self.assertEqual(vector_response.status_code, 200)
        result = vector_response.json()
        # Image file should be skipped
        self.assertGreater(len(result["skipped_files"]), 0)

    def test_stored_file_includes_vector_status_fields(self):
        """Verify StoredFile response includes vector_status fields."""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "向量状态字段测试", "product_line": "公路声屏障"},
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("技术标书.docx", _docx_bytes("招标文件 技术标准 高速 公路 声屏障"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))],
        )
        self.client.post(f"/api/projects/{project_id}/analyze")

        classification = self.client.get(f"/api/projects/{project_id}/classification").json()
        files = classification.get("files", [])

        self.assertGreater(len(files), 0)
        for f in files:
            # All parsed files should have vector_status fields
            self.assertIn("vector_status", f)
            self.assertIn("vector_chunk_count", f)
            self.assertIn("vector_error_message", f)

    def test_analyze_project_returns_text_preview_for_parsed_files(self):
        """验证 analyze_project 返回的 files 中每个 parsed 文件包含 text_preview。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "文本预览测试",
                "project_location": "北京",
                "product_line": "公路声屏障",
            },
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("技术标书.docx", _docx_bytes("招标文件 技术标准 高速 公路 声屏障"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
            ],
        )

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        files_by_name = {item["filename"]: item for item in classification["files"]}
        docx_file = files_by_name["技术标书.docx"]
        self.assertEqual(docx_file["parse_status"], "parsed")
        self.assertIn("text_preview", docx_file)
        self.assertTrue(len(docx_file["text_preview"]) > 0)
        # text_preview 是截断文本，不应等于完整文本（如果有的话）
        self.assertLessEqual(len(docx_file["text_preview"]), 1500)

    def test_analyze_project_pending_file_has_empty_text_preview(self):
        """验证 pending_enhancement 文件 text_preview 为空字符串。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "待增强文件测试", "product_line": "公路声屏障"},
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("总平面图.png", b"\x89PNG\r\n\x1a\n", "image/png"))],
        )

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        png_file = next(item for item in classification["files"] if item["filename"] == "总平面图.png")
        self.assertEqual(png_file["parse_status"], "pending_enhancement")
        self.assertIn("text_preview", png_file)
        self.assertEqual(png_file["text_preview"], "")

    def test_get_parsed_text_returns_file_text(self):
        """验证 GET /api/projects/{id}/files/{file_id}/parsed-text 返回完整文本。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "全文接口测试",
                "project_location": "北京",
                "product_line": "公路声屏障",
            },
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[
                ("files", ("技术标书.docx", _docx_bytes("招标文件 技术标准 高速 公路 声屏障"), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
            ],
        )
        self.client.post(f"/api/projects/{project_id}/analyze")

        classification = self.client.get(f"/api/projects/{project_id}/classification").json()
        file_id = next(item["file_id"] for item in classification["files"] if item["filename"] == "技术标书.docx")

        response = self.client.get(f"/api/projects/{project_id}/files/{file_id}/parsed-text")
        self.assertEqual(response.status_code, 200)
        result = response.json()
        self.assertEqual(result["file_id"], file_id)
        self.assertEqual(result["parse_status"], "parsed")
        self.assertIn("招标文件", result["text"])
        self.assertEqual(result["error_message"], "")

    def test_get_parsed_text_returns_empty_for_unparsed_file(self):
        """验证未解析文件返回空文本和非 parsed 状态，不报错误。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "未解析文件测试", "product_line": "公路声屏障"},
        ).json()["project_id"]

        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("总平面图.png", b"\x89PNG\r\n\x1a\n", "image/png"))],
        )
        self.client.post(f"/api/projects/{project_id}/analyze")

        classification = self.client.get(f"/api/projects/{project_id}/classification").json()
        file_id = next(item["file_id"] for item in classification["files"] if item["filename"] == "总平面图.png")

        response = self.client.get(f"/api/projects/{project_id}/files/{file_id}/parsed-text")
        self.assertEqual(response.status_code, 200)
        result = response.json()
        self.assertEqual(result["file_id"], file_id)
        self.assertEqual(result["parse_status"], "pending_enhancement")
        self.assertEqual(result["text"], "")
        self.assertNotEqual(result["error_message"], "")

    def test_get_parsed_text_404_for_nonexistent_project(self):
        """验证不存在的项目返回 404。"""
        response = self.client.get("/api/projects/99999/files/1/parsed-text")
        self.assertEqual(response.status_code, 404)

    def test_get_parsed_text_404_for_nonexistent_file(self):
        """验证不存在的文件返回 404。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "不存在文件测试", "product_line": "公路声屏障"},
        ).json()["project_id"]

        response = self.client.get(f"/api/projects/{project_id}/files/99999/parsed-text")
        self.assertEqual(response.status_code, 404)

    # ---- M5 案例确认校验测试 ----

    def test_review_rejects_invalid_confirmed_case_id(self):
        """验证传入无效 confirmed_case_id 时返回 400。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "案例校验项目", "product_line": "轨道交通声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": "999",
                "notes": "传入不存在的案例 ID",
            },
        )
        self.assertEqual(response.status_code, 400)
        self.assertIn("confirmed_case_id", response.json()["detail"])

    def test_review_rejects_abc_case_id(self):
        """验证传入 abc 等无效字符串 case_id 时返回 400。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "案例校验字符串项目", "product_line": "轨道交通声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": "abc",
                "notes": "传入非法字符串案例 ID",
            },
        )
        self.assertEqual(response.status_code, 400)

    def test_review_accepts_valid_recommended_case_id(self):
        """验证传入推荐案例中的 case_id 时正常通过。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "有效案例项目", "product_line": "轨道交通声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        recommended = classification.get("case_selection", {}).get("recommended_cases", [])
        if not recommended:
            self.skipTest("无推荐案例，跳过")
        valid_case_id = recommended[0]["case_id"]

        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": valid_case_id,
                "notes": "使用推荐案例",
            },
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["case_selection"]["confirmed_case_id"], valid_case_id)

    def test_review_accepts_demo_case_id(self):
        """验证传入演示固定案例库中的 case_id（如整数 1）时正常通过。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "演示案例项目", "product_line": "轨道交通声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": 1,
                "notes": "使用演示固定案例",
            },
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["case_selection"]["confirmed_case_id"], 1)

    # ---- LLM confidence 校验测试 ----

    def test_llm_confidence_88_falls_back_to_rules(self):
        """验证 LLM 返回 confidence=88（百分制）时 fallback 到规则。"""
        from app.m1m2_classifier import _normalize_llm_result

        data = {
            "project_type": "metro",
            "confidence": 88,
            "matched_keywords": ["地铁"],
            "evidence": [],
            "reasoning_summary": "test",
        }
        with self.assertRaises(RuntimeError) as ctx:
            _normalize_llm_result(data)
        self.assertIn("超出 0-1 范围", str(ctx.exception))

    def test_llm_confidence_nan_falls_back_to_rules(self):
        """验证 LLM 返回 confidence=NaN 时 fallback。"""
        import math
        from app.m1m2_classifier import _normalize_llm_result

        data = {
            "project_type": "metro",
            "confidence": float("nan"),
            "matched_keywords": ["地铁"],
            "evidence": [],
            "reasoning_summary": "test",
        }
        with self.assertRaises(RuntimeError) as ctx:
            _normalize_llm_result(data)
        self.assertIn("非有限值", str(ctx.exception))

    def test_llm_confidence_inf_falls_back_to_rules(self):
        """验证 LLM 返回 confidence=Infinity 时 fallback。"""
        from app.m1m2_classifier import _normalize_llm_result

        data = {
            "project_type": "metro",
            "confidence": float("inf"),
            "matched_keywords": ["地铁"],
            "evidence": [],
            "reasoning_summary": "test",
        }
        with self.assertRaises(RuntimeError) as ctx:
            _normalize_llm_result(data)
        self.assertIn("非有限值", str(ctx.exception))

    def test_llm_confidence_negative_falls_back_to_rules(self):
        """验证 LLM 返回 confidence=-0.5 时 fallback。"""
        from app.m1m2_classifier import _normalize_llm_result

        data = {
            "project_type": "metro",
            "confidence": -0.5,
            "matched_keywords": ["地铁"],
            "evidence": [],
            "reasoning_summary": "test",
        }
        with self.assertRaises(RuntimeError) as ctx:
            _normalize_llm_result(data)
        self.assertIn("超出 0-1 范围", str(ctx.exception))

    def test_llm_confidence_valid_value_accepted(self):
        """验证 LLM 返回合法 confidence=0.85 时正常接受。"""
        from app.m1m2_classifier import _normalize_llm_result

        data = {
            "project_type": "metro",
            "confidence": 0.85,
            "matched_keywords": ["地铁"],
            "evidence": [],
            "reasoning_summary": "test",
        }
        result = _normalize_llm_result(data)
        self.assertEqual(result.project_type, "metro")
        self.assertAlmostEqual(result.confidence, 0.85)

    # ---- existing_rail_transit 规则收紧测试 ----

    def test_highway_reconstruction_not_classified_as_existing_rail_transit(self):
        """验证'高速公路改造项目'不应被判为 existing_rail_transit。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "沪宁高速公路声屏障改造项目",
                "product_line": "公路声屏障",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("高速公路降噪改造.pdf", "高速 公路 全封闭 声屏障 改造".encode(), "application/pdf"))],
        )

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        self.assertNotEqual(
            classification["detected_project_type"], "existing_rail_transit",
            "高速公路改造项目不应被判为既有线轨交",
        )
        self.assertEqual(classification["detected_project_type"], "highway")

    def test_metro_reconstruction_without_strong_constraints_leans_metro(self):
        """验证'地铁声屏障改造'无既有线/运营线/天窗等强约束时更偏向 metro。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "某城市地铁声屏障改造工程",
                "product_line": "地铁声屏障",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目改造.pdf", "地铁 声屏障 改造 车站 区间".encode(), "application/pdf"))],
        )

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        self.assertEqual(
            classification["detected_project_type"], "metro",
            "地铁声屏障改造无强约束时应偏向 metro",
        )

    def test_existing_rail_transit_with_strong_constraints_still_works(self):
        """验证同时出现轨交上下文和既有线/天窗等强约束时仍判为 existing_rail_transit。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "既有轨道交通线路声屏障改造工程",
                "product_line": "轨交既有线改造",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("既有线夜间施工窗口.pdf", "既有线 轨道交通 改造 夜间 施工窗口".encode(), "application/pdf"))],
        )

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        self.assertEqual(classification["detected_project_type"], "existing_rail_transit")

    def test_metro_with_short_window_classified_as_existing_rail_transit(self):
        """验证'地铁 + 短窗口'（强约束）应判为 existing_rail_transit，而非 metro。"""
        project_id = self.client.post(
            "/api/projects",
            json={
                "project_name": "某城市地铁声屏障工程",
                "product_line": "地铁声屏障",
            },
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁短窗口施工.pdf", "地铁 轨道交通 声屏障 短窗口 天窗点作业".encode(), "application/pdf"))],
        )

        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        self.assertEqual(
            classification["detected_project_type"], "existing_rail_transit",
            "地铁项目含短窗口/天窗等强约束时应判为既有线轨交",
        )

    # ---- M5 文件夹案例库扫描测试 ----

    def test_m5_folder_cases_appear_in_api_cases(self):
        """验证 GET /api/cases 返回 M5 文件夹中的固定案例。"""
        response = self.client.get("/api/cases")
        self.assertEqual(response.status_code, 200)
        cases = response.json()

        # 应包含 M5 文件夹案例
        m5_cases = [c for c in cases if c.get("source_type") == "fixed_m5"]
        self.assertGreaterEqual(len(m5_cases), 3, "应至少包含 3 个 M5 文件夹案例")

        # 验证字段完整性
        required_fields = ["case_id", "title", "filename", "project_type", "source_path", "module_id", "source_type"]
        for case in m5_cases:
            for field in required_fields:
                self.assertIn(field, case, f"M5 案例缺少字段: {field}")
            self.assertTrue(case["case_id"].startswith("fixed_m5_case:"), f"case_id 格式错误: {case['case_id']}")
            self.assertEqual(case["module_id"], "M5")
            self.assertEqual(case["source_type"], "fixed_m5")

        # 验证类型映射
        type_map = {c["project_type"] for c in m5_cases}
        self.assertIn("highway", type_map, "应有 highway 类型案例")
        self.assertIn("railway", type_map, "应有 railway 类型案例")
        self.assertIn("metro", type_map, "应有 metro 类型案例")

    def test_m5_folder_case_id_is_stable(self):
        """验证 M5 案例 case_id 基于文件名 hash，多次调用结果一致。"""
        from app.m5_case_scanner import scan_m5_cases

        first = scan_m5_cases()
        second = scan_m5_cases()

        self.assertEqual(len(first), len(second))
        for a, b in zip(first, second):
            self.assertEqual(a["case_id"], b["case_id"])
            self.assertEqual(a["title"], b["title"])

    def test_m5_recommend_highway_for_highway_project(self):
        """验证 highway 项目类型推荐公路开头的 M5 案例。"""
        from app.m5_case_scanner import recommend_m5_case

        case = recommend_m5_case("highway")
        self.assertIsNotNone(case, "highway 应有推荐案例")
        self.assertTrue(case["filename"].startswith("公路"), f"文件名应以'公路'开头: {case['filename']}")
        self.assertEqual(case["project_type"], "highway")

    def test_m5_recommend_railway_for_railway_project(self):
        """验证 railway 项目类型推荐铁路开头的 M5 案例。"""
        from app.m5_case_scanner import recommend_m5_case

        case = recommend_m5_case("railway")
        self.assertIsNotNone(case, "railway 应有推荐案例")
        self.assertTrue(case["filename"].startswith("铁路"), f"文件名应以'铁路'开头: {case['filename']}")
        self.assertEqual(case["project_type"], "railway")

    def test_m5_recommend_metro_for_metro_project(self):
        """验证 metro 项目类型推荐轨道交通开头的 M5 案例。"""
        from app.m5_case_scanner import recommend_m5_case

        case = recommend_m5_case("metro")
        self.assertIsNotNone(case, "metro 应有推荐案例")
        self.assertTrue(case["filename"].startswith("轨道交通"), f"文件名应以'轨道交通'开头: {case['filename']}")
        self.assertEqual(case["project_type"], "metro")

    def test_m5_recommend_metro_for_existing_rail_transit(self):
        """验证 existing_rail_transit 暂时推荐轨道交通开头的 M5 案例。"""
        from app.m5_case_scanner import recommend_m5_case

        case = recommend_m5_case("existing_rail_transit")
        self.assertIsNotNone(case, "existing_rail_transit 应有推荐案例")
        self.assertTrue(case["filename"].startswith("轨道交通"), f"文件名应以'轨道交通'开头: {case['filename']}")

    def test_m5_case_id_accepted_in_review(self):
        """验证 M5 文件夹案例的 case_id 能通过 review 校验。"""
        from app.m5_case_scanner import scan_m5_cases

        m5_cases = scan_m5_cases()
        if not m5_cases:
            self.skipTest("M5 文件夹无案例")

        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "M5 校验测试", "product_line": "公路声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("高速项目简介.pdf", b"highway noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        # 用 M5 文件夹案例的 case_id 做 review
        m5_case_id = m5_cases[0]["case_id"]
        response = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "highway",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": m5_case_id,
                "notes": "使用 M5 文件夹案例",
            },
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["case_selection"]["confirmed_case_id"], m5_case_id)

    def test_analyze_recommends_m5_folder_case_for_metro(self):
        """验证 analyze 返回的 recommended_cases 包含 M5 文件夹推荐案例。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "地铁 M5 推荐测试", "product_line": "地铁声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁项目简介.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        recommended = classification.get("case_selection", {}).get("recommended_cases", [])
        m5_cases = [c for c in recommended if c.get("case_id", "").startswith("fixed_m5_case:")]
        self.assertGreater(len(m5_cases), 0, "recommended_cases 应包含 M5 文件夹案例")
        # 通过 source_path 验证是轨道交通文件
        self.assertIn("轨道交通", m5_cases[0]["source_path"])

    def test_analyze_recommended_cases_returns_exactly_one_m5_case(self):
        """验证 analyze 对 4 类 project_type 的 recommended_cases 长度为 1，且以 fixed_m5_case: 开头。"""
        test_cases = [
            ("高速 M5 长度测试", "公路声屏障", [("files", ("高速.pdf", b"highway noise barrier", "application/pdf"))]),
            ("铁路 M5 长度测试", "铁路声屏障", [("files", ("铁路.pdf", b"railway noise barrier", "application/pdf"))]),
            ("地铁 M5 长度测试", "地铁声屏障", [("files", ("地铁.pdf", b"metro line noise barrier", "application/pdf"))]),
            ("既有线 M5 长度测试", "轨交既有线改造", [("files", ("既有线.pdf", "既有线 轨道交通 改造 夜间 施工窗口".encode(), "application/pdf"))]),
        ]
        for name, product_line, files in test_cases:
            with self.subTest(name=name):
                project_id = self.client.post(
                    "/api/projects",
                    json={"project_name": name, "product_line": product_line},
                ).json()["project_id"]
                self.client.post(f"/api/projects/{project_id}/files", files=files)
                classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

                recommended = classification.get("case_selection", {}).get("recommended_cases", [])
                self.assertEqual(len(recommended), 1, f"{name}: recommended_cases 应恰好 1 个，实际 {len(recommended)}")
                self.assertTrue(
                    recommended[0]["case_id"].startswith("fixed_m5_case:"),
                    f"{name}: case_id 应以 fixed_m5_case: 开头，实际 {recommended[0]['case_id']}",
                )

    # ---- M5 fixed_m5_case 生成链路测试 ----

    def test_generate_with_fixed_m5_case_copies_source_pptx(self):
        """验证选择 fixed_m5_case 后 generate 的 M5 章节文件与源 PPTX 一致。"""
        from app.m5_case_scanner import recommend_m5_case

        # 先确定一个 fixed_m5_case 及其源文件
        m5_case = recommend_m5_case("highway")
        self.assertIsNotNone(m5_case, "应有 highway 类型 M5 案例")
        source_bytes = Path(m5_case["source_path"]).read_bytes()

        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "M5 源文件复制测试", "product_line": "公路声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("高速.pdf", b"highway noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()
        recommended = classification["case_selection"]["recommended_cases"]
        self.assertEqual(len(recommended), 1)
        fixed_case_id = recommended[0]["case_id"]
        self.assertTrue(fixed_case_id.startswith("fixed_m5_case:"))

        review_resp = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "highway",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": fixed_case_id,
                "notes": "使用 fixed_m5_case",
            },
        )
        self.assertEqual(review_resp.status_code, 200)

        generate_resp = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_resp.status_code, 202)
        generated = generate_resp.json()
        m5_module = next(m for m in generated["modules"] if m["module_id"] == "M5")
        self.assertEqual(m5_module["status"], "rendered")
        self.assertTrue(m5_module["chapter_ppt_path"])

        # 验证 M5 章节文件内容与源 PPTX 一致
        chapter_bytes = Path(m5_module["chapter_ppt_path"]).read_bytes()
        self.assertEqual(chapter_bytes, source_bytes, "M5 章节文件应与源 PPTX 字节一致")

    def test_generate_with_m5_demo_case_id_uses_fallback_template(self):
        """验证 m5_demo / case_id=1 等非 fixed_m5_case 仍使用 M5示例.pptx。"""
        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "M5 demo 兼容测试", "product_line": "地铁声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("地铁.pdf", b"metro line noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        # 使用旧 case_id=1（演示案例）
        review_resp = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "metro",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": 1,
                "notes": "使用旧演示案例",
            },
        )
        self.assertEqual(review_resp.status_code, 200)

        generate_resp = self.client.post(f"/api/projects/{project_id}/generate")
        self.assertEqual(generate_resp.status_code, 202)
        generated = generate_resp.json()
        m5_module = next(m for m in generated["modules"] if m["module_id"] == "M5")
        self.assertEqual(m5_module["status"], "rendered")
        self.assertTrue(m5_module["chapter_ppt_path"])
        self.assertTrue(Path(m5_module["chapter_ppt_path"]).exists())

    def test_renderer_rejects_fixed_m5_source_path_outside_m5_dir(self):
        """验证 renderer 拒绝 M5 目录外的 source_path。"""
        from ppt_engine.renderer import _render_m5_case_template

        case_data = {
            "case_id": "fixed_m5_case:fake",
            "source_path": str(Path(__file__).resolve()),  # 指向测试文件自身
            "filename": "test_api.py",
            "title": "非法路径",
            "source_type": "fixed_m5",
        }
        with tempfile.TemporaryDirectory() as tmp:
            with self.assertRaises(ValueError) as ctx:
                _render_m5_case_template(case_data, {}, Path(tmp))
            self.assertIn("超出 M5 目录范围", str(ctx.exception))

    def test_generate_fails_when_fixed_m5_case_disappears_after_review(self):
        """验证 review 后源案例失效时 generate 应失败，不静默回退到 M5示例.pptx。"""
        from app.m5_case_scanner import recommend_m5_case

        m5_case = recommend_m5_case("railway")
        self.assertIsNotNone(m5_case, "应有 railway 类型 M5 案例")
        fixed_case_id = m5_case["case_id"]

        project_id = self.client.post(
            "/api/projects",
            json={"project_name": "M5 源案例失效测试", "product_line": "铁路声屏障"},
        ).json()["project_id"]
        self.client.post(
            f"/api/projects/{project_id}/files",
            files=[("files", ("铁路.pdf", b"railway noise barrier", "application/pdf"))],
        )
        classification = self.client.post(f"/api/projects/{project_id}/analyze").json()

        review_resp = self.client.post(
            f"/api/projects/{project_id}/classification/review",
            json={
                "confirmed_project_type": "railway",
                "template_selection": classification["template_selection"],
                "confirmed_case_id": fixed_case_id,
                "notes": "确认铁路案例",
            },
        )
        self.assertEqual(review_resp.status_code, 200)

        # 模拟 M5 扫描解析不到该 case（文件被移动/删除）
        with patch("app.ppt_generation.get_m5_case_by_id", return_value=None):
            generate_resp = self.client.post(f"/api/projects/{project_id}/generate")

        self.assertEqual(generate_resp.status_code, 202)
        generated = generate_resp.json()
        self.assertEqual(generated["task_status"], "失败")
        m5_module = next(m for m in generated["modules"] if m["module_id"] == "M5")
        self.assertEqual(m5_module["status"], "failed")
        # 确认错误信息中包含提示
        self.assertIn("fixed_m5_case", m5_module["error_message"])


class RemovedM3PartialRenderTest(unittest.TestCase):
    """验证 M3 文字/图片拆分测试接口已下线。"""

    def setUp(self):
        from fastapi.testclient import TestClient

        app_module = importlib.import_module("app.main")
        self.client = TestClient(app_module.app)

    def test_m3_text_render_endpoint_removed(self):
        response = self.client.post("/api/test/m3-render", json={"project_name": "已下线"})
        self.assertEqual(response.status_code, 404)

    def test_m3_text_render_download_endpoint_removed(self):
        response = self.client.get("/api/test/m3-render/download/nonexistent_file.pptx")
        self.assertEqual(response.status_code, 404)

    def test_m3_image_render_endpoint_removed(self):
        response = self.client.post("/api/test/m3-image-render")
        self.assertEqual(response.status_code, 404)

    def test_m3_image_render_download_endpoint_removed(self):
        response = self.client.get("/api/test/m3-image-render/download/nonexistent_file.pptx")
        self.assertEqual(response.status_code, 404)


if __name__ == "__main__":
    unittest.main()
