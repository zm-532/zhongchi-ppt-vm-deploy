"""M3 图片替换独立测试（不接入正式生产流程）。"""

import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation

from ppt_engine.renderer import (
    M3_IMAGE_PURPOSES,
    M3_IMAGE_TEST_TEMPLATE_FILENAME,
    M3_TEMPLATE_FILENAME,
    PPT_TEMPLATE_ROOT,
    render_m3_image_test_ppt,
    validate_m3_image_template_placeholders,
)


def _png_bytes(color: tuple[int, int, int] = (30, 90, 160), size: tuple[int, int] = (640, 360)) -> bytes:
    buffer = BytesIO()
    Image.new("RGB", size, color).save(buffer, format="PNG")
    return buffer.getvalue()


class M3ImageRendererTest(unittest.TestCase):
    def test_m3_image_test_template_exists_separately(self):
        image_template = PPT_TEMPLATE_ROOT / M3_IMAGE_TEST_TEMPLATE_FILENAME
        text_template = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME

        self.assertTrue(image_template.exists(), f"M3 图片测试模板不存在：{image_template}")
        self.assertTrue(text_template.exists(), f"M3 文字模板不存在：{text_template}")
        self.assertNotEqual(image_template.name, text_template.name)

    def test_m3_image_purposes_are_fixed(self):
        self.assertEqual(
            M3_IMAGE_PURPOSES,
            {
                "project_scope_map": "项目建设范围图",
                "project_line_map": "项目线路图",
                "survey_route_map": "踏勘路线/点位图",
                "site_survey_photos": "现场踏勘照片组",
                "key_difficulty_evidence": "重难点证据图",
            },
        )

    def test_m3_image_template_has_all_placeholders(self):
        image_template = PPT_TEMPLATE_ROOT / M3_IMAGE_TEST_TEMPLATE_FILENAME
        issues = validate_m3_image_template_placeholders(image_template)
        self.assertEqual(issues, [])

    def test_render_m3_image_test_ppt_produces_openable_pptx(self):
        images = {
            "project_scope_map": [_png_bytes((200, 30, 30))],
            "project_line_map": [_png_bytes((30, 200, 30))],
            "survey_route_map": [_png_bytes((30, 30, 200))],
            "site_survey_photos": [
                _png_bytes((200, 180, 30)),
                _png_bytes((120, 60, 200)),
            ],
            "key_difficulty_evidence": [_png_bytes((20, 160, 180))],
        }
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_image_test_ppt("图片替换测试项目", images, temp_dir)
            self.assertTrue(output_path.exists())
            prs = Presentation(str(output_path))
            self.assertGreater(len(prs.slides), 0)

    def test_render_m3_image_test_ppt_does_not_modify_text_template(self):
        text_template = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
        before = text_template.stat().st_mtime

        with tempfile.TemporaryDirectory() as temp_dir:
            render_m3_image_test_ppt(
                "不影响文字模板",
                {"project_scope_map": [_png_bytes()]},
                temp_dir,
            )

        self.assertEqual(text_template.stat().st_mtime, before)

    def test_render_m3_image_test_rejects_invalid_purpose(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError):
                render_m3_image_test_ppt("非法用途", {"bad_purpose": [_png_bytes()]}, temp_dir)

    def test_render_m3_image_test_rejects_too_many_single_purpose_images(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError):
                render_m3_image_test_ppt(
                    "单图字段过量",
                    {"project_scope_map": [_png_bytes(), _png_bytes()]},
                    temp_dir,
                )

    def test_render_m3_image_test_rejects_corrupt_image(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError):
                render_m3_image_test_ppt("损坏图片", {"project_scope_map": [b"not-an-image"]}, temp_dir)


if __name__ == "__main__":
    unittest.main()
