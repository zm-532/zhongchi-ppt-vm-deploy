"""M3 完整测试渲染（文字 + 图片，占位符替换，不接入正式流程）。"""

import tempfile
import unittest
import zipfile
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from ppt_engine.renderer import (
    M3_FULL_SECTIONS,
    M3_FULL_TEST_TEMPLATE_FILENAME,
    PPT_TEMPLATE_ROOT,
    render_m3_full_test_ppt,
    validate_m3_full_template_placeholders,
)


def _png_bytes(color: tuple[int, int, int] = (30, 90, 160), size: tuple[int, int] = (640, 360)) -> bytes:
    buffer = BytesIO()
    Image.new("RGB", size, color).save(buffer, format="PNG")
    return buffer.getvalue()


def _xlsx_with_dimension(dimension: str) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "xl/workbook.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets><sheet name="Sheet1" sheetId="1" r:id="rId1"/></sheets>
</workbook>""",
        )
        archive.writestr(
            "xl/_rels/workbook.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <dimension ref="{dimension}"/>
  <sheetData>
    <row r="1"><c r="A1" t="inlineStr"><is><t>测试</t></is></c></row>
  </sheetData>
</worksheet>""",
        )
    return buffer.getvalue()


TEXTS = {
    section["text_field"]: f"{section['title']}测试文字"
    for section in M3_FULL_SECTIONS
}

TABLE_TEMPLATE_ROOT = PPT_TEMPLATE_ROOT / "M3表格模板"


def _first_table_on_slide(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape.table
    raise AssertionError("slide has no table")


def _cell_has_all_borders(cell) -> bool:
    tc_pr = cell._tc.tcPr
    if tc_pr is None:
        return False
    line_tags = {child.tag.rsplit("}", 1)[-1] for child in tc_pr}
    return {"lnL", "lnR", "lnT", "lnB"}.issubset(line_tags)


class M3FullRendererTest(unittest.TestCase):
    def test_m3_full_template_exists_and_has_all_sections(self):
        template = PPT_TEMPLATE_ROOT / M3_FULL_TEST_TEMPLATE_FILENAME
        self.assertTrue(template.exists(), f"M3 完整测试模板不存在：{template}")
        prs = Presentation(str(template))
        self.assertEqual(len(prs.slides), 9)

    def test_m3_full_template_has_all_text_and_image_placeholders(self):
        template = PPT_TEMPLATE_ROOT / M3_FULL_TEST_TEMPLATE_FILENAME
        self.assertEqual(validate_m3_full_template_placeholders(template), [])

    def test_m3_full_template_validation_rejects_duplicate_placeholders(self):
        prs = Presentation()
        layout = prs.slide_layouts[6]
        for _ in range(2):
            slide = prs.slides.add_slide(layout)
            slide.shapes.add_textbox(0, 0, 1000000, 1000000).text = "{{m3_basic_summary}}"
            slide.shapes.add_textbox(0, 1000000, 1000000, 1000000).text = "{{image:m3_basic}}"
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "duplicate.pptx"
            prs.save(str(path))
            issues = validate_m3_full_template_placeholders(path)
        self.assertIn("m3_basic_summary", issues)
        self.assertIn("image:m3_basic", issues)

    def test_render_m3_full_single_image_per_section_outputs_9_slides(self):
        images = {
            section["image_field"]: [_png_bytes()]
            for section in M3_FULL_SECTIONS
        }
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt("M3完整测试", TEXTS, images, temp_dir)
            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 9)

    def test_render_m3_full_duplicates_section_when_two_images_uploaded(self):
        images = {
            "image:m3_basic": [_png_bytes((200, 40, 40)), _png_bytes((40, 200, 40))],
        }
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt("M3完整多图测试", TEXTS, images, temp_dir)
            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 10)

    def test_render_m3_full_counts_multiple_sections_with_extra_images(self):
        images = {
            "image:m3_basic": [_png_bytes(), _png_bytes()],
            "image:m3_line": [_png_bytes(), _png_bytes(), _png_bytes()],
        }
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt("M3完整多部分多图", TEXTS, images, temp_dir)
            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 12)

    def test_render_m3_full_places_table_page_before_image_page(self):
        table_path = TABLE_TEMPLATE_ROOT / "敏感点路段.xlsx"
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt(
                "M3完整表格图片顺序",
                TEXTS,
                {"image:m3_sensitive_points": [_png_bytes()]},
                temp_dir,
                {"image:m3_sensitive_points": ["敏感点图片说明"]},
                {"image:m3_sensitive_points": [table_path.read_bytes()]},
            )
            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 10)
            table_slide_cells = []
            for shape in prs.slides[2].shapes:
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        table_slide_cells.extend(cell.text for cell in row.cells)
            table_slide_text = "\n".join(table_slide_cells)
            image_slide_text = "\n".join(shape.text for shape in prs.slides[3].shapes if hasattr(shape, "text"))
            self.assertIn("措施里程", table_slide_text)
            self.assertIn("DR2K0+099~DR2K0+241", table_slide_text)
            self.assertNotIn("敏感点图片说明", table_slide_text)
            self.assertIn("敏感点图片说明", image_slide_text)

    def test_render_m3_full_table_uses_grid_borders_and_larger_font(self):
        table_path = TABLE_TEMPLATE_ROOT / "敏感点路段.xlsx"
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt(
                "M3完整表格样式",
                TEXTS,
                {},
                temp_dir,
                None,
                {"image:m3_sensitive_points": [table_path.read_bytes()]},
            )
            prs = Presentation(str(output_path))
            table = _first_table_on_slide(prs.slides[2])
            self.assertTrue(_cell_has_all_borders(table.cell(0, 0)))
            self.assertTrue(_cell_has_all_borders(table.cell(1, 0)))

            header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
            body_run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
            self.assertEqual(header_run.font.size, Pt(8))
            self.assertEqual(body_run.font.size, Pt(8))

    def test_render_m3_full_table_cells_are_vertically_centered(self):
        table_path = TABLE_TEMPLATE_ROOT / "敏感点路段.xlsx"
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt(
                "M3完整表格上下居中",
                TEXTS,
                {},
                temp_dir,
                None,
                {"image:m3_sensitive_points": [table_path.read_bytes()]},
            )
            prs = Presentation(str(output_path))
            table = _first_table_on_slide(prs.slides[2])
            self.assertEqual(table.cell(0, 0).vertical_anchor, MSO_VERTICAL_ANCHOR.MIDDLE)
            self.assertEqual(table.cell(1, 0).vertical_anchor, MSO_VERTICAL_ANCHOR.MIDDLE)

    def test_render_m3_full_rejects_xlsx_with_too_many_columns(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaisesRegex(ValueError, "Excel 表格超出可渲染范围"):
                render_m3_full_test_ppt(
                    "M3完整超大表格",
                    TEXTS,
                    {},
                    temp_dir,
                    None,
                    {"image:m3_sensitive_points": [_xlsx_with_dimension("A1:M2")]},
                )

    def test_render_m3_full_uses_page_texts_for_duplicated_image_pages(self):
        images = {
            "image:m3_basic": [_png_bytes(), _png_bytes((40, 200, 40))],
        }
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt(
                "M3完整逐页描述",
                TEXTS,
                images,
                temp_dir,
                {"image:m3_basic": ["第一张描述", "第二张描述"]},
            )
            prs = Presentation(str(output_path))
            first_slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
            second_slide_text = "\n".join(shape.text for shape in prs.slides[1].shapes if hasattr(shape, "text"))
            self.assertIn("第一张描述", first_slide_text)
            self.assertIn("第二张描述", second_slide_text)

    def test_render_m3_full_removes_text_when_page_text_is_empty(self):
        images = {
            "image:m3_basic": [_png_bytes()],
        }
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt(
                "M3完整空描述",
                TEXTS,
                images,
                temp_dir,
                {"image:m3_basic": [""]},
            )
            prs = Presentation(str(output_path))
            slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
            self.assertNotIn("{{m3_basic_summary}}", slide_text)
            self.assertNotIn("项目基本情况测试文字", slide_text)

    def test_render_m3_full_removes_text_on_table_page_when_page_text_is_empty(self):
        table_path = TABLE_TEMPLATE_ROOT / "现场勘查情况.xlsx"
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_full_test_ppt(
                "M3完整表格空描述",
                TEXTS,
                {},
                temp_dir,
                {"image:m3_investigation": [""]},
                {"image:m3_investigation": [table_path.read_bytes()]},
            )
            prs = Presentation(str(output_path))
            slide_text = "\n".join(shape.text for shape in prs.slides[6].shapes if hasattr(shape, "text"))
            self.assertNotIn("{{m3_investigation_summary}}", slide_text)
            self.assertNotIn("现场勘察情况测试文字", slide_text)
            self.assertNotIn("待补充", slide_text)

    def test_render_m3_full_does_not_modify_source_template(self):
        template = PPT_TEMPLATE_ROOT / M3_FULL_TEST_TEMPLATE_FILENAME
        before = template.stat().st_mtime
        with tempfile.TemporaryDirectory() as temp_dir:
            render_m3_full_test_ppt("不修改源模板", TEXTS, {"image:m3_basic": [_png_bytes()]}, temp_dir)
        self.assertEqual(template.stat().st_mtime, before)

    def test_render_m3_full_rejects_invalid_image_purpose(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError):
                render_m3_full_test_ppt("非法用途", TEXTS, {"bad_purpose": [_png_bytes()]}, temp_dir)

    def test_render_m3_full_rejects_corrupt_image(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError):
                render_m3_full_test_ppt("损坏图片", TEXTS, {"image:m3_basic": [b"not-an-image"]}, temp_dir)


if __name__ == "__main__":
    unittest.main()
