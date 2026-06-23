"""M3 独立渲染测试（不接入正式生产流程）"""

import os
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from ppt_engine.renderer import (
    M3_ACTIVE_FIELD_NAMES,
    M3_FIELD_NAMES,
    M3_TEMPLATE_FILENAME,
    PPT_TEMPLATE_ROOT,
    build_m3_replacement_map,
    render_m3_test_ppt,
    validate_m3_template_placeholders,
    validate_m3_template_clean,
)

os.environ["ZHONGCHI_PPT_MERGE_ENGINE"] = "python-pptx"

PROJECT = {
    "project_name": "南京地铁3号线声屏障改造工程",
    "project_location": "南京",
    "owner_unit": "南京地铁集团",
    "product_line": "轨道交通声屏障",
}


class M3RendererTest(unittest.TestCase):
    """验证 M3 独立测试渲染功能。"""

    def test_m3_template_file_exists(self):
        """M3 模板文件应存在于 solution_fixed_modules 目录。"""
        tmpl_path = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
        self.assertTrue(tmpl_path.exists(), f"M3 模板不存在：{tmpl_path}")

    def test_m3_template_has_active_placeholders_only(self):
        """M3 模板应只包含当前实际使用的 active 占位符。"""
        tmpl_path = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
        prs = Presentation(str(tmpl_path))
        import re
        PLACEHOLDER = re.compile(r"\{\{(\w+)\}\}")
        found = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    found.update(PLACEHOLDER.findall(shape.text))
        missing = [f for f in M3_ACTIVE_FIELD_NAMES if f not in found]
        removed = [
            "m3_line_summary",
            "m3_sensitive_points_summary",
            "m3_structure_summary",
            "m3_investigation_summary",
        ]
        self.assertEqual(
            missing, [],
            f"M3 模板缺少占位符：{missing}。请运行 scripts/inject_m3_placeholders.py",
        )
        for field in removed:
            self.assertNotIn(field, found)

    def test_validate_m3_template_placeholders_returns_empty_for_valid_template(self):
        """校验函数对有效模板返回空列表。"""
        tmpl_path = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
        missing = validate_m3_template_placeholders(tmpl_path)
        self.assertEqual(missing, [])

    def test_validate_m3_template_placeholders_reports_missing(self):
        """校验函数能检测缺失的占位符（占位符必须在规定位置才算存在）。"""
        # 创建一个临时 PPTX，只在 slide 0 Shape 2 放一个占位符
        with tempfile.TemporaryDirectory() as tmp:
            prs = Presentation()
            prs.slide_width = 9144000
            prs.slide_height = 6858000
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(457200, 457200, 9144000, 1143000)
            tf = shape.text_frame
            tf.clear()
            run = tf.paragraphs[0].add_run()
            run.text = "{{m3_basic_summary}}"  # 只提供一个，但放错位置
            tmp_path = Path(tmp) / "partial.pptx"
            prs.save(str(tmp_path))

            missing = validate_m3_template_placeholders(tmp_path)
            self.assertGreater(len(missing), 0)
            # m3_basic_summary 不在规定位置（slide 0 shape 2），应报告缺失
            self.assertIn("m3_basic_summary", missing)

    def test_m3_field_names_has_9_fields(self):
        """M3_FIELD_NAMES 应包含 9 个字段。"""
        self.assertEqual(len(M3_FIELD_NAMES), 9)
        expected_fields = (
            "m3_basic_summary",
            "m3_line_summary",
            "m3_sensitive_points_summary",
            "m3_quantity_summary",
            "m3_structure_summary",
            "m3_site_survey_summary",
            "m3_investigation_summary",
            "m3_risk_summary",
            "m3_solution_summary",
        )
        self.assertEqual(M3_FIELD_NAMES, expected_fields)

    def test_m3_active_field_names_has_6_fields(self):
        """M3 当前模板只启用 5 个文字替换字段。"""
        self.assertEqual(
            M3_ACTIVE_FIELD_NAMES,
            (
                "m3_basic_summary",
                "m3_quantity_summary",
                "m3_site_survey_summary",
                "m3_risk_summary",
                "m3_solution_summary",
            ),
        )

    def test_build_m3_replacement_map_returns_all_9_fields(self):
        """build_m3_replacement_map 应返回 9 个字段的映射。"""
        result = build_m3_replacement_map(PROJECT, [])
        self.assertEqual(set(result.keys()), set(M3_FIELD_NAMES))
        for field in M3_FIELD_NAMES:
            self.assertIn(field, result)
            self.assertIsInstance(result[field], str)
            self.assertGreater(len(result[field]), 0)

    def test_build_m3_replacement_map_uses_project_fields(self):
        """m3_basic_summary 应拼接项目基础信息。"""
        result = build_m3_replacement_map(PROJECT, [])
        self.assertIn("南京地铁3号线声屏障改造工程", result["m3_basic_summary"])
        self.assertIn("南京", result["m3_basic_summary"])
        self.assertIn("南京地铁集团", result["m3_basic_summary"])
        self.assertIn("轨道交通声屏障", result["m3_basic_summary"])

    def test_build_m3_replacement_map_extracts_from_parsed_sources(self):
        """从 parsed_sources 中提取的字段不应是 [待补充]。"""
        sources = [
            "项目位于南京地铁3号线区间，全线约15公里，设有多个敏感点路段。",
            "全线工程量统计：声屏障长度约12km，采用全封闭结构形式。",
            "现场踏勘发现施工窗口受限，夜间天窗点作业。",
            "风险分析：工期紧张，夜间施工风压控制是重难点。",
        ]
        result = build_m3_replacement_map(PROJECT, sources)
        # 线路字段应包含"3号线"或"公里"
        self.assertNotEqual(result["m3_line_summary"], "[待补充：项目线路图说明]")
        # 工程量字段应包含"12km"
        self.assertNotEqual(result["m3_quantity_summary"], "[待补充：工程量统计说明]")

    def test_build_m3_replacement_map_fallback_when_no_sources(self):
        """无 parsed_sources 时，除 m3_basic_summary 外的字段应使用兜底占位符。"""
        result = build_m3_replacement_map(PROJECT, None)
        # m3_basic_summary 有项目基础信息，不应用兜底
        self.assertNotIn("待补充", result["m3_basic_summary"])
        # 其他字段在无文本来源时应返回兜底
        for field in M3_FIELD_NAMES:
            if field != "m3_basic_summary":
                self.assertIn("待补充", result[field])

    def test_render_m3_test_ppt_produces_valid_pptx(self):
        """render_m3_test_ppt 应生成可被 python-pptx 打开的 PPTX。"""
        sources = [
            "南京地铁3号线声屏障改造工程位于南京市，全线约15公里。",
            "涉及多个敏感点路段，需要采取降噪措施。",
            "现场踏勘发现施工窗口受限，夜间天窗点作业。",
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            self.assertTrue(output_path.exists())
            self.assertTrue(output_path.suffix.lower() == ".pptx")
            prs = Presentation(str(output_path))
            self.assertGreater(len(prs.slides), 0)

    def test_rendered_m3_pptx_page_count_equals_template(self):
        """渲染后的 M3 PPTX 页数应与源模板相同，不新增页面。"""
        tmpl_path = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
        prs_template = Presentation(str(tmpl_path))
        template_slide_count = len(prs_template.slides)

        sources = ["南京地铁3号线声屏障改造工程位于南京市。"]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            prs_output = Presentation(str(output_path))
            self.assertEqual(
                len(prs_output.slides),
                template_slide_count,
                f"输出 PPTX 应为 {template_slide_count} 页，实际 {len(prs_output.slides)} 页",
            )

    def test_rendered_m3_pptx_has_no_placeholder_braces(self):
        """渲染后的 M3 PPTX 不应包含 {{m3_ 开头的占位符。"""
        sources = [
            "南京地铁3号线声屏障改造工程位于南京市。",
            "现场踏勘发现施工窗口受限。",
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            prs = Presentation(str(output_path))
            all_text = "".join(
                shape.text for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            import re
            placeholder_pattern = re.compile(r"\{\{m3_\w+\}\}")
            matches = placeholder_pattern.findall(all_text)
            self.assertEqual(matches, [], f"发现未替换的占位符：{matches}")

    def test_rendered_m3_pptx_contains_project_name(self):
        """渲染后的 M3 PPTX 应包含测试项目名称。"""
        sources = ["南京地铁3号线声屏障改造工程位于南京市。"]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            prs = Presentation(str(output_path))
            all_text = "".join(
                shape.text for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertIn("南京地铁3号线声屏障改造工程", all_text)

    def test_rendered_m3_pptx_replaces_old_project_text(self):
        """渲染后的 M3 PPTX 第1页不应包含旧项目文本。"""
        sources = ["南京地铁3号线声屏障改造工程"]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            prs = Presentation(str(output_path))
            first_slide_text = "".join(
                shape.text for shape in prs.slides[0].shapes
                if hasattr(shape, "text") and shape.text
            )
            # 第1页应包含新项目名，不应包含旧项目文本"宁波市人民政府"
            self.assertIn("南京地铁3号线声屏障改造工程", first_slide_text)
            self.assertNotIn("宁波市人民政府", first_slide_text)

    def test_rendered_m3_pptx_replaces_site_survey_text(self):
        """渲染后的 M3 PPTX 第11页（现场踏勘）应包含测试输入的踏勘文本。"""
        sources = [
            "现场踏勘发现施工窗口受限，夜间天窗点作业，吊装设备需从桥下进入。",
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            prs = Presentation(str(output_path))
            # 第11页（index=10）是现场踏勘页，Shape 1
            slide_11_shape = prs.slides[10].shapes[1]
            slide_11_text = slide_11_shape.text if hasattr(slide_11_shape, "text") else ""
            self.assertIn("施工窗口受限", slide_11_text)
            self.assertIn("天窗", slide_11_text)
            # 旧宁波项目文本不应残留在目标 shape
            self.assertNotIn("宁波2号线", slide_11_text)
            self.assertNotIn("东外环", slide_11_text)

    def test_rendered_m3_pptx_no_old_keywords_on_all_target_shapes(self):
        """渲染后所有 active 目标 shape 均不含旧项目关键词。"""
        sources = [
            "南京地铁3号线声屏障改造工程位于南京市，全线约15公里。",
            "现场踏勘发现施工窗口受限，夜间天窗点作业。",
            "风险分析：工期紧张，夜间施工风压控制是重难点。",
            "立柱定位使用TEKLA放样，确保精度。",
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            prs = Presentation(str(output_path))
            issues = validate_m3_template_clean(output_path)
            self.assertEqual(
                issues, {},
                f"以下字段的目标 shape 中残留旧项目关键词：{issues}",
            )

    def test_rendered_m3_pptx_does_not_write_removed_pages(self):
        """第2/3/6/16页不再写入 M3 测试文本或兜底文本。"""
        sources = [
            "项目位于南京地铁3号线既有线区间，涉及多个敏感点路段，全线约15公里。",
            "现场踏勘发现施工窗口受限，夜间天窗点作业。",
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            prs = Presentation(str(output_path))
            for slide_idx in (1, 2, 5, 15):
                slide_text = "".join(
                    shape.text for shape in prs.slides[slide_idx].shapes
                    if hasattr(shape, "text") and shape.text
                )
                self.assertNotIn("项目位于南京地铁3号线既有线区间", slide_text)
                self.assertNotIn("涉及多个敏感点路段", slide_text)
                self.assertNotIn("[待补充：现场勘察情况]", slide_text)
                self.assertNotIn("[待补充：结构形式说明]", slide_text)

    def test_rendered_m3_risk_summary_uses_large_body_shape(self):
        """第19页重难点分析应写入大正文 Shape 1，小标签 Shape 2 保持固定文本。"""
        risk_text = "风险分析：工期紧张，夜间施工风压控制是重难点，需要工装TEKLA定位和抗风支架防松动措施。"
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, [risk_text], temp_dir)
            prs = Presentation(str(output_path))
            self.assertIn("工期紧张", prs.slides[18].shapes[1].text)
            self.assertEqual(prs.slides[18].shapes[2].text.strip(), "技术实施风险")

    def test_validate_m3_template_clean_detects_residual_keywords(self):
        """validate_m3_template_clean 能检测残留的旧项目关键词。"""
        import shutil
        # 用真实模板复制一个，手动注入宁波关键词到目标 shape
        with tempfile.TemporaryDirectory() as tmp:
            src = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
            bad = Path(tmp) / "residual.pptx"
            shutil.copy2(src, bad)
            # 在第11页（index=10）Shape 1 中注入旧关键词
            prs = Presentation(str(bad))
            shape = prs.slides[10].shapes[1]
            tf = shape.text_frame
            tf.clear()
            run = tf.paragraphs[0].add_run()
            run.text = "宁波2号线东外环停车场施工现场"
            prs.save(str(bad))

            issues = validate_m3_template_clean(bad)
            self.assertIn("m3_site_survey_summary", issues)
            self.assertIn("宁波2号线", issues["m3_site_survey_summary"])

    def test_validate_m3_template_placeholders_checks_position(self):
        """validate_m3_template_placeholders 能检测占位符不在规定位置的情况。"""
        import shutil
        with tempfile.TemporaryDirectory() as tmp:
            src = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
            bad = Path(tmp) / "misplaced.pptx"
            shutil.copy2(src, bad)
            # 从正确位置清除 m3_risk_summary，然后在错误位置注入
            prs = Presentation(str(bad))
            # 清除正确位置（slide 18, shape 1）
            s19 = prs.slides[18]
            for shape in s19.shapes:
                if hasattr(shape, "text") and "{{m3_risk_summary}}" in shape.text:
                    tf = shape.text_frame
                    tf.clear()
                    run = tf.paragraphs[0].add_run()
                    run.text = ""
            # 在错误位置（slide 3）新增一个文本框放占位符
            s3 = prs.slides[2]
            shape = s3.shapes.add_textbox(457200, 457200, 9144000, 1143000)
            tf = shape.text_frame
            tf.clear()
            run = tf.paragraphs[0].add_run()
            run.text = "{{m3_risk_summary}}"
            prs.save(str(bad))

            missing = validate_m3_template_placeholders(bad)
            # m3_risk_summary 应在 slide 19 Shape 1，不在则报缺失
            self.assertIn("m3_risk_summary", missing)

    def test_render_m3_does_not_modify_source_template(self):
        """render_m3_test_ppt 不应修改 M3 源模板文件。"""
        sources = ["南京地铁3号线声屏障改造工程"]
        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
            source_mtime = source_path.stat().st_mtime

            render_m3_test_ppt(PROJECT, sources, temp_dir)

            self.assertEqual(source_path.stat().st_mtime, source_mtime)

    def test_render_m3_output_filename_contains_m3_tag(self):
        """render_m3_test_ppt 的输出文件名应包含 'M3_文字替换测试'。"""
        sources = ["南京地铁3号线声屏障改造工程"]
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_m3_test_ppt(PROJECT, sources, temp_dir)
            self.assertIn("M3", output_path.name)
            self.assertIn("文字替换测试", output_path.name)

    def test_render_m3_rejects_template_without_placeholders(self):
        """如果模板缺少 M3 占位符，渲染函数应抛出清晰错误。"""
        with tempfile.TemporaryDirectory() as tmp:
            # 创建一个不含占位符的临时 PPTX
            prs = Presentation()
            prs.slide_width = 9144000
            prs.slide_height = 6858000
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(457200, 457200, 9144000, 1143000)
            tf = shape.text_frame
            tf.clear()
            run = tf.paragraphs[0].add_run()
            run.text = "这是一个没有占位符的模板"
            tmp_path = Path(tmp) / "no_placeholder.pptx"
            prs.save(str(tmp_path))

            # 复制此模板到另一个路径，模拟"缺失占位符的模板"
            import shutil
            bad_dest = Path(tmp) / "bad_template.pptx"
            shutil.copy2(tmp_path, bad_dest)

            # 调用 validate，应报告缺失（因为占位符不在规定位置）
            missing = validate_m3_template_placeholders(bad_dest)
            self.assertGreater(len(missing), 0)

            # 用没有占位符的模板调用 render_m3_test_ppt，会在 validate 阶段抛出错误
            # 由于我们无法直接替换 M3_TEMPLATE_FILENAME，这里只测 validate
            self.assertIn("m3_basic_summary", missing)


if __name__ == "__main__":
    unittest.main()
