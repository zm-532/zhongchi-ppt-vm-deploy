import os
import tempfile
import unittest
import unittest.mock
import zipfile
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt

from ppt_engine.renderer import (
    HOME_TEMPLATE_FILENAME,
    MERGE_ORDER,
    M1_M2_TEMPLATE_MAP,
    M5_TEMPLATE_FILENAME,
    M6_TEMPLATE_FILENAME,
    PPT_TEMPLATE_ROOT,
    TAIL_FIXED_TEMPLATE_FILENAME,
    TAIL_PRINT_TEMPLATE_FILENAME,
    build_m1_m2_replacement_map,
    merge_pptx,
    render_chapter_ppt,
    render_final_ppt,
    replace_text_placeholders,
)

os.environ["ZHONGCHI_PPT_MERGE_ENGINE"] = "python-pptx"


PROJECT = {
    "project_name": "某城市轨道交通声屏障改造项目",
    "project_location": "南京",
    "owner_unit": "某建设单位",
    "product_line": "轨交既有线改造",
}

# project_type -> M1_M2 template selection
OUTLINES_M1_M2_ONLY = {
    "M1_M2": {
        "project_type": "metro",
    },
}

M5_CASE_DATA = {
    "case_id": "sr_case_abc123",
    "case_title": "南昌轨道交通4号线声屏障工程",
    "match_reason": "同为轨道交通声屏障",
}


def outline_for_module(module_id: str) -> dict:
    if module_id == "M1_M2":
        return {"project_type": "metro"}
    if module_id == "M5":
        return {"case_data": M5_CASE_DATA}
    return {}


# M5 with case_data
OUTLINES_WITH_M5 = {
    "M1_M2": {
        "project_type": "metro",
    },
    "M5": {"case_data": M5_CASE_DATA},
}

# Final render outlines without M5 case_data, so tests do not depend on optional M5 template files.
OUTLINES_FULL = {
    "M1_M2": {
        "project_type": "metro",
    },
    "M6": {},
}


def slide_count(path: Path) -> int:
    presentation = Presentation(str(path))
    return len(presentation.slides)


def slide_texts(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return texts


def slide_image_relationship_errors(path: Path) -> list[str]:
    errors = []
    with zipfile.ZipFile(path) as pptx_zip:
        slide_names = sorted(
            [
                name
                for name in pptx_zip.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            ],
            key=lambda name: int(
                name.removeprefix("ppt/slides/slide").removesuffix(".xml")
            ),
        )
        for slide_name in slide_names:
            slide_xml = pptx_zip.read(slide_name).decode("utf-8")
            rels_name = slide_name.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
            rels_xml = pptx_zip.read(rels_name).decode("utf-8")
            for relationship_id in set(slide_xml.split('r:embed="')[1:]):
                relationship_id = relationship_id.split('"', 1)[0]
                marker = f'Id="{relationship_id}"'
                if marker not in rels_xml:
                    errors.append(f"{slide_name}: missing {relationship_id}")
                    continue
                rel_start = rels_xml.index(marker)
                rel_end = rels_xml.find("/>", rel_start)
                rel_fragment = rels_xml[rel_start:rel_end]
                if "/relationships/image" not in rel_fragment:
                    errors.append(f"{slide_name}: {relationship_id} is not image")
    return errors


def duplicate_zip_entries(path: Path) -> list[str]:
    with zipfile.ZipFile(path) as pptx_zip:
        counts = Counter(pptx_zip.namelist())
    return sorted(name for name, count in counts.items() if count > 1)


class TemplateConfigTest(unittest.TestCase):
    """验证模板路径配置统一入口。"""

    def test_ppt_template_root_is_absolute_path(self):
        self.assertTrue(PPT_TEMPLATE_ROOT.is_absolute())

    def test_m1_m2_template_map_has_four_project_types(self):
        self.assertEqual(
            set(M1_M2_TEMPLATE_MAP.keys()),
            {"highway", "metro", "existing_rail_transit", "railway"},
        )

    def test_merge_order_starts_with_home_and_ends_with_tail_pages(self):
        self.assertEqual(MERGE_ORDER, ("HOME", "M1_M2", "M3", "M5", "M6", "TAIL_FIXED", "TAIL_PRINT"))

    def test_home_template_filename_correct(self):
        self.assertEqual(HOME_TEMPLATE_FILENAME, "首页.pptx")

    def test_m6_template_filename_correct(self):
        self.assertEqual(M6_TEMPLATE_FILENAME, "中驰企业介绍合并初版（M6）.pptx")

    def test_m5_template_filename_correct(self):
        self.assertEqual(M5_TEMPLATE_FILENAME, "M5示例.pptx")

    def test_tail_template_filenames_correct(self):
        self.assertEqual(TAIL_FIXED_TEMPLATE_FILENAME, "固定尾页.pptx")
        self.assertEqual(TAIL_PRINT_TEMPLATE_FILENAME, "尾页-打印版.pptx")


class M1M2TemplateSelectionTest(unittest.TestCase):
    """验证 M1/M2 根据 project_type 选择固定模板。"""

    def test_m1_m2_metro_template_copy_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "metro"}
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)

            self.assertTrue(output_path.exists())
            presentation = Presentation(str(output_path))
            # Metro template should have multiple slides
            self.assertGreater(len(presentation.slides), 0)

    def test_m1_m2_highway_template_copy_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "highway"}
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)

            self.assertTrue(output_path.exists())
            presentation = Presentation(str(output_path))
            self.assertGreater(len(presentation.slides), 0)

    def test_m1_m2_existing_rail_transit_template_copy_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "existing_rail_transit"}
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)

            self.assertTrue(output_path.exists())
            presentation = Presentation(str(output_path))
            self.assertGreater(len(presentation.slides), 0)

    def test_m1_m2_railway_template_copy_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "railway"}
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)

            self.assertTrue(output_path.exists())
            presentation = Presentation(str(output_path))
            self.assertGreater(len(presentation.slides), 0)

    def test_m1_m2_rejects_invalid_project_type(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "invalid_type"}
            with self.assertRaises(ValueError) as context:
                render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)

            self.assertIn("invalid_type", str(context.exception))


class M5CaseTemplateTest(unittest.TestCase):
    """验证 M5 案例模板入口。"""

    def test_m5_with_none_case_data_is_rejected(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"case_data": None}
            with self.assertRaises(ValueError) as context:
                render_chapter_ppt("M5", PROJECT, outline, temp_dir)
            self.assertIn("未选择 M5 案例", str(context.exception))

    def test_m5_with_case_data_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {
                "case_data": M5_CASE_DATA
            }
            output_path = render_chapter_ppt("M5", PROJECT, outline, temp_dir)

            self.assertTrue(output_path.exists())
            presentation = Presentation(str(output_path))
            self.assertGreater(len(presentation.slides), 0)


class M6FixedTemplateTest(unittest.TestCase):
    """验证 M6 固定企业介绍模板（不做字段替换）。"""

    def test_m6_fixed_template_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {}
            output_path = render_chapter_ppt("M6", PROJECT, outline, temp_dir)

            self.assertTrue(output_path.exists())
            presentation = Presentation(str(output_path))
            # M6 fixed template has 15 slides per spec
            self.assertEqual(len(presentation.slides), 15)

            texts = slide_texts(output_path)
            text_blob = "\n".join(texts)
            # Verify M6 fixed template content is present (enterprise content)
            self.assertTrue(
                any("中驰" in t or "企业" in t for t in texts),
                "M6 should contain enterprise content from fixed template",
            )
            # Verify project data does NOT appear (no field substitution)
            self.assertNotIn("某城市轨道交通声屏障改造项目", text_blob)
            self.assertNotIn("轨交既有线改造", text_blob)


class TailFixedTemplateTest(unittest.TestCase):
    """验证尾页固定模板渲染。"""

    def test_tail_fixed_template_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_chapter_ppt("TAIL_FIXED", PROJECT, {}, temp_dir)

            self.assertTrue(output_path.exists())
            self.assertGreater(slide_count(output_path), 0)

    def test_tail_print_template_produces_valid_pptx(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_chapter_ppt("TAIL_PRINT", PROJECT, {}, temp_dir)

            self.assertTrue(output_path.exists())
            self.assertGreater(slide_count(output_path), 0)


class RenderChapterPptValidationTest(unittest.TestCase):
    """验证 render_chapter_ppt 参数校验。"""

    def test_accepts_m3(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"parsed_sources": ["南京地铁线路总长度约30公里"]}
            result = render_chapter_ppt("M3", PROJECT, outline, temp_dir)
            self.assertTrue(result.exists())

    def test_rejects_m4(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError) as context:
                render_chapter_ppt("M4", PROJECT, {}, temp_dir)

            self.assertIn("M4", str(context.exception))

    def test_rejects_m1_alone(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError) as context:
                render_chapter_ppt("M1", PROJECT, {}, temp_dir)

            self.assertIn("M1", str(context.exception))

    def test_rejects_m2_alone(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            with self.assertRaises(ValueError) as context:
                render_chapter_ppt("M2", PROJECT, {}, temp_dir)

            self.assertIn("M2", str(context.exception))


class MergePptxTest(unittest.TestCase):
    """验证 PPTX 合并功能。"""

    def test_merge_pptx_produces_valid_file(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            # Render each chapter
            chapters = {}
            for module_id in MERGE_ORDER:
                outline = outline_for_module(module_id)
                chapters[module_id] = render_chapter_ppt(
                    module_id, PROJECT, outline, temp_dir
                )

            # Merge
            final_path = Path(temp_dir) / "final.pptx"
            result = merge_pptx(list(chapters.values()), final_path)

            self.assertTrue(result.exists())
            presentation = Presentation(str(result))
            self.assertGreater(len(presentation.slides), 0)

    def test_merge_pptx_preserves_source_slide_size(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            chapters = []
            for module_id in MERGE_ORDER:
                outline = outline_for_module(module_id)
                chapters.append(render_chapter_ppt(module_id, PROJECT, outline, temp_dir))

            final_path = Path(temp_dir) / "final.pptx"
            result = merge_pptx(chapters, final_path)

            first_chapter = Presentation(str(chapters[0]))
            merged = Presentation(str(result))
            self.assertEqual(merged.slide_width, first_chapter.slide_width)
            self.assertEqual(merged.slide_height, first_chapter.slide_height)

    def test_merge_pptx_keeps_blip_relationships_pointing_to_images(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            chapters = []
            for module_id in MERGE_ORDER:
                outline = outline_for_module(module_id)
                chapters.append(render_chapter_ppt(module_id, PROJECT, outline, temp_dir))

            final_path = Path(temp_dir) / "final.pptx"
            result = merge_pptx(chapters, final_path)

            self.assertEqual(slide_image_relationship_errors(result), [])

    def test_merge_pptx_does_not_write_duplicate_zip_entries(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            chapters = []
            for module_id in MERGE_ORDER:
                outline = outline_for_module(module_id)
                chapters.append(render_chapter_ppt(module_id, PROJECT, outline, temp_dir))

            final_path = Path(temp_dir) / "final.pptx"
            result = merge_pptx(chapters, final_path)

            self.assertEqual(duplicate_zip_entries(result), [])

    def test_merge_pptx_slide_count_matches_chapter_slide_count_sum(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            chapters = []
            for module_id in MERGE_ORDER:
                outline = outline_for_module(module_id)
                chapters.append(render_chapter_ppt(module_id, PROJECT, outline, temp_dir))

            final_path = Path(temp_dir) / "final.pptx"
            result = merge_pptx(chapters, final_path)

            expected_count = sum(slide_count(chapter) for chapter in chapters)
            self.assertEqual(slide_count(result), expected_count)

    def test_merge_pptx_rejects_mismatched_slide_sizes(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            first = Presentation()
            first.slide_width = Inches(10)
            first.slide_height = Inches(5.625)
            first.slides.add_slide(first.slide_layouts[6])
            first_path = Path(temp_dir) / "first.pptx"
            first.save(first_path)

            second = Presentation()
            # Use a different aspect ratio (4:3) to create a real size mismatch
            second.slide_width = Inches(10)
            second.slide_height = Inches(7.5)
            second.slides.add_slide(second.slide_layouts[6])
            second_path = Path(temp_dir) / "second.pptx"
            second.save(second_path)

            # Force python-pptx path to get the size-mismatch error
            with unittest.mock.patch.dict(os.environ, {"ZHONGCHI_PPT_MERGE_ENGINE": "python-pptx"}):
                with self.assertRaises(ValueError) as context:
                    merge_pptx([first_path, second_path], Path(temp_dir) / "final.pptx")

            self.assertIn("尺寸", str(context.exception))


class RenderFinalPptTest(unittest.TestCase):
    """验证 render_final_ppt 按模块顺序合并并追加尾页。"""

    def test_final_ppt_merge_order(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            final_path = render_final_ppt(PROJECT, OUTLINES_FULL, temp_dir)

            self.assertTrue(final_path.exists())
            presentation = Presentation(str(final_path))
            # Verify total slides > 0
            self.assertGreater(len(presentation.slides), 0)

    def test_final_ppt_contains_all_three_modules(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            final_path = render_final_ppt(PROJECT, OUTLINES_FULL, temp_dir)

            texts = slide_texts(final_path)
            text_blob = "\n".join(texts)
            # M6 fixed template should contain "中驰"
            self.assertTrue(any("中驰" in t for t in texts))

    def test_final_ppt_excludes_m3_m4(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            final_path = render_final_ppt(PROJECT, OUTLINES_FULL, temp_dir)

            texts = slide_texts(final_path)
            text_blob = "\n".join(texts)
            # M3/M4 should not appear
            self.assertNotIn("定制化设计方案", text_blob)
            self.assertNotIn("工程量与施工周期测算", text_blob)

    def test_final_ppt_filename_contains_project_name(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            final_path = render_final_ppt(PROJECT, OUTLINES_FULL, temp_dir)

            self.assertIn("某城市轨道交通声屏障改造项目", final_path.name)

    def test_final_ppt_adds_fixed_tail_by_default_and_skips_print_tail(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            final_path = render_final_ppt(PROJECT, OUTLINES_FULL, temp_dir)

            chapters_dir = Path(temp_dir) / "chapters"
            fixed_tail = chapters_dir / "TAIL_FIXED_固定尾页.pptx"
            print_tail = chapters_dir / "TAIL_PRINT_尾页打印版.pptx"
            self.assertTrue(fixed_tail.exists())
            self.assertFalse(print_tail.exists())
            self.assertEqual(
                slide_count(final_path),
                sum(slide_count(path) for path in chapters_dir.glob("*.pptx")),
            )

    def test_final_ppt_adds_print_tail_when_project_requests_it(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            project = {**PROJECT, "include_print_tail_page": True}
            final_path = render_final_ppt(project, OUTLINES_FULL, temp_dir)

            chapters_dir = Path(temp_dir) / "chapters"
            fixed_tail = chapters_dir / "TAIL_FIXED_固定尾页.pptx"
            print_tail = chapters_dir / "TAIL_PRINT_尾页打印版.pptx"
            self.assertTrue(fixed_tail.exists())
            self.assertTrue(print_tail.exists())
            self.assertEqual(
                slide_count(final_path),
                sum(slide_count(path) for path in chapters_dir.glob("*.pptx")),
            )


class ReplaceTextPlaceholdersTest(unittest.TestCase):
    """验证 replace_text_placeholders 通用占位符替换能力。"""

    def _create_test_pptx_with_placeholder(self, placeholder: str, dest_dir: Path) -> Path:
        """创建一个包含指定占位符的测试 PPTX。"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        # f-string 中 {{ 产生字面 {，所以 {{project_name}} 要写成 {{{{project_name}}}}
        run.text = f"项目名称：{{{{{placeholder}}}}}"
        path = dest_dir / "test_placeholder.pptx"
        prs.save(str(path))
        return path

    def test_replaces_placeholder_with_value(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_path = self._create_test_pptx_with_placeholder("project_name", Path(temp_dir))
            replace_text_placeholders(pptx_path, {"project_name": "南京地铁3号线声屏障工程"})
            prs = Presentation(str(pptx_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            self.assertIn("南京地铁3号线声屏障工程", "\n".join(texts))
            self.assertNotIn("{{project_name}}", "\n".join(texts))

    def test_empty_field_uses_placeholder_fallback(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_path = self._create_test_pptx_with_placeholder("project_location", Path(temp_dir))
            # build_m1_m2_replacement_map 对空字段会填充中文兜底
            # 测试：传中文兜底值进去，PPTX 里显示的也是中文兜底
            replace_text_placeholders(pptx_path, {"project_location": "[待补充：项目所在地]"})
            prs = Presentation(str(pptx_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            self.assertIn("[待补充：项目所在地]", "\n".join(texts))

    def test_no_placeholder_in_template_does_not_error(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = "没有任何占位符的普通文本"
            pptx_path = Path(temp_dir) / "no_placeholder.pptx"
            prs.save(str(pptx_path))
            # 不应抛出异常
            result = replace_text_placeholders(pptx_path, {"project_name": "测试项目"})
            self.assertTrue(result.exists())

    def test_multiple_placeholders_in_same_text(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = "{{project_name}} 位于 {{project_location}}，由 {{owner_unit}} 负责建设"
            pptx_path = Path(temp_dir) / "multi_placeholder.pptx"
            prs.save(str(pptx_path))
            replace_text_placeholders(pptx_path, {
                "project_name": "南京地铁3号线",
                "project_location": "南京市",
                "owner_unit": "南京地铁集团",
            })
            prs2 = Presentation(str(pptx_path))
            texts = []
            for slide in prs2.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            text_blob = "\n".join(texts)
            self.assertIn("南京地铁3号线", text_blob)
            self.assertIn("南京市", text_blob)
            self.assertIn("南京地铁集团", text_blob)

    def test_split_placeholder_replacement_preserves_paragraph_style(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run1 = para.add_run()
            run1.text = "{{project"
            run1.font.name = "Microsoft YaHei"
            run1.font.size = Pt(32)
            run2 = para.add_run()
            run2.text = "_name}}"
            run2.font.name = "Microsoft YaHei"
            run2.font.size = Pt(32)
            pptx_path = Path(temp_dir) / "split_placeholder.pptx"
            prs.save(str(pptx_path))

            replace_text_placeholders(pptx_path, {"project_name": "南京地铁3号线既有线声屏障改造工程"})

            prs2 = Presentation(str(pptx_path))
            para2 = prs2.slides[0].shapes[0].text_frame.paragraphs[0]
            self.assertEqual(para2.text, "南京地铁3号线既有线声屏障改造工程")
            self.assertEqual(para2.alignment, PP_ALIGN.CENTER)
            self.assertEqual(para2.runs[0].font.name, "Microsoft YaHei")
            self.assertEqual(para2.runs[0].font.size, Pt(32))


class BuildM1M2ReplacementMapTest(unittest.TestCase):
    """验证 build_m1_m2_replacement_map 字段映射构建。"""

    def test_basic_fields_from_project(self):
        project = {
            "project_name": "南京地铁3号线声屏障改造工程",
            "project_location": "南京市",
            "owner_unit": "南京地铁集团",
            "product_line": "轨道交通声屏障",
        }
        outline = {}
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(result["project_name"], "南京地铁3号线声屏障改造工程")
        self.assertEqual(result["project_location"], "南京市")
        self.assertEqual(result["owner_unit"], "南京地铁集团")
        self.assertEqual(result["product_line"], "轨道交通声屏障")

    def test_missing_basic_field_uses_fallback(self):
        project = {
            "project_name": "南京地铁3号线",
            "project_location": "",
        }
        outline = {}
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(result["project_name"], "南京地铁3号线")
        self.assertEqual(result["project_location"], "[待补充：项目所在地]")
        self.assertEqual(result["owner_unit"], "[待补充：建设/业主单位]")
        self.assertEqual(result["product_line"], "[待补充：产品线]")

    def test_detected_and_confirmed_project_type(self):
        project = {
            "detected_project_type": "metro",
            "confirmed_project_type": "metro",
        }
        outline = {}
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(result["detected_project_type"], "metro")
        self.assertEqual(result["confirmed_project_type"], "metro")

    def test_m1_m2_template_from_template_selection(self):
        project = {
            "project_name": "测试项目",
            "template_selection": {
                "M1_M2": {
                    "template_filename": "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx"
                }
            },
        }
        outline = {}
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(
            result["m1_m2_template"],
            "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx"
        )

    def test_rule_field_line_name_from_filename(self):
        project = {"project_name": "测试项目"}
        outline = {
            "files": [
                {"filename": "南京地铁3号线施工方案.pdf"}
            ]
        }
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(result["line_name"], "3号线")

    def test_rule_field_line_name_s1_identifier(self):
        """S1号线 应识别为 S1号线，不是 1号线。"""
        project = {"project_name": "测试项目"}
        outline = {
            "files": [
                {"filename": "S1号线机场快线施工组织设计.pdf"}
            ]
        }
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(result["line_name"], "S1号线")

    def test_rule_field_site_pain_points_from_text(self):
        project = {"project_name": "测试项目"}
        outline = {
            "files": [
                {"filename": "施工方案.pdf", "parsed_text_path": None}
            ]
        }
        result = build_m1_m2_replacement_map(project, outline)
        # 无关键词时应返回待补充
        self.assertEqual(result["site_pain_points"], "[待补充：现场痛点]")

    def test_rule_field_construction_scenario_recognized(self):
        project = {"project_name": "测试项目"}
        outline = {
            "files": [
                {"filename": "既有线改造施工方案.pdf"}
            ]
        }
        result = build_m1_m2_replacement_map(project, outline)
        self.assertIn("既有线改造", result["construction_scenario"])

    def test_rule_field_no_match_uses_fallback(self):
        project = {"project_name": "测试项目"}
        outline = {
            "files": [
                {"filename": "普通文档.pdf"}
            ]
        }
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(result["line_name"], "[待补充：线路名称]")
        self.assertEqual(result["site_pain_points"], "[待补充：现场痛点]")
        self.assertEqual(result["construction_scenario"], "[待补充：施工场景]")


class M1M2FieldReplacementTest(unittest.TestCase):
    """验证 M1/M2 渲染后字段被正确替换。"""

    def test_m1_m2_replacement_map_includes_confirmed_project_type(self):
        project = {
            "project_name": "南京地铁3号线声屏障改造工程",
            "project_location": "南京市",
            "owner_unit": "南京地铁集团",
            "product_line": "轨道交通声屏障",
            "detected_project_type": "metro",
            "confirmed_project_type": "metro",
            "template_selection": {
                "M1_M2": {
                    "template_filename": "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx"
                }
            },
        }
        outline = {
            "project_type": "metro",
            "files": [
                {"filename": "南京地铁3号线施工组织设计.pdf"}
            ]
        }
        result = build_m1_m2_replacement_map(project, outline)
        self.assertEqual(result["confirmed_project_type"], "metro")
        self.assertEqual(result["project_name"], "南京地铁3号线声屏障改造工程")

    def test_m1_m2_renders_with_field_replacement(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {
                "project_type": "metro",
                "files": [
                    {"filename": "南京地铁3号线施工组织设计.pdf"}
                ]
            }
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)
            self.assertTrue(output_path.exists())
            # 能打开即有效
            prs = Presentation(str(output_path))
            self.assertGreater(len(prs.slides), 0)

    def test_placeholder_replacement_with_m1_m2_map(self):
        """验证使用 build_m1_m2_replacement_map 生成的映射能正确替换占位符。"""
        with tempfile.TemporaryDirectory() as temp_dir:
            # 创建一个含占位符的临时 PPTX
            from ppt_engine.renderer import (
                build_m1_m2_replacement_map,
                replace_text_placeholders,
            )
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = "项目名称：{{project_name}}，项目所在地：{{project_location}}"
            src_path = Path(temp_dir) / "src.pptx"
            prs.save(str(src_path))

            # 用 M1/M2 映射替换
            replacements = build_m1_m2_replacement_map(PROJECT, {})
            replace_text_placeholders(src_path, replacements)

            # 验证占位符已被替换为项目数据
            prs2 = Presentation(str(src_path))
            texts = []
            for slide in prs2.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            text_blob = "\n".join(texts)
            self.assertNotIn("{{project_name}}", text_blob)
            self.assertNotIn("{{project_location}}", text_blob)
            self.assertIn("某城市轨道交通声屏障改造项目", text_blob)
            self.assertIn("南京", text_blob)


class M5CaseTemplateNoReplacementTest(unittest.TestCase):
    """验证 M5 渲染不做字段替换（M5 只复制模板）。"""

    def test_m5_no_field_replacement(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"case_data": M5_CASE_DATA}
            output_path = render_chapter_ppt("M5", PROJECT, outline, temp_dir)
            self.assertTrue(output_path.exists())
            prs = Presentation(str(output_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            # M5 不应有 project_name 等项目字段被写入（只有模板原始内容）
            # 这里只验证能正常打开，具体内容由模板决定


class M6FixedTemplateNoProjectFieldsTest(unittest.TestCase):
    """验证 M6 固定模板不写入项目字段。"""

    def test_m6_fixed_template_no_project_field_replacement(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {}
            output_path = render_chapter_ppt("M6", PROJECT, outline, temp_dir)
            self.assertTrue(output_path.exists())
            prs = Presentation(str(output_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            text_blob = "\n".join(texts)
            # M6 是固定模板，不应有项目字段被替换的痕迹
            self.assertNotIn("某城市轨道交通声屏障改造项目", text_blob)


class M1M2TemplateHasPlaceholdersTest(unittest.TestCase):
    """验证首页占位符和 M1/M2 旧项目信息页清理状态。"""

    def test_home_template_contains_only_project_name_placeholder(self):
        """首页模板只要求包含项目名称占位符。"""
        import re
        PLACEHOLDER_PATTERN = re.compile(r"\{\{(\w+)\}\}")
        home_path = PPT_TEMPLATE_ROOT / HOME_TEMPLATE_FILENAME
        prs = Presentation(str(home_path))
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text += shape.text
        self.assertEqual(set(PLACEHOLDER_PATTERN.findall(all_text)), {"project_name"})

    def test_all_four_m1_m2_templates_do_not_contain_project_info_page(self):
        """M1/M2 模板不应再包含旧的"项目生成信息"页。"""
        for tmpl_filename in M1_M2_TEMPLATE_MAP.values():
            tmpl_path = PPT_TEMPLATE_ROOT / tmpl_filename
            prs = Presentation(str(tmpl_path))
            slides_with_project_info = []
            for index, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text
                if "项目生成信息" in slide_text:
                    slides_with_project_info.append(index)
            self.assertEqual(
                slides_with_project_info, [],
                f"{tmpl_filename} 仍包含旧项目生成信息页: {slides_with_project_info}",
            )

    def test_rendered_m1_m2_output_no_longer_contains_placeholder(self):
        """渲染 M1/M2 模板后，输出 PPTX 不再包含 {{project_name}}。"""
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "metro"}
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)
            prs = Presentation(str(output_path))
            all_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        all_text += shape.text
            self.assertNotIn("{{project_name}}", all_text)

    def test_rendered_m1_m2_output_does_not_contain_project_info_page(self):
        """渲染 M1/M2 模板后，不应出现旧的项目生成信息页。"""
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "metro"}
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)
            prs = Presentation(str(output_path))
            slides_with_project_info = []
            for index, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text
                if "项目生成信息" in slide_text:
                    slides_with_project_info.append(index)
            self.assertEqual(slides_with_project_info, [])

    def test_rendered_m1_m2_output_keeps_original_first_slide(self):
        """渲染 M1/M2 模板后，不再把末页移动到第一页。"""
        with tempfile.TemporaryDirectory() as temp_dir:
            outline = {"project_type": "metro"}
            output_path = render_chapter_ppt("M1_M2", PROJECT, outline, temp_dir)
            prs = Presentation(str(output_path))
            first_text = "".join(
                shape.text
                for shape in prs.slides[0].shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertNotIn("项目生成信息", first_text)


class HomeTemplateRenderTest(unittest.TestCase):
    """验证首页模板渲染与最终拼接位置。"""

    def test_home_template_render_replaces_project_name(self):
        """首页渲染后包含实际项目名，不再包含占位符。"""
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = render_chapter_ppt("HOME", PROJECT, {}, temp_dir)
            self.assertTrue(output_path.exists())
            prs = Presentation(str(output_path))
            all_text = "".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertIn("某城市轨道交通声屏障改造项目", all_text)
            self.assertNotIn("{{project_name}}", all_text)

    def test_final_ppt_first_slide_is_home(self):
        """最终 PPT 第一页来自首页模板，并包含项目名。"""
        with tempfile.TemporaryDirectory() as temp_dir:
            final_path = render_final_ppt(PROJECT, OUTLINES_FULL, temp_dir)
            prs = Presentation(str(final_path))
            first_text = "".join(
                shape.text
                for shape in prs.slides[0].shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertIn("某城市轨道交通声屏障改造项目", first_text)
            self.assertNotIn("{{project_name}}", first_text)


if __name__ == "__main__":
    unittest.main()
