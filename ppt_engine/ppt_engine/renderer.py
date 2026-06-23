import os
import re
import shutil
import zipfile
from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

# =============================================================================
# 模板路径配置（统一入口，禁止散落硬编码）
# =============================================================================
PPT_TEMPLATE_ROOT = Path(
    os.environ.get(
        "ZHONGCHI_PPT_TEMPLATE_ROOT",
        Path(__file__).resolve().parents[1] / "templates" / "solution_fixed_modules",
    )
)

# M1/M2 固定模板映射表：project_type -> 模板文件名
M1_M2_TEMPLATE_MAP = {
    "highway": "公路全封闭声屏障（M1_&_M2）.pptx",
    "metro": "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx",
    "existing_rail_transit": "铁路_&_轨道交通既有线声屏障_（M1_&_M2）.pptx",
    "railway": "铁路声屏障行业背景与技术发展（M1_&_M2）.pptx",
}

# M5 演示案例模板文件名
M5_TEMPLATE_FILENAME = "M5示例.pptx"

# M6 固定模板文件名
M6_TEMPLATE_FILENAME = "中驰企业介绍合并初版（M6）.pptx"

# 固定尾页模板文件名
TAIL_FIXED_TEMPLATE_FILENAME = "尾页.pptx"

# 尾页打印版模板文件名
TAIL_PRINT_TEMPLATE_FILENAME = "尾页-打印版.pptx"

# M3 固定模板文件名
M3_TEMPLATE_FILENAME = "M3_项目深化方案模板.pptx"

# 首页固定模板文件名
HOME_TEMPLATE_FILENAME = "首页.pptx"

# M3 完整测试模板文件名。仅用于功能测试，不接入正式完整流程。
M3_FULL_TEST_TEMPLATE_FILENAME = "M3_项目深化方案模板_最终版本.pptx"

M3_FULL_SECTIONS = [
    {"title": "项目基本情况", "text_field": "m3_basic_summary", "image_field": "image:m3_basic"},
    {"title": "项目线路图", "text_field": "m3_line_summary", "image_field": "image:m3_line"},
    {"title": "敏感点路段", "text_field": "m3_sensitive_points_summary", "image_field": "image:m3_sensitive_points"},
    {"title": "工程量统计", "text_field": "m3_quantity_summary", "image_field": "image:m3_quantity"},
    {"title": "结构形式", "text_field": "m3_structure_summary", "image_field": "image:m3_structure"},
    {"title": "现场踏勘", "text_field": "m3_site_survey_summary", "image_field": "image:m3_site_survey"},
    {"title": "现场勘察情况", "text_field": "m3_investigation_summary", "image_field": "image:m3_investigation"},
    {"title": "项目重难点分析", "text_field": "m3_risk_summary", "image_field": "image:m3_risk"},
    {"title": "重难点应对措施", "text_field": "m3_solution_summary", "image_field": "image:m3_solution"},
]

M3_FULL_IMAGE_FIELDS = {section["image_field"] for section in M3_FULL_SECTIONS}
M3_FULL_TEXT_FIELDS = {section["text_field"] for section in M3_FULL_SECTIONS}
M3_FULL_TABLE_MAX_ROWS = 60
M3_FULL_TABLE_MAX_COLUMNS = 12

# M3 占位符字段名列表（9个字段）
M3_FIELD_NAMES = (
    "m3_basic_summary",
    "m3_line_summary",
    "m3_sensitive_points_summary",
    "m3_quantity_summary",
    "m3_structure_summary",
    "m3_site_survey_summary",
    "m3_investigation_summary",
    "m3_risk_summary",
    "m3_solution_summary",
)

M3_ACTIVE_FIELD_NAMES = (
    "m3_basic_summary",
    "m3_quantity_summary",
    "m3_site_survey_summary",
    "m3_risk_summary",
    "m3_solution_summary",
)


# =============================================================================
# 合并顺序与模块标题
# =============================================================================
# 最终 PPT 合并顺序：首页 -> M1/M2 -> M3 -> M5 -> M6 -> 固定尾页 -> 打印版尾页
MERGE_ORDER = ("HOME", "M1_M2", "M3", "M5", "M6", "TAIL_FIXED", "TAIL_PRINT")

MODULE_TITLES = {
    "HOME": "首页",
    "M1_M2": "行业背景与技术标准",
    "M3": "项目深化方案",
    "M5": "同类型案例匹配",
    "M6": "企业背书与荣誉",
    "TAIL_FIXED": "固定尾页",
    "TAIL_PRINT": "尾页打印版",
}

# M1/M2 模板的 project_type 枚举
PROJECT_TYPES = ("highway", "metro", "existing_rail_transit", "railway")

PRIMARY_BLUE = RGBColor(0x15, 0x65, 0xC0)
DEEP_BLUE = RGBColor(0x1A, 0x3A, 0x6B)
ACCENT_GOLD = RGBColor(0xCA, 0x8A, 0x04)
LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)
BODY_TEXT = RGBColor(0x1F, 0x29, 0x37)
EMU_PER_POINT = 12700

# 占位符正则（用于识别未知占位符并替换为中文兜底）
PLACEHOLDER_PATTERN = re.compile(r"\{\{(\w+)\}\}")

# 规则识别候选占位符列表
RULE_FIELD_NAMES = ("line_name", "site_pain_points", "construction_scenario")

# 字段中文标签映射（用于兜底占位符显示）
FIELD_LABELS: dict[str, str] = {
    "project_name": "项目名称",
    "project_location": "项目所在地",
    "owner_unit": "建设/业主单位",
    "product_line": "产品线",
    "detected_project_type": "系统识别项目类型",
    "confirmed_project_type": "人工确认项目类型",
    "m1_m2_template": "M1/M2模板",
    "line_name": "线路名称",
    "site_pain_points": "现场痛点",
    "construction_scenario": "施工场景",
}


# =============================================================================
# 辅助函数
# =============================================================================


def _missing_placeholder(field_name: str) -> str:
    """返回字段缺失时的中文兜底占位符。"""
    label = FIELD_LABELS.get(field_name, field_name)
    return f"[待补充：{label}]"


def _validate_project_type(project_type: str) -> None:
    """验证 project_type 是否为合法枚举值。"""
    if project_type not in PROJECT_TYPES:
        raise ValueError(
            f"无效的 project_type：{project_type}，可选值：{PROJECT_TYPES}。"
        )


def _validate_module(module_id: str) -> None:
    """验证 module_id 是否为合法模块。"""
    if module_id not in MERGE_ORDER:
        raise ValueError("本阶段 PPT 引擎只支持 HOME/M1_M2/M3/M5/M6/TAIL_FIXED/TAIL_PRINT，M4 暂未接入。")


def _safe_text(value: Any, field_name: str) -> str:
    if value is None or value == "":
        return _missing_placeholder(field_name)
    return str(value)


def _replace_placeholders_in_text(text: str, replacements: dict[str, str]) -> str:
    """替换字符串中的已知和未知占位符。"""
    new_text = text
    for placeholder, value in replacements.items():
        pattern = f"{{{{{placeholder}}}}}"
        if pattern in new_text:
            new_text = new_text.replace(pattern, value)
    for match in PLACEHOLDER_PATTERN.finditer(new_text):
        field_name = match.group(1)
        if field_name not in replacements:
            new_text = new_text.replace(match.group(0), _missing_placeholder(field_name))
    return new_text


# =============================================================================
# 规则识别字段提取
# =============================================================================


def _extract_line_name(text: str) -> str:
    """从文本中提取线路名称。

    匹配模式：
    - \\d+号线，如"3号线"、"10号线"
    - S\\d+号线，如"S1号线"
    - 昌平线
    - K数字+数字~K数字+数字，如"K12+100~K15+200"
    """
    patterns = [
        r'(S\d+)号线',
        r'(\d+)号线',
        r'昌平线',
        r'(K\d+\+\d+~K\d+\+\d+)',
    ]
    for pattern in patterns:
        match = re.search(pattern, text)
        if match:
            return match.group(0)
    return ""


def _extract_site_pain_points(text: str) -> str:
    """从文本中归纳现场痛点。

    关键词映射：
    - 噪声、噪音、降噪 -> 噪声治理
    - 夜间、天窗、施工窗口、短窗口 -> 施工窗口受限
    - 工期、赶工、交付 -> 工期紧张
    - 改造、既有线、运营线 -> 既有线改造约束
    """
    pain_points: list[str] = []
    lower_text = text.lower()

    if any(kw in lower_text for kw in ['噪声', '噪音', '降噪']):
        pain_points.append("噪声治理")
    if any(kw in lower_text for kw in ['夜间', '天窗', '施工窗口', '短窗口']):
        pain_points.append("施工窗口受限")
    if any(kw in lower_text for kw in ['工期', '赶工', '交付']):
        pain_points.append("工期紧张")
    if any(kw in lower_text for kw in ['改造', '既有线', '运营线']):
        pain_points.append("既有线改造约束")

    return "、".join(pain_points) if pain_points else ""


def _extract_construction_scenario(text: str) -> str:
    """从文本中判断施工场景。

    关键词映射：
    - 既有线、运营线、改造 -> 既有线改造
    - 加装、增设、新增 -> 加装工程
    - 新建 -> 新建工程
    - 高速、公路 -> 公路声屏障
    - 地铁、轨道交通、轨交 -> 轨道交通声屏障
    - 铁路 -> 铁路声屏障
    """
    scenarios: list[str] = []
    lower_text = text.lower()

    if any(kw in lower_text for kw in ['既有线', '运营线', '改造']):
        scenarios.append("既有线改造")
    if any(kw in lower_text for kw in ['加装', '增设', '新增']):
        scenarios.append("加装工程")
    if '新建' in lower_text:
        scenarios.append("新建工程")
    if any(kw in lower_text for kw in ['高速', '公路']) and '声屏障' in lower_text:
        scenarios.append("公路声屏障")
    if any(kw in lower_text for kw in ['地铁', '轨道交通', '轨交']) and '声屏障' in lower_text:
        scenarios.append("轨道交通声屏障")
    if '铁路' in lower_text and '声屏障' in lower_text:
        scenarios.append("铁路声屏障")

    # 去重，保持顺序
    seen = set()
    unique = []
    for s in scenarios:
        if s not in seen:
            seen.add(s)
            unique.append(s)
    return "、".join(unique) if unique else ""


def _rule_based_field_value(field_name: str, project: dict[str, Any], outline: dict[str, Any]) -> str:
    """对规则识别字段（line_name/site_pain_points/construction_scenario）进行提取。

    优先从 project.project_name / project_location / owner_unit / product_line 提取，
    其次从 outline（classification_result.files）中的文件名和解析文本提取。
    规则识别不到则返回空字符串，由调用方统一填充 [待补充：字段名]。
    """
    # 收集候选文本
    candidates: list[str] = []
    for key in ("project_name", "project_location", "owner_unit", "product_line"):
        val = project.get(key)
        if val:
            candidates.append(str(val))

    # 从 outline 中获取文件信息（优先）
    files = outline.get("files", []) if isinstance(outline, dict) else []
    # 同时从 project.classification_result.files 获取（兼容旧流程）
    if not files:
        files = project.get("classification_result", {}).get("files", [])
    for f in files:
        fn = f.get("filename", "")
        if fn:
            candidates.append(fn)
        parsed_text_path = f.get("parsed_text_path", "")
        if parsed_text_path and Path(parsed_text_path).exists():
            try:
                candidates.append(Path(parsed_text_path).read_text(encoding="utf-8", errors="ignore"))
            except OSError:
                pass

    combined = " ".join(candidates)

    if field_name == "line_name":
        return _extract_line_name(combined)
    if field_name == "site_pain_points":
        return _extract_site_pain_points(combined)
    if field_name == "construction_scenario":
        return _extract_construction_scenario(combined)
    return ""


# =============================================================================
# M1/M2 字段替换映射构建
# =============================================================================


def build_m1_m2_replacement_map(
    project: dict[str, Any],
    outline: dict[str, Any],
) -> dict[str, str]:
    """构建 M1/M2 模板字段替换映射。

    字段值来源优先级：
    1. 项目基础信息：project_name、project_location、owner_unit、product_line
    2. 系统识别和人工确认：detected_project_type、confirmed_project_type、template_selection
    3. 规则识别字段：line_name、site_pain_points、construction_scenario
    4. 兜底：[待补充：字段名]

    Args:
        project: 项目基础信息字典。
        outline: 章节配置字典，可能包含 files（用于规则识别）。

    Returns:
        占位符 -> 替换值的字典。
    """
    replacements: dict[str, str] = {}

    # 1. 项目基础信息字段
    basic_fields = [
        ("project_name", project.get("project_name")),
        ("project_location", project.get("project_location")),
        ("owner_unit", project.get("owner_unit")),
        ("product_line", project.get("product_line")),
    ]
    for field_name, value in basic_fields:
        if value is not None and value != "":
            replacements[field_name] = str(value)
        else:
            replacements[field_name] = _missing_placeholder(field_name)

    # 2. 系统识别和人工确认字段
    replacements["detected_project_type"] = _safe_text(
        project.get("detected_project_type"), "系统识别项目类型"
    )
    replacements["confirmed_project_type"] = _safe_text(
        project.get("confirmed_project_type"), "人工确认项目类型"
    )

    # 模板文件名
    template_selection = project.get("template_selection", {})
    m1_m2_template = template_selection.get("M1_M2", {}).get("template_filename", "")
    replacements["m1_m2_template"] = m1_m2_template or _missing_placeholder("m1_m2_template")

    # 3. 规则识别字段（line_name、site_pain_points、construction_scenario）
    for field_name in RULE_FIELD_NAMES:
        value = _rule_based_field_value(field_name, project, outline)
        if value:
            replacements[field_name] = value
        else:
            replacements[field_name] = _missing_placeholder(field_name)

    return replacements


def _set_run(run, size: int, color: RGBColor = BODY_TEXT, bold: bool = False) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _add_textbox(slide, left, top, width, height, text: str, size: int, color: RGBColor = BODY_TEXT, bold: bool = False):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    _set_run(run, size=size, color=color, bold=bold)
    return shape


def _add_header(slide, module_id: str, page_label: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.55), Inches(0.35), Inches(7.2), Inches(0.3), "中驰智能PPT Demo", 12, DEEP_BLUE, True)
    _add_textbox(slide, Inches(10.4), Inches(0.35), Inches(2.3), Inches(0.3), f"{module_id} | {page_label}", 11, ACCENT_GOLD, True)


def _add_cover_slide(prs: Presentation, module_id: str, project: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, module_id, "章节封面")

    _add_textbox(slide, Inches(0.9), Inches(1.35), Inches(9.2), Inches(0.7), MODULE_TITLES[module_id], 32, DEEP_BLUE, True)
    _add_textbox(
        slide,
        Inches(0.95),
        Inches(2.25),
        Inches(8.8),
        Inches(0.55),
        _safe_text(project.get("project_name"), "项目名称"),
        20,
        BODY_TEXT,
        False,
    )
    _add_textbox(
        slide,
        Inches(0.95),
        Inches(3.05),
        Inches(10.8),
        Inches(0.45),
        f"项目所在地：{_safe_text(project.get('project_location'), '项目所在地')}    产品线：{_safe_text(project.get('product_line'), '产品线')}",
        15,
        BODY_TEXT,
        False,
    )

    accent = slide.shapes.add_shape(1, Inches(0.95), Inches(4.4), Inches(2.1), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GOLD
    accent.line.fill.background()


def _add_outline_slide(prs: Presentation, module_id: str, slide_data: dict[str, Any], page_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, module_id, f"第 {page_number} 页")

    _add_textbox(
        slide,
        Inches(0.75),
        Inches(0.95),
        Inches(11.6),
        Inches(0.55),
        _safe_text(slide_data.get("page_title"), "页面标题"),
        25,
        DEEP_BLUE,
        True,
    )
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(1.75),
        Inches(11.2),
        Inches(0.65),
        _safe_text(slide_data.get("core_insight"), "核心观点"),
        17,
        BODY_TEXT,
        True,
    )

    bullet_points = slide_data.get("bullet_points") or ["[待补充：本页要点]"]
    bullet_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.75), Inches(10.9), Inches(2.6))
    frame = bullet_box.text_frame
    frame.clear()
    for index, bullet in enumerate(bullet_points[:5]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(bullet)
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = BODY_TEXT

    missing_fields = slide_data.get("missing_fields") or []
    if missing_fields:
        marker_text = "  ".join(f"[待补充：{field}]" for field in missing_fields)
        _add_textbox(slide, Inches(0.85), Inches(6.0), Inches(11.5), Inches(0.4), marker_text, 13, ACCENT_GOLD, True)


# =============================================================================
# 通用占位符替换
# =============================================================================


def replace_text_placeholders(pptx_path: str | Path, replacements: dict[str, str]) -> Path:
    """对 PPTX 中所有文本占位符进行替换。

    支持任意 PPTX 中的 {{placeholder}} 形式占位符替换。
    - 模板中没有占位符时不报错
    - 字段为空时使用 [待补充：字段名] 兜底（由调用方保证）
    - 替换后 PPTX 仍能被 python-pptx 和 PowerPoint 打开
    - 不处理图片替换
    - 支持普通文本框中完整出现的占位符（如 {{project_name}}）

    Args:
        pptx_path: 源 PPTX 文件路径（会被原地修改并保存到同名路径）。
        replacements: 占位符 -> 替换值的字典。

    Returns:
        替换后的 PPTX 文件路径（与输入相同）。
    """
    path = Path(pptx_path)
    prs = Presentation(str(path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            changed_in_runs = False
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    new_text = _replace_placeholders_in_text(original_text, replacements)
                    if new_text != original_text:
                        run.text = new_text
                        changed_in_runs = True
                if not changed_in_runs and "{{" in paragraph.text and paragraph.runs:
                    original_text = paragraph.text
                    new_text = _replace_placeholders_in_text(original_text, replacements)
                    if new_text != original_text:
                        paragraph.runs[0].text = new_text
                        for run in paragraph.runs[1:]:
                            run.text = ""

    prs.save(str(path))
    return path


def _copy_fixed_template(source_path: Path, dest_path: Path) -> None:
    """Copy fixed template PPTX to destination. Creates parent directories if needed."""
    dest_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_path, dest_path)


def _render_m6_fixed_template(project: dict[str, Any], output_dir: Path) -> Path:
    """渲染 M6 固定企业介绍模板（不做字段替换）。

    直接复制固定模板到输出目录，不做任何字段替换或素材替换。
    模板缺失时抛出异常，不生成静默占位页。
    """
    source_template = PPT_TEMPLATE_ROOT / M6_TEMPLATE_FILENAME
    if not source_template.exists():
        raise FileNotFoundError(
            f"M6 固定模板未找到：{source_template}。"
            "请确认 PPT_TEMPLATE_ROOT 配置正确且模板文件存在。"
        )

    file_path = output_dir / f"M6_{MODULE_TITLES['M6']}.pptx"
    _copy_fixed_template(source_template, file_path)
    return file_path


def _render_home_fixed_template(project: dict[str, Any], output_dir: Path) -> Path:
    """渲染首页固定模板，只替换项目名称占位符。"""
    source_template = PPT_TEMPLATE_ROOT / HOME_TEMPLATE_FILENAME
    if not source_template.exists():
        raise FileNotFoundError(
            f"首页固定模板未找到：{source_template}。"
            "请确认 PPT_TEMPLATE_ROOT 配置正确且模板文件存在。"
        )

    file_path = output_dir / f"HOME_{MODULE_TITLES['HOME']}.pptx"
    _copy_fixed_template(source_template, file_path)
    replace_text_placeholders(file_path, {
        "project_name": _safe_text(project.get("project_name"), "项目名称"),
    })
    return file_path


def _render_tail_fixed_template(output_dir: Path) -> Path:
    """渲染固定尾页模板（不做字段替换）。"""
    source_template = PPT_TEMPLATE_ROOT / TAIL_FIXED_TEMPLATE_FILENAME
    if not source_template.exists():
        raise FileNotFoundError(
            f"固定尾页模板未找到：{source_template}。"
            "请确认 PPT_TEMPLATE_ROOT 配置正确且模板文件存在。"
        )

    file_path = output_dir / f"TAIL_FIXED_{MODULE_TITLES['TAIL_FIXED']}.pptx"
    _copy_fixed_template(source_template, file_path)
    return file_path


def _render_tail_print_template(output_dir: Path) -> Path:
    """渲染尾页打印版模板（不做字段替换）。"""
    source_template = PPT_TEMPLATE_ROOT / TAIL_PRINT_TEMPLATE_FILENAME
    if not source_template.exists():
        raise FileNotFoundError(
            f"尾页打印版模板未找到：{source_template}。"
            "请确认 PPT_TEMPLATE_ROOT 配置正确且模板文件存在。"
        )

    file_path = output_dir / f"TAIL_PRINT_{MODULE_TITLES['TAIL_PRINT']}.pptx"
    _copy_fixed_template(source_template, file_path)
    return file_path


def _render_m1_m2_fixed(
    project_type: str, project: dict[str, Any], outline: dict[str, Any], output_dir: Path
) -> Path:
    """渲染 M1/M2 固定模板并执行字段替换。

    根据 confirmed_project_type 从 M1_M2_TEMPLATE_MAP 选择对应模板，
    复制模板文件后执行 M1/M2 字段替换。
    模板缺失时抛出异常，不生成静默占位页。
    """
    _validate_project_type(project_type)
    template_filename = M1_M2_TEMPLATE_MAP[project_type]
    source_template = PPT_TEMPLATE_ROOT / template_filename

    if not source_template.exists():
        raise FileNotFoundError(
            f"M1/M2 固定模板未找到：{source_template}。"
            f"请确认 project_type={project_type} 对应的模板文件存在。"
        )

    file_path = output_dir / f"M1_M2_{MODULE_TITLES['M1_M2']}.pptx"
    _copy_fixed_template(source_template, file_path)

    # 执行 M1/M2 字段替换
    replacements = build_m1_m2_replacement_map(project, outline)
    replace_text_placeholders(file_path, replacements)
    return file_path


def _render_m5_case_template(
    case_data: dict[str, Any] | None, project: dict[str, Any], output_dir: Path
) -> Path:
    """渲染 M5 案例模板。

    选择规则：
    - case_data.source_type == "fixed_m5" 且 source_path 存在 → 使用该 source_path
    - source_path 不存在或路径不在 M5 目录内 → 抛出清晰错误
    - 非 fixed_m5（m5_demo、case_id=1 等）→ 使用 M5示例.pptx
    """
    source_template: Path

    if case_data and case_data.get("source_type") == "fixed_m5":
        raw_path = Path(case_data["source_path"]).resolve()
        m5_dir = (PPT_TEMPLATE_ROOT / "M5").resolve()
        # 路径安全校验：必须在 M5 目录内且后缀是 .pptx
        if not str(raw_path).startswith(str(m5_dir) + os.sep) and raw_path != m5_dir:
            raise ValueError(
                f"fixed_m5 source_path 超出 M5 目录范围：{raw_path}。"
                f"预期在 {m5_dir} 内。"
            )
        if raw_path.suffix.lower() != ".pptx":
            raise ValueError(
                f"fixed_m5 source_path 后缀不是 .pptx：{raw_path}。"
            )
        if not raw_path.exists():
            raise FileNotFoundError(
                f"fixed_m5 案例文件不存在：{raw_path}。"
                "请确认 M5 文件夹中该案例文件存在。"
            )
        source_template = raw_path
    else:
        # 旧兼容：m5_demo、case_id=1、非 fixed_m5_case 等
        source_template = PPT_TEMPLATE_ROOT / M5_TEMPLATE_FILENAME
        if not source_template.exists():
            raise FileNotFoundError(
                f"M5 案例模板未找到：{source_template}。"
                "请确认 PPT_TEMPLATE_ROOT 配置正确且模板文件存在。"
            )

    file_path = output_dir / f"M5_{MODULE_TITLES['M5']}.pptx"
    _copy_fixed_template(source_template, file_path)
    return file_path


def _render_m3_from_template(
    project: dict[str, Any], parsed_sources: list[str] | None, output_path: Path
) -> Path:
    """复制 M3 模板并执行 active 文本占位符替换。"""
    source_template = PPT_TEMPLATE_ROOT / M3_TEMPLATE_FILENAME
    if not source_template.exists():
        raise FileNotFoundError(f"M3 模板未找到：{source_template}。")

    missing = validate_m3_template_placeholders(source_template)
    if missing:
        raise ValueError(
            f"M3 模板缺少占位符：{missing}。"
            "请先运行 scripts/inject_m3_placeholders.py 补充占位符。"
        )

    _copy_fixed_template(source_template, output_path)
    replacements = build_m3_replacement_map(project, parsed_sources or [])
    replace_text_placeholders(output_path, replacements)
    return output_path


def render_m3_module(
    project: dict[str, Any], parsed_sources: list[str] | None, output_dir: str | Path
) -> Path:
    """渲染正式流程 M3 项目深化方案 PPTX。"""
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    dest_path = output_path / "M3_项目深化方案.pptx"
    return _render_m3_from_template(project, parsed_sources, dest_path)


def render_chapter_ppt(
    module_id: str,
    project: dict[str, Any],
    outline: dict[str, Any],
    output_dir: str | Path,
) -> Path:
    """渲染单个章节 PPTX。

    Args:
        module_id: 模块标识，支持 HOME / M1_M2 / M3 / M5 / M6 / TAIL_FIXED / TAIL_PRINT。
        project: 项目基础信息字典。
        outline: 章节配置。对于 M1_M2，outline 应包含 project_type；对于 M3，可包含 parsed_sources；
                 对于 M5，应包含 case_data；对于 M6，outline 被忽略（固定模板）。
        output_dir: 输出目录路径。

    Returns:
        生成的章节 PPTX 文件路径。
    """
    _validate_module(module_id)
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    if module_id == "HOME":
        # HOME：首页固定模板，只替换项目名称
        return _render_home_fixed_template(project, output_path)

    if module_id == "M1_M2":
        # M1/M2：根据 project_type 选择固定模板并执行字段替换
        project_type = outline.get("project_type", "highway")
        return _render_m1_m2_fixed(project_type, project, outline, output_path)

    if module_id == "M5":
        # M5：案例模板填充入口，必须有有效的 case_data 才能渲染
        case_data = outline.get("case_data")
        if not case_data or not case_data.get("case_id"):
            raise ValueError("未选择 M5 案例，禁止生成 M5 案例模板。")
        return _render_m5_case_template(case_data, project, output_path)

    if module_id == "M3":
        m3_materials = outline.get("m3_materials")
        if m3_materials:
            return render_m3_full_project_ppt(
                project,
                m3_materials.get("texts") or {},
                m3_materials.get("images_by_purpose") or {},
                output_path,
                m3_materials.get("page_texts") or None,
                m3_materials.get("tables_by_purpose") or None,
            )
        return render_m3_module(project, outline.get("parsed_sources", []), output_path)

    if module_id == "M6":
        # M6：直接使用固定企业介绍模板，不做字段替换
        return _render_m6_fixed_template(project, output_path)

    if module_id == "TAIL_FIXED":
        return _render_tail_fixed_template(output_path)

    if module_id == "TAIL_PRINT":
        return _render_tail_print_template(output_path)

    # 不应到达此处（_validate_module 已拦截）
    raise ValueError(f"不支持的模块：{module_id}")


def _scale_shape_xml(element, scale_x: float, scale_y: float) -> None:
    """按目标画布比例缩放复制出来的 shape XML。"""
    if scale_x == 1 and scale_y == 1:
        return
    for node in element.iter():
        tag = str(node.tag)
        if tag.endswith("}off"):
            if "x" in node.attrib:
                node.set("x", str(round(int(node.attrib["x"]) * scale_x)))
            if "y" in node.attrib:
                node.set("y", str(round(int(node.attrib["y"]) * scale_y)))
        elif tag.endswith("}ext"):
            if "cx" in node.attrib:
                node.set("cx", str(round(int(node.attrib["cx"]) * scale_x)))
            if "cy" in node.attrib:
                node.set("cy", str(round(int(node.attrib["cy"]) * scale_y)))


def _append_slide(target: Presentation, source_slide, scale_x: float = 1, scale_y: float = 1) -> None:
    """将 source_slide 完整追加到 target，包括其 OPC 关系（图片/媒体/图表等）。

    策略：直接操作底层 XML，将源 slide 的 <p:cSld>（含背景、色图映射、全部 shape）
    整体注入目标 slide，保留完整的母版引用和背景继承关系。
    仅在 python-pptx API 可完整保留来源视觉效果时才使用其高级 API，
    否则降级为直接 XML 注入以确保母版/背景/图片关系不丢失。
    """
    target_slide = target.slides.add_slide(target.slide_layouts[6])
    source_element = source_slide.element

    # ── 1. 复制图片关系（先于 XML 注入，以便新 rId 可立即用于注入的 XML）──
    relationship_id_map: dict[str, str] = {}
    target_slide_part = target_slide.part

    if hasattr(source_slide.part, "rels"):
        for rId, rel in source_slide.part.rels.items():
            if rel.reltype.endswith("/relationships/image") and not rel.is_external:
                # 从源 slide part 的 blob（二进制内容）获取图片数据，加入目标
                try:
                    blob = rel.target_part.blob
                    _, new_rId = target_slide_part.get_or_add_image_part(BytesIO(blob))
                    relationship_id_map[rId] = new_rId
                except Exception:
                    # 图片复制失败：保留旧 rId，让 PowerPoint 尝试修复
                    pass

    # ── 2. 提取源 slide 的 <p:cSld> 并注入目标 slide ─────────────────────────
    # 找到 <p:cSld> 元素（所有视觉内容的根节点）
    cSld = source_element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
    )
    if cSld is None:
        cSld = source_element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")

    if cSld is not None:
        # 深度复制 cSld（包含 <p:spTree> 全部 shape 和 <p:bg> 背景节点）
        new_cSld = deepcopy(cSld)

        # 2a. 缩放所有 shape 的坐标和尺寸
        _scale_shape_xml(new_cSld, scale_x, scale_y)

        # 2b. 重映射所有 rId（embed/link 等属性）
        # 扫描 new_cSld 树，将旧 rId 替换为 relationship_id_map 中的新 rId；
        # 若某 rId 未被映射（关系丢失），则移除该属性（避免 dangling ref）
        for elem in new_cSld.iter():
            for attr_name, attr_val in list(elem.attrib.items()):
                if isinstance(attr_val, str) and attr_val.startswith("rId"):
                    if attr_val in relationship_id_map:
                        elem.set(attr_name, relationship_id_map[attr_val])
                    # 未映射的非图片 rId（如 slideLayout/slideMaster 引用）保留原值，
                    # python-pptx 会自动解析为相对于目标 presentation 的路径。

        # 2c. 找到目标 slide 的 <p:cSld> 并替换（保留 slide 壳层，仅替换内容）
        target_cSld = target_slide.element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
        )
        if target_cSld is None:
            target_cSld = target_slide.element.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
            )
        if target_cSld is not None:
            parent = target_cSld.getparent()
            if parent is not None:
                idx = list(parent).index(target_cSld) if target_cSld in parent else -1
                parent.remove(target_cSld)
                parent.insert(idx, new_cSld)


def _presentation_size(path: str | Path) -> tuple[int, int]:
    presentation = Presentation(str(path))
    return int(presentation.slide_width), int(presentation.slide_height)


def _validate_chapter_sizes(chapter_paths: list[str | Path]) -> tuple[int, int]:
    first_width, first_height = _presentation_size(chapter_paths[0])
    first_ratio = first_width / first_height
    for chapter_path in chapter_paths[1:]:
        width, height = _presentation_size(chapter_path)
        ratio = width / height
        if abs(ratio - first_ratio) > 0.001:
            raise ValueError(
                "章节 PPTX 宽高比不一致，当前合并器暂不支持跨比例合并："
                f"{chapter_path} 的尺寸为 {width}x{height}，"
                f"首个章节尺寸为 {first_width}x{first_height}。"
            )
    return first_width, first_height


def _merge_pptx_with_powerpoint(chapter_paths: list[str | Path], output_path: str | Path) -> bool:
    """使用 PowerPoint COM 原样合并幻灯片；不可用时返回 False。"""
    if os.name != "nt":
        return False

    try:
        import win32com.client  # type: ignore[import-not-found]
    except ImportError:
        return False

    final_path = Path(output_path).resolve()
    final_path.parent.mkdir(parents=True, exist_ok=True)

    app = None
    target = None
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        target = app.Presentations.Add(WithWindow=False)

        # PowerPoint COM 使用 points；12700 EMU = 1 point。
        first_width, first_height = _presentation_size(chapter_paths[0])
        target.PageSetup.SlideWidth = first_width / EMU_PER_POINT
        target.PageSetup.SlideHeight = first_height / EMU_PER_POINT

        insert_index = 0
        for chapter_path in chapter_paths:
            source_path = str(Path(chapter_path).resolve())
            slide_count = len(Presentation(source_path).slides)
            if slide_count == 0:
                continue
            target.Slides.InsertFromFile(source_path, insert_index, 1, slide_count)
            insert_index += slide_count

        target.SaveAs(str(final_path))
        return True
    except Exception as exc:  # noqa: BLE001
        import logging
        logging.getLogger("ppt_engine.merge").exception(
            "PowerPoint COM 合并异常，合并路径=%s，输出=%s",
            chapter_paths, output_path,
        )
        _last_powerpoint_merge_error = str(exc)
        return False
    finally:
        if target is not None:
            try:
                target.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass


def merge_pptx(chapter_paths: list[str | Path], output_path: str | Path) -> Path:
    if not chapter_paths:
        raise ValueError("至少需要一个章节 PPTX 用于合并。")

    _validate_chapter_sizes(chapter_paths)

    merge_engine = os.environ.get("ZHONGCHI_PPT_MERGE_ENGINE", "auto").lower()
    if merge_engine not in {"auto", "powerpoint", "python-pptx"}:
        raise ValueError("ZHONGCHI_PPT_MERGE_ENGINE 只允许 auto/powerpoint/python-pptx。")

    if merge_engine in {"auto", "powerpoint"}:
        merged_by_powerpoint = _merge_pptx_with_powerpoint(chapter_paths, output_path)
        if merged_by_powerpoint:
            return Path(output_path)
        # PowerPoint COM 不可用时，根据模式决定行为
        if merge_engine == "powerpoint":
            raise RuntimeError(
                "PowerPoint COM 合并不可用。当前配置为 ZHONGCHI_PPT_MERGE_ENGINE=powerpoint，"
                "要求使用 Windows + Microsoft PowerPoint + pywin32 进行原生合并。"
                "请确认：1) Windows 系统；2) 已安装 Microsoft PowerPoint；3) 已安装 pywin32（pip install pywin32）。"
                "如需在非 Windows 环境开发调试，请临时改为 python-pptx 模式（仅限测试，不保证合并质量）。"
            )

    # python-pptx 合并路径：仅作为显式配置的开发和测试兜底，不用于正式生成。
    # 正式流程不应走此路径，因为 python-pptx 无法完整保留母版/布局/图片/背景/主题关系，
    # 会导致 PowerPoint 打开时提示修复、M6 出现空白页、图片关系丢失等问题。
    # 注意：即使在此模式下，高保真也仅对简单模板成立；复杂模板（多层母版嵌套、
    # 跨PPT图片引用、MSO_FILL.BACKGROUND 背景继承）极易出现关系断裂或空白页。
    if merge_engine == "python-pptx":
        chapters = [Presentation(str(chapter_path)) for chapter_path in chapter_paths]
        first_chapter = chapters[0]

        merged = Presentation()
        merged.slide_width = first_chapter.slide_width
        merged.slide_height = first_chapter.slide_height

        if len(merged.slides) > 0:
            xml_slides = merged.slides._sldIdLst
            for slide_id in list(xml_slides):
                xml_slides.remove(slide_id)

        for chapter in chapters:
            scale_x = merged.slide_width / chapter.slide_width
            scale_y = merged.slide_height / chapter.slide_height
            for slide in chapter.slides:
                _append_slide(merged, slide, scale_x, scale_y)

        final_path = Path(output_path)
        final_path.parent.mkdir(parents=True, exist_ok=True)
        merged.save(final_path)
        return final_path

    # merge_engine == "auto" 且 PowerPoint COM 不可用时：抛出清晰错误，不静默回退。
    # auto 模式仅在明确希望系统自动选择时才使用，正式生产建议显式指定 powerpoint。
    raise RuntimeError(
        f"ZHONGCHI_PPT_MERGE_ENGINE={merge_engine}，PowerPoint COM 合并失败。"
        "请改为 powerpoint（需 Windows + PowerPoint + pywin32）或 python-pptx（仅测试用，质量无法保证）。"
    )


def render_final_ppt(
    project: dict[str, Any],
    outlines: dict[str, dict[str, Any]],
    output_dir: str | Path,
) -> Path:
    """按 HOME -> M1/M2 -> M3 -> M5 -> M6 -> 尾页顺序合并章节 PPTX。

    Args:
        project: 项目基础信息字典。
        outlines: 章节配置字典，key 为模块标识，value 为模块配置。
                  HOME/M6 被忽略；M1_M2 需要包含 project_type；M3 可包含 parsed_sources；
                  M5 需要包含 case_data。
        output_dir: 输出目录路径。

    Returns:
        生成的最终 PPTX 文件路径。
    """
    output_path = Path(output_dir)
    chapters_dir = output_path / "chapters"
    chapter_paths = []
    for module_id in MERGE_ORDER:
        outline = outlines.get(module_id, {})
        if module_id == "M5":
            case_data = outline.get("case_data")
            if not case_data or not case_data.get("case_id"):
                continue
        if module_id == "TAIL_PRINT" and not project.get("include_print_tail_page", False):
            continue
        chapter_path = render_chapter_ppt(module_id, project, outline, chapters_dir)
        chapter_paths.append(chapter_path)
    final_name = f"{_safe_text(project.get('project_name'), '项目名称')}_中驰智能PPT_Demo.pptx"
    return merge_pptx(chapter_paths, output_path / final_name)


# =============================================================================
# M3 独立测试渲染（不接入正式生产流程）
# =============================================================================


M3_ACTIVE_PLACEHOLDER_SPEC = {
    # field_name: (slide_idx_0based, shape_idx)
    "m3_basic_summary":            (0,  2),
    "m3_quantity_summary":          (4,  2),
    "m3_site_survey_summary":       (10,  1),
    "m3_risk_summary":             (18,  1),
    "m3_solution_summary":          (22,  2),
}

M3_PLACEHOLDER_SPEC = M3_ACTIVE_PLACEHOLDER_SPEC

# 旧模板中可能残留的宁波项目关键词（黑名单，渲染后不应出现在目标 shape 中）
M3_OLD_PROJECT_KEYWORDS = [
    "宁波",
    "宁波2号线",
    "东外环",
    "五里牌",
    "清水浦",
    "虞仁荣",
    "宁波市人民政府",
]

def validate_m3_template_placeholders(pptx_path: Path) -> list[str]:
    """校验 M3 模板包含 active M3 占位符，且每个占位符在其规定位置出现恰好一次。

    校验逻辑：
    1. 每个字段的占位符必须在对应 slide/shape 中存在，且该 shape 的全文必须与占位符完全相等（不允许混入其他内容或其他占位符）。
    2. 每个占位符在整个 deck 中只能出现一次（不允许重复注入到多个 shape）。

    Returns:
        缺失或位置/内容错误的占位符字段名列表。若为空则全部存在且正确。
    """
    prs = Presentation(str(pptx_path))
    PLACEHOLDER = re.compile(r"\{\{(\w+)\}\}")
    missing = []

    # 第1步：逐字段严格检查规定位置的内容是否完全匹配
    for field, (slide_idx, shape_idx) in M3_ACTIVE_PLACEHOLDER_SPEC.items():
        if slide_idx >= len(prs.slides):
            missing.append(field)
            continue
        slide = prs.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            missing.append(field)
            continue
        shape = slide.shapes[shape_idx]
        if not hasattr(shape, "text"):
            missing.append(field)
            continue
        expected = f"{{{{{field}}}}}"
        if shape.text.strip() != expected:
            missing.append(field)

    # 第2步：全局唯一性检查 — 每个占位符只能出现 1 次
    seen: dict[str, int] = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text:
                continue
            found = PLACEHOLDER.findall(shape.text)
            for f in found:
                seen[f] = seen.get(f, 0) + 1
    for field in M3_ACTIVE_FIELD_NAMES:
        count = seen.get(field, 0)
        if count != 1:
            if field not in missing:
                missing.append(field)
    for field in seen:
        if field.startswith("m3_") and field not in M3_ACTIVE_FIELD_NAMES:
            missing.append(field)

    return missing


def validate_m3_template_clean(pptx_path: Path) -> dict[str, list[str]]:
    """校验渲染后的 M3 PPTX：目标 shape 中不包含旧项目关键词。

    Returns:
        {field_name: [残留关键词列表]}。若为空 dict 则全部 clean。
    """
    prs = Presentation(str(pptx_path))
    PLACEHOLDER = re.compile(r"\{\{(\w+)\}\}")
    issues: dict[str, list[str]] = {}

    for field, (slide_idx, shape_idx) in M3_ACTIVE_PLACEHOLDER_SPEC.items():
        if slide_idx >= len(prs.slides) or shape_idx >= len(prs.slides[slide_idx].shapes):
            continue
        shape = prs.slides[slide_idx].shapes[shape_idx]
        if not hasattr(shape, "text"):
            continue
        text = shape.text
        # 跳过仍含占位符的 shape（尚未渲染）
        if PLACEHOLDER.search(text):
            continue
        residual = [kw for kw in M3_OLD_PROJECT_KEYWORDS if kw in text]
        if residual:
            issues[field] = residual

    return issues


def _iter_text_shapes(prs: Presentation):
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text:
                yield slide_idx, shape_idx, slide, shape


def _find_shape_by_exact_text(prs: Presentation, marker: str):
    matches = []
    for slide_idx, shape_idx, slide, shape in _iter_text_shapes(prs):
        if shape.text.strip() == marker:
            matches.append((slide_idx, shape_idx, slide, shape))
    return matches


def validate_m3_full_template_placeholders(pptx_path: Path) -> list[str]:
    """校验 M3 完整测试模板包含 9 个 section 所需文字槽和图片槽。"""
    prs = Presentation(str(pptx_path))
    issues: list[str] = []

    found: dict[str, int] = {}
    for _, _, _, shape in _iter_text_shapes(prs):
        for match in re.findall(r"\{\{([^{}]+)\}\}", shape.text):
            found[match] = found.get(match, 0) + 1

    section_slide_indices: set[int] = set()
    for section in M3_FULL_SECTIONS:
        text_field = section["text_field"]
        image_field = section["image_field"]
        if found.get(text_field, 0) != 1:
            issues.append(text_field)
        if found.get(image_field, 0) != 1:
            issues.append(image_field)
        matches = _find_m3_full_section_slide_indices(prs, text_field, image_field)
        if len(matches) != 1:
            issues.append(f"{text_field}/{image_field}")
        section_slide_indices.update(matches)

    if len(section_slide_indices) != len(M3_FULL_SECTIONS) or len(prs.slides) != len(M3_FULL_SECTIONS):
        issues.append(f"slide_count:{len(prs.slides)}")

    expected = M3_FULL_TEXT_FIELDS | M3_FULL_IMAGE_FIELDS
    for marker in found:
        if marker.startswith("m3_") or marker.startswith("image:m3_"):
            if marker not in expected:
                issues.append(marker)

    return issues


def _find_m3_full_section_slide_indices(prs: Presentation, text_field: str, image_field: str) -> list[int]:
    text_marker = f"{{{{{text_field}}}}}"
    image_marker = f"{{{{{image_field}}}}}"
    indices: list[int] = []
    for slide_index, slide in enumerate(prs.slides):
        slide_text = "\n".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        if text_marker in slide_text and image_marker in slide_text:
            indices.append(slide_index)
    return indices


def _validate_m3_full_image_inputs(images_by_purpose: dict[str, list[bytes]]) -> dict[str, list[bytes]]:
    normalized: dict[str, list[bytes]] = {}
    for purpose, blobs in images_by_purpose.items():
        if purpose not in M3_FULL_IMAGE_FIELDS:
            raise ValueError(f"非法图片用途：{purpose}")
        checked: list[bytes] = []
        for blob in blobs:
            try:
                with Image.open(BytesIO(blob)) as image:
                    image.verify()
            except Exception as exc:
                raise ValueError("图片文件无效或已损坏") from exc
            checked.append(blob)
        normalized[purpose] = checked
    return normalized


def _add_cover_picture(slide, blob: bytes, left, top, width, height):
    """按 cover 策略插入图片：填满占位框，居中裁切，不改变占位框尺寸。"""
    with Image.open(BytesIO(blob)) as image:
        img_width, img_height = image.size
    if img_width <= 0 or img_height <= 0:
        raise ValueError("图片文件无效或已损坏")

    slot_ratio = width / height
    image_ratio = img_width / img_height
    picture = slide.shapes.add_picture(BytesIO(blob), left, top, width=width, height=height)

    if image_ratio > slot_ratio:
        crop = (1 - (slot_ratio / image_ratio)) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < slot_ratio:
        crop = (1 - (image_ratio / slot_ratio)) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def _remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def _find_shape_on_slide_by_exact_text(slide, marker: str):
    matches = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip() == marker:
            matches.append(shape)
    return matches


def replace_m3_full_image_on_slide(slide, image_field: str, blob: bytes) -> None:
    marker = f"{{{{{image_field}}}}}"
    matches = _find_shape_on_slide_by_exact_text(slide, marker)
    if len(matches) != 1:
        raise ValueError(f"M3 完整测试模板缺少唯一图片槽：{marker}")
    shape = matches[0]
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    _remove_shape(shape)
    _add_cover_picture(slide, blob, left, top, width, height)


def replace_m3_full_text_on_slide(slide, text_field: str, value: str) -> None:
    marker = f"{{{{{text_field}}}}}"
    replaced = False
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if marker in run.text:
                    run.text = run.text.replace(marker, value)
                    replaced = True
    if not replaced:
        raise ValueError(f"M3 完整测试模板缺少文字槽：{marker}")


_XLSX_NS = {
    "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "rel": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _column_number(column: str) -> int:
    value = 0
    for char in column:
        value = value * 26 + (ord(char) - ord("A") + 1)
    return value


def _split_cell_ref(ref: str) -> tuple[int, int]:
    match = re.match(r"^([A-Z]+)(\d+)$", ref)
    if not match:
        raise ValueError(f"无效单元格引用：{ref}")
    return int(match.group(2)), _column_number(match.group(1))


def _read_xlsx_shared_strings(archive: zipfile.ZipFile) -> list[str]:
    if "xl/sharedStrings.xml" not in archive.namelist():
        return []
    root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
    strings: list[str] = []
    for item in root.findall("main:si", _XLSX_NS):
        strings.append("".join(node.text or "" for node in item.iter(f"{{{_XLSX_NS['main']}}}t")))
    return strings


def _xlsx_cell_text(cell, shared_strings: list[str]) -> str:
    cell_type = cell.get("t")
    if cell_type == "inlineStr":
        return "".join(node.text or "" for node in cell.iter(f"{{{_XLSX_NS['main']}}}t")).strip()
    value = cell.find("main:v", _XLSX_NS)
    if value is None or value.text is None:
        return ""
    if cell_type == "s":
        try:
            return shared_strings[int(value.text)].strip()
        except (ValueError, IndexError):
            return value.text.strip()
    return value.text.strip()


def _first_xlsx_worksheet_path(archive: zipfile.ZipFile) -> str:
    workbook = ElementTree.fromstring(archive.read("xl/workbook.xml"))
    relationships = ElementTree.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
    rel_targets = {rel.get("Id"): rel.get("Target", "") for rel in relationships}
    first_sheet = workbook.find("main:sheets/main:sheet", _XLSX_NS)
    if first_sheet is None:
        raise ValueError("Excel 表格缺少工作表")
    rel_id = first_sheet.get(f"{{{_XLSX_NS['rel']}}}id")
    target = rel_targets.get(rel_id, "")
    if not target:
        raise ValueError("Excel 表格缺少工作表关系")
    target = target.lstrip("/")
    return target if target.startswith("xl/") else f"xl/{target}"


def _read_xlsx_table(blob: bytes) -> tuple[list[list[str]], list[tuple[int, int, int, int]]]:
    try:
        with zipfile.ZipFile(BytesIO(blob)) as archive:
            shared_strings = _read_xlsx_shared_strings(archive)
            worksheet_path = _first_xlsx_worksheet_path(archive)
            root = ElementTree.fromstring(archive.read(worksheet_path))
    except (KeyError, zipfile.BadZipFile, ElementTree.ParseError) as exc:
        raise ValueError("Excel 表格文件无效或已损坏") from exc

    cells: dict[tuple[int, int], str] = {}
    max_row = 0
    max_col = 0
    for cell in root.findall(".//main:c", _XLSX_NS):
        ref = cell.get("r")
        if not ref:
            continue
        row, col = _split_cell_ref(ref)
        text = _xlsx_cell_text(cell, shared_strings).replace("\xa0", " ")
        cells[(row, col)] = text
        max_row = max(max_row, row)
        max_col = max(max_col, col)

    dimension = root.find("main:dimension", _XLSX_NS)
    if dimension is not None and dimension.get("ref"):
        refs = dimension.get("ref", "").split(":")
        if refs:
            row, col = _split_cell_ref(refs[-1])
            max_row = max(max_row, row)
            max_col = max(max_col, col)

    if max_row <= 0 or max_col <= 0:
        raise ValueError("Excel 表格没有可渲染内容")

    if max_row > M3_FULL_TABLE_MAX_ROWS or max_col > M3_FULL_TABLE_MAX_COLUMNS:
        raise ValueError(
            f"Excel 表格超出可渲染范围：{max_row} 行 x {max_col} 列，"
            f"最大支持 {M3_FULL_TABLE_MAX_ROWS} 行 x {M3_FULL_TABLE_MAX_COLUMNS} 列"
        )

    merges: list[tuple[int, int, int, int]] = []
    merge_cells = root.find("main:mergeCells", _XLSX_NS)
    if merge_cells is not None:
        for merge in merge_cells.findall("main:mergeCell", _XLSX_NS):
            ref = merge.get("ref", "")
            if ":" not in ref:
                continue
            start, end = ref.split(":", 1)
            start_row, start_col = _split_cell_ref(start)
            end_row, end_col = _split_cell_ref(end)
            merges.append((start_row, start_col, end_row, end_col))

    rows = [
        [cells.get((row_index, col_index), "") for col_index in range(1, max_col + 1)]
        for row_index in range(1, max_row + 1)
    ]
    rows, kept_row_indices, kept_col_indices = _trim_empty_table_edges(rows)
    if not rows:
        raise ValueError("Excel 表格没有可渲染内容")
    index_by_row = {original: current for current, original in enumerate(kept_row_indices, start=1)}
    index_by_col = {original: current for current, original in enumerate(kept_col_indices, start=1)}
    trimmed_merges: list[tuple[int, int, int, int]] = []
    for start_row, start_col, end_row, end_col in merges:
        if (
            start_row in index_by_row
            and end_row in index_by_row
            and start_col in index_by_col
            and end_col in index_by_col
        ):
            trimmed_merges.append(
                (
                    index_by_row[start_row],
                    index_by_col[start_col],
                    index_by_row[end_row],
                    index_by_col[end_col],
                )
            )
    return rows, trimmed_merges


def _trim_empty_table_edges(rows: list[list[str]]) -> tuple[list[list[str]], list[int], list[int]]:
    row_indices = [
        index
        for index, row in enumerate(rows, start=1)
        if any(cell.strip() for cell in row)
    ]
    col_indices = [
        index
        for index in range(1, len(rows[0]) + 1)
        if any(row[index - 1].strip() for row in rows)
    ]
    if not row_indices or not col_indices:
        return [], [], []
    trimmed = [
        [rows[row_index - 1][col_index - 1] for col_index in col_indices]
        for row_index in row_indices
    ]
    return trimmed, row_indices, col_indices


def _set_table_cell_borders(cell, color: str = "666666", width_pt: float = 0.75) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for line_tag in ("lnL", "lnR", "lnT", "lnB"):
        existing = tc_pr.find(qn(f"a:{line_tag}"))
        if existing is not None:
            tc_pr.remove(existing)
        line = OxmlElement(f"a:{line_tag}")
        line.set("w", str(int(Pt(width_pt))))
        solid_fill = OxmlElement("a:solidFill")
        srgb_clr = OxmlElement("a:srgbClr")
        srgb_clr.set("val", color)
        solid_fill.append(srgb_clr)
        line.append(solid_fill)
        prst_dash = OxmlElement("a:prstDash")
        prst_dash.set("val", "solid")
        line.append(prst_dash)
        tc_pr.append(line)


def replace_m3_full_table_on_slide(slide, image_field: str, blob: bytes) -> None:
    marker = f"{{{{{image_field}}}}}"
    matches = _find_shape_on_slide_by_exact_text(slide, marker)
    if len(matches) != 1:
        raise ValueError(f"M3 完整测试模板缺少唯一表格槽：{marker}")

    rows, merges = _read_xlsx_table(blob)
    shape = matches[0]
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    _remove_shape(shape)

    row_count = len(rows)
    col_count = len(rows[0])
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = table_shape.table

    if row_count:
        row_height = int(height / row_count)
        for row in table.rows:
            row.height = row_height
    if col_count:
        col_width = int(width / col_count)
        for col in table.columns:
            col.width = col_width

    skipped_cells: set[tuple[int, int]] = set()
    for start_row, start_col, end_row, end_col in merges:
        if start_row > row_count or start_col > col_count or end_row > row_count or end_col > col_count:
            continue
        if start_row == end_row and start_col == end_col:
            continue
        try:
            table.cell(start_row - 1, start_col - 1).merge(table.cell(end_row - 1, end_col - 1))
        except ValueError:
            continue
        for row_index in range(start_row - 1, end_row):
            for col_index in range(start_col - 1, end_col):
                if (row_index, col_index) != (start_row - 1, start_col - 1):
                    skipped_cells.add((row_index, col_index))

    font_size = Pt(8 if row_count > 15 or col_count > 6 else 10)
    for row_index, values in enumerate(rows):
        for col_index, value in enumerate(values):
            if (row_index, col_index) in skipped_cells:
                continue
            cell = table.cell(row_index, col_index)
            _set_table_cell_borders(cell)
            cell.text = value
            cell.margin_left = Pt(2)
            cell.margin_right = Pt(2)
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(18, 117, 181) if row_index == 0 else RGBColor(255, 255, 255)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = font_size
                    run.font.color.rgb = RGBColor(255, 255, 255) if row_index == 0 else RGBColor(0, 0, 0)


def _render_m3_full_ppt(
    project_name: str,
    texts: dict[str, str],
    images_by_purpose: dict[str, list[bytes]],
    output_dir: str | Path,
    dest_filename: str,
    page_texts: dict[str, list[str]] | None = None,
    tables_by_purpose: dict[str, list[bytes]] | None = None,
) -> Path:
    """渲染 M3 九部分文字 + 图片 PPTX，支持多图扩页。"""
    source_template = PPT_TEMPLATE_ROOT / M3_FULL_TEST_TEMPLATE_FILENAME
    if not source_template.exists():
        raise FileNotFoundError(f"M3 完整模板未找到：{source_template}。")

    issues = validate_m3_full_template_placeholders(source_template)
    if issues:
        raise ValueError(f"M3 完整模板占位槽异常：{issues}")

    normalized_images = _validate_m3_full_image_inputs(images_by_purpose)
    normalized_tables = tables_by_purpose or {}
    for purpose in normalized_tables:
        if purpose not in M3_FULL_IMAGE_FIELDS:
            raise ValueError(f"非法表格用途：{purpose}")
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    dest_path = output_path / dest_filename

    source = Presentation(str(source_template))
    target = Presentation()
    target.slide_width = source.slide_width
    target.slide_height = source.slide_height

    generated: list[tuple[int, str, str, str, int | None]] = []
    for section in M3_FULL_SECTIONS:
        image_field = section["image_field"]
        text_field = section["text_field"]
        section_slide_indices = _find_m3_full_section_slide_indices(source, text_field, image_field)
        if not section_slide_indices:
            raise ValueError(f"M3 完整模板缺少 section 模板页：{section['title']}")
        section_template_slide = source.slides[section_slide_indices[0]]
        blobs = normalized_images.get(image_field, [])
        table_blobs = normalized_tables.get(image_field, [])
        if table_blobs or blobs:
            for table_index in range(len(table_blobs)):
                _append_slide(target, section_template_slide)
                generated.append((len(target.slides) - 1, image_field, text_field, "table", table_index))
            for image_index in range(len(blobs)):
                _append_slide(target, section_template_slide)
                generated.append((len(target.slides) - 1, image_field, text_field, "image", image_index))
            continue
        else:
            _append_slide(target, section_template_slide)
            generated.append((len(target.slides) - 1, image_field, text_field, "empty", None))

    target.save(str(dest_path))

    rendered = Presentation(str(dest_path))
    section_titles = {section["text_field"]: section["title"] for section in M3_FULL_SECTIONS}
    for slide_index, image_field, text_field, media_kind, media_index in generated:
        slide = rendered.slides[slide_index]
        if page_texts is not None and media_kind in {"table", "image"} and media_index is not None:
            values = page_texts.get(image_field, [])
            text_value = values[media_index] if media_index < len(values) else ""
        else:
            text_value = _safe_text(texts.get(text_field, ""), section_titles[text_field])
        replace_m3_full_text_on_slide(slide, text_field, text_value)
        if media_kind == "table" and media_index is not None:
            replace_m3_full_table_on_slide(
                slide,
                image_field,
                normalized_tables[image_field][media_index],
            )
        elif media_kind == "image" and media_index is not None:
            replace_m3_full_image_on_slide(
                slide,
                image_field,
                normalized_images[image_field][media_index],
            )
    rendered.save(str(dest_path))
    return dest_path


def render_m3_full_test_ppt(
    project_name: str,
    texts: dict[str, str],
    images_by_purpose: dict[str, list[bytes]],
    output_dir: str | Path,
    page_texts: dict[str, list[str]] | None = None,
    tables_by_purpose: dict[str, list[bytes]] | None = None,
) -> Path:
    """渲染 M3 完整功能测试 PPTX：9 部分文字 + 图片，支持多图扩页。"""
    safe_name = _safe_pptx_filename(project_name)
    return _render_m3_full_ppt(
        project_name,
        texts,
        images_by_purpose,
        output_dir,
        f"M3_完整测试_{safe_name}.pptx",
        page_texts,
        tables_by_purpose,
    )


def render_m3_full_project_ppt(
    project: dict[str, Any],
    texts: dict[str, str],
    images_by_purpose: dict[str, list[bytes]],
    output_dir: str | Path,
    page_texts: dict[str, list[str]] | None = None,
    tables_by_purpose: dict[str, list[bytes]] | None = None,
) -> Path:
    """渲染正式流程 M3 完整资料 PPTX。"""
    return _render_m3_full_ppt(
        _safe_text(project.get("project_name"), "项目"),
        texts,
        images_by_purpose,
        output_dir,
        "M3_项目深化方案.pptx",
        page_texts,
        tables_by_purpose,
    )


def _find_snippet(texts: list[str], keywords: list[str]) -> str:
    """从文本列表中查找包含任一关键词的片段，返回该片段（截断至200字）。

    找不到匹配时返回空字符串。
    """
    for text in texts:
        lower = text.lower()
        if any(kw in lower for kw in keywords):
            return text[:200]
    return ""


def build_m3_replacement_map(
    project: dict[str, Any],
    parsed_sources: list[str] | None = None,
) -> dict[str, str]:
    """构建 M3 模板字段替换映射（第一版：不做 LLM，不接向量库）。

    字段生成规则：
    - m3_basic_summary：使用 project_name/project_location/owner_unit/product_line 拼接
    - 其他字段：从 parsed_sources 中按关键词匹配提取，找不到则用 [待补充：...] 兜底

    Args:
        project: 项目基础信息字典。
        parsed_sources: 可选的文本来源列表（如解析后的文档片段）。

    Returns:
        占位符 -> 替换值的字典，共 9 个字段。
    """
    replacements: dict[str, str] = {}

    # m3_basic_summary：项目基本信息拼接
    parts = [
        project.get("project_name", ""),
        project.get("project_location", ""),
        project.get("owner_unit", ""),
        project.get("product_line", ""),
    ]
    basic = "，".join(p for p in parts if p)
    replacements["m3_basic_summary"] = basic if basic else _missing_placeholder("项目基本情况")

    # 将 parsed_sources 合并为一个文本列表
    texts = parsed_sources or []

    # m3_line_summary：线路图说明
    snippet = _find_snippet(texts, ["线路", "区间", "里程", "站"])
    replacements["m3_line_summary"] = snippet if snippet else _missing_placeholder("项目线路图说明")

    # m3_sensitive_points_summary：敏感点路段说明
    snippet = _find_snippet(texts, ["敏感点", "声屏障", "里程", "措施"])
    replacements["m3_sensitive_points_summary"] = snippet if snippet else _missing_placeholder("敏感点路段说明")

    # m3_quantity_summary：工程量统计说明
    snippet = _find_snippet(texts, ["工程量", "数量", "长度", "m", "㎡", "吨"])
    replacements["m3_quantity_summary"] = snippet if snippet else _missing_placeholder("工程量统计说明")

    # m3_structure_summary：结构形式说明
    snippet = _find_snippet(texts, ["全封闭", "直立", "护栏吸声板", "结构形式"])
    replacements["m3_structure_summary"] = snippet if snippet else _missing_placeholder("结构形式说明")

    # m3_site_survey_summary：现场踏勘说明
    snippet = _find_snippet(texts, ["踏勘", "现场", "吊装", "施工窗口", "轨行区"])
    replacements["m3_site_survey_summary"] = snippet if snippet else _missing_placeholder("现场踏勘说明")

    # m3_investigation_summary：现场勘察情况
    snippet = _find_snippet(texts, ["勘察", "材料上桥", "工作面", "施工条件", "安全防护"])
    replacements["m3_investigation_summary"] = snippet if snippet else _missing_placeholder("现场勘察情况")

    # m3_risk_summary：项目重难点分析
    snippet = _find_snippet(texts, ["风险", "重难点", "工期", "安全", "风压", "维保", "技术实施"])
    replacements["m3_risk_summary"] = snippet if snippet else _missing_placeholder("项目重难点分析")

    # m3_solution_summary：重难点应对措施
    snippet = _find_snippet(texts, ["措施", "方案", "工装", "TEKLA", "定位", "抗风支架", "防松动"])
    replacements["m3_solution_summary"] = snippet if snippet else _missing_placeholder("重难点应对措施")

    return replacements


def _safe_pptx_filename(name: str) -> str:
    """清洗 PPTX 文件名，移除 Windows 非法字符。"""
    illegal = '\\/:*?"<>|'
    for ch in illegal:
        name = name.replace(ch, "_")
    return name or "unnamed"


def render_m3_test_ppt(
    project: dict[str, Any],
    parsed_sources: list[str],
    output_dir: str | Path,
) -> Path:
    """渲染 M3 独立测试 PPTX。

    流程：校验源模板占位符 -> 复制模板 -> 字段替换 -> 返回输出路径。
    不接入正式生产流程，不修改 render_project_ppt() 的行为。

    Args:
        project: 项目基础信息字典。
        parsed_sources: 模拟资料文本列表。
        output_dir: 输出目录路径。

    Returns:
        生成的 M3 测试 PPTX 文件路径。
    """
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    # 生成输出文件名（清洗非法字符）
    project_name = _safe_text(project.get("project_name"), "项目名称")
    safe_name = _safe_pptx_filename(project_name)
    output_filename = f"M3_文字替换测试_{safe_name}.pptx"
    dest_path = output_path / output_filename
    return _render_m3_from_template(project, parsed_sources, dest_path)
