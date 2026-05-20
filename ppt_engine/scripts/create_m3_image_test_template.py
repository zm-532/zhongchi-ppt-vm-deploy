"""创建 M3 图片替换独立测试模板。

只复制并修改图片测试专用模板，不改正式 M3 模板。
"""

from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_DIR = ROOT / "templates" / "solution_fixed_modules"
SOURCE = TEMPLATE_DIR / "M3_项目深化方案模板.pptx"
DEST = TEMPLATE_DIR / "M3_项目深化方案模板_图片替换测试.pptx"

SLOTS = [
    (0, "{{image:project_scope_map}}", 7.0, 1.3, 5.0, 3.0),
    (1, "{{image:project_line_map}}", 7.0, 1.3, 5.0, 3.0),
    (10, "{{image:survey_route_map}}", 7.0, 1.3, 5.0, 3.0),
    (11, "{{image:site_survey_photos:1}}", 0.9, 1.5, 5.4, 2.6),
    (11, "{{image:site_survey_photos:2}}", 6.8, 1.5, 5.4, 2.6),
    (11, "{{image:site_survey_photos:3}}", 0.9, 4.3, 5.4, 2.6),
    (11, "{{image:site_survey_photos:4}}", 6.8, 4.3, 5.4, 2.6),
    (18, "{{image:key_difficulty_evidence}}", 7.0, 1.3, 5.0, 3.0),
]


def _clear_existing_image_markers(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip().startswith("{{image:"):
                element = shape._element
                element.getparent().remove(element)


def main() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(f"源 M3 模板不存在：{SOURCE}")

    shutil.copy2(SOURCE, DEST)
    prs = Presentation(str(DEST))
    _clear_existing_image_markers(prs)

    for slide_idx, marker, left, top, width, height in SLOTS:
        if slide_idx >= len(prs.slides):
            raise ValueError(f"模板页数不足，缺少 slide index {slide_idx}")
        slide = prs.slides[slide_idx]
        shape = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.name = marker
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xF3, 0xF7, 0xFB)
        line = shape.line
        line.color.rgb = RGBColor(0x15, 0x65, 0xC0)
        line.width = Pt(1.25)

        tf = shape.text_frame
        tf.clear()
        paragraph = tf.paragraphs[0]
        paragraph.alignment = 1
        run = paragraph.add_run()
        run.text = marker
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = RGBColor(0x15, 0x65, 0xC0)

    prs.save(str(DEST))
    print(DEST)


if __name__ == "__main__":
    main()
