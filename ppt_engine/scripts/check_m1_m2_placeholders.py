# -*- coding: utf-8 -*-
"""检查首页与 M1/M2 固化模板状态。

检查项：
- 首页模板存在且只包含 {{project_name}}。
- M1/M2 模板不再包含旧的"项目生成信息"页。
"""
from pptx import Presentation
from pathlib import Path
import re

TEMPLATE_ROOT = Path(__file__).resolve().parents[1] / "templates" / "solution_fixed_modules"

M1_M2_TEMPLATES = [
    "公路全封闭声屏障（M1_&_M2）.pptx",
    "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx",
    "铁路_&_轨道交通既有线声屏障_（M1_&_M2）.pptx",
    "铁路声屏障行业背景与技术发展（M1_&_M2）.pptx",
]

HOME_TEMPLATE = "首页.pptx"
PLACEHOLDER_PATTERN = re.compile(r"\{\{(\w+)\}\}")


def _collect_ppt_text(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        texts.append("\n".join(slide_text))
    return texts


def check_home_template() -> dict:
    path = TEMPLATE_ROOT / HOME_TEMPLATE
    prs = Presentation(str(path))
    all_text = "\n".join(_collect_ppt_text(prs))
    placeholders_found = sorted(set(PLACEHOLDER_PATTERN.findall(all_text)))
    return {
        "filename": HOME_TEMPLATE,
        "slide_count": len(prs.slides),
        "placeholders": placeholders_found,
        "is_valid": placeholders_found == ["project_name"],
    }


def check_template(filename: str) -> dict:
    path = TEMPLATE_ROOT / filename
    prs = Presentation(str(path))
    placeholders_found = set()
    project_info_slides = []

    for index, text in enumerate(_collect_ppt_text(prs), 1):
        if "项目生成信息" in text:
            project_info_slides.append(index)
        for m in PLACEHOLDER_PATTERN.finditer(text):
            placeholders_found.add(m.group(1))

    return {
        "filename": filename,
        "slide_count": len(prs.slides),
        "has_placeholders": bool(placeholders_found),
        "placeholders": sorted(placeholders_found),
        "project_info_slides": project_info_slides,
        "is_valid": not project_info_slides,
    }


if __name__ == "__main__":
    home_result = check_home_template()
    print(f"=== {home_result['filename']} ===")
    print(f"  Slides: {home_result['slide_count']}")
    print(f"  Placeholders: {home_result['placeholders']}")
    print(f"  Valid homepage placeholders: {home_result['is_valid']}")
    print()

    for tmpl in M1_M2_TEMPLATES:
        result = check_template(tmpl)
        print(f"=== {result['filename']} ===")
        print(f"  Slides: {result['slide_count']}")
        print(f"  Has placeholders: {result['has_placeholders']}")
        if result["placeholders"]:
            print(f"  Placeholders: {result['placeholders']}")
        print(f"  项目生成信息 slides: {result['project_info_slides']}")
        print(f"  Valid M1/M2 cleanup: {result['is_valid']}")
        print()
