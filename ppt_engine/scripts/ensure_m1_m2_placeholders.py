# -*- coding: utf-8 -*-
"""确保 M1/M2 固化模板包含占位符的辅助脚本。

在每个 M1/M2 模板末尾插入一页"项目生成信息"页，包含所有可替换字段占位符。
可重复运行——若检测到"项目生成信息"页面已存在，则跳过插入。
"""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE_ROOT = Path(__file__).resolve().parents[1] / "templates" / "solution_fixed_modules"

M1_M2_TEMPLATES = [
    "公路全封闭声屏障（M1_&_M2）.pptx",
    "轨道交通地铁全封闭声屏障（M1_&_M2）.pptx",
    "铁路_&_轨道交通既有线声屏障_（M1_&_M2）.pptx",
    "铁路声屏障行业背景与技术发展（M1_&_M2）.pptx",
]

# 颜色
PRIMARY_BLUE = RGBColor(0x15, 0x65, 0xC0)
DEEP_BLUE = RGBColor(0x1A, 0x3A, 0x6B)
LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)
BODY_TEXT = RGBColor(0x1F, 0x29, 0x37)

PLACEHOLDER_FIELDS = [
    ("项目名称", "{{project_name}}"),
    ("项目所在地", "{{project_location}}"),
    ("建设/业主单位", "{{owner_unit}}"),
    ("产品线", "{{product_line}}"),
    ("线路名称", "{{line_name}}"),
    ("现场痛点", "{{site_pain_points}}"),
    ("施工场景", "{{construction_scenario}}"),
]


def _add_textbox(slide, left, top, width, height, text: str, size: int,
                 color: RGBColor = BODY_TEXT, bold: bool = False):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return shape


def _has_project_info_page(prs: Presentation) -> bool:
    """检查演示文稿是否已包含"项目生成信息"页。"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and "项目生成信息" in (shape.text or ""):
                return True
    return False


def insert_project_info_slide(pptx_path: Path) -> bool:
    """向 PPTX 末尾插入"项目生成信息"页。

    Returns:
        True 表示执行了插入操作，False 表示已存在页面（跳过）。
    """
    prs = Presentation(str(pptx_path))

    if _has_project_info_page(prs):
        return False

    # 使用空白布局（某些模板只有 DEFAULT 布局，尝试 layouts[6] 再回退到 layouts[0]）
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    # 背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    # 蓝色顶部导航条
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.333), Inches(0.18),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()

    # 左上角标题
    _add_textbox(
        slide,
        Inches(0.55), Inches(0.35), Inches(7.2), Inches(0.3),
        "中驰智能PPT Demo", 12, DEEP_BLUE, True,
    )

    # 页面标签
    _add_textbox(
        slide,
        Inches(10.4), Inches(0.35), Inches(2.3), Inches(0.3),
        "M1/M2 | 项目生成信息", 11, RGBColor(0xCA, 0x8A, 0x04), True,
    )

    # 页面标题
    _add_textbox(
        slide,
        Inches(0.7), Inches(0.7), Inches(10.0), Inches(0.55),
        "项目生成信息", 26, DEEP_BLUE, True,
    )

    # 分隔线
    accent = slide.shapes.add_shape(
        1, Inches(0.7), Inches(1.35), Inches(2.0), Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY_BLUE
    accent.line.fill.background()

    # 正文字段列表
    body_top = Inches(1.65)
    line_height = Inches(0.55)
    for label, placeholder in PLACEHOLDER_FIELDS:
        _add_textbox(
            slide,
            Inches(0.9), body_top, Inches(10.5), Inches(0.45),
            f"{label}：{placeholder}",
            16, BODY_TEXT, False,
        )
        body_top = body_top + line_height

    prs.save(str(pptx_path))
    return True


if __name__ == "__main__":
    for tmpl in M1_M2_TEMPLATES:
        path = TEMPLATE_ROOT / tmpl
        inserted = insert_project_info_slide(path)
        status = "已插入" if inserted else "已存在，跳过"
        print(f"{tmpl}: {status}")