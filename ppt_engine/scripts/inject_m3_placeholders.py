r"""
修改 M3 模板，向真实章节页面中注入占位符。

幂等设计：
    每次运行前会先清除所有已有的 {{m3_...}} 占位符，
    确保每个占位符在正确位置只出现一次，且不污染标题。

Usage:
    cd D:\中驰股份\code\ppt_engine
    uv run python scripts/inject_m3_placeholders.py
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


TEMPLATE_PATH = Path(__file__).resolve().parents[1] / "templates" / "solution_fixed_modules" / "M3_项目深化方案模板.pptx"

# 占位符注入规范：
# - 字段名 -> (目标页 0index, 目标shape索引或位置描述)
# Shape 0 是页面标题，永远不动
# 所有字段均使用 "replace" 模式：清空目标 shape 后写入占位符，
# 不保留旧模板正文，确保渲染时旧内容不会和新内容共存。
# target_shape_idx 指定目标 shape；"replace" 本身只用于第1页 Shape 2 的简写。
M3_FIELD_PLACEMENTS = {
    "m3_basic_summary":            (0,  2, "第1页旧项目介绍段落（Shape 2）"),
    "m3_quantity_summary":          (4,  2, "第5页工程量说明（Shape 2）"),
    "m3_site_survey_summary":       (10, 1, "第11页现场踏勘（Shape 1）"),
    "m3_risk_summary":             (18,  1, "第19页重难点分析正文（Shape 1）"),
    "m3_solution_summary":          (22,  2, "第23页应对措施（Shape 2）"),
}

PLACEHOLDER_PATTERN = re.compile(r"\{\{(\w+)\}\}")


def clear_all_m3_placeholders(prs: Presentation) -> None:
    """清除模板中所有 {{m3_xxx}} 占位符，恢复原始文本。"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    # 清除含占位符的文本
                    new_text = PLACEHOLDER_PATTERN.sub("", run.text)
                    if new_text != run.text:
                        run.text = new_text


def restore_static_labels(prs: Presentation) -> None:
    """恢复被占位符覆盖过的固定模板标签。"""
    if len(prs.slides) > 18 and len(prs.slides[18].shapes) > 2:
        shape = prs.slides[18].shapes[2]
        if hasattr(shape, "text_frame"):
            tf = shape.text_frame
            tf.clear()
            tf.paragraphs[0].text = "技术实施风险"


def inject_placeholders() -> None:
    """向 M3 模板注入占位符（幂等：每次运行先清除旧占位符再精确注入）。"""
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"M3 模板不存在：{TEMPLATE_PATH}")

    prs = Presentation(str(TEMPLATE_PATH))
    print(f"打开模板：{TEMPLATE_PATH}，共 {len(prs.slides)} 页")

    # 幂等第一步：清除所有旧占位符
    clear_all_m3_placeholders(prs)
    restore_static_labels(prs)
    print("已清除所有旧占位符")

    # 注入新的占位符
    for field_name, (slide_idx, target, note) in M3_FIELD_PLACEMENTS.items():
        if slide_idx >= len(prs.slides):
            print(f"  [跳过] {field_name}：slide {slide_idx} 超出范围")
            continue

        slide = prs.slides[slide_idx]
        placeholder = f"{{{{{field_name}}}}}"
        target_shape_idx = target  # 所有字段都使用 shape 索引直接指定

        if target_shape_idx < len(slide.shapes):
            shape = slide.shapes[target_shape_idx]
            if hasattr(shape, "text_frame"):
                tf = shape.text_frame
                tf.clear()
                run = tf.paragraphs[0].add_run()
                run.text = placeholder
                print(f"  [替换] slide {slide_idx+1} Shape {target_shape_idx} → {placeholder}（{note}）")
            else:
                print(f"  [警告] slide {slide_idx+1} Shape {target_shape_idx} 无 text_frame")
        else:
            print(f"  [警告] slide {slide_idx+1} 没有 Shape {target_shape_idx}，新增文本框")
            new_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.5))
            tf = new_shape.text_frame
            run = tf.paragraphs[0].add_run()
            run.text = placeholder
            print(f"  [新增] slide {slide_idx+1} 新文本框 → {placeholder}（{note}）")

    prs.save(str(TEMPLATE_PATH))
    print(f"\n已保存模板：{TEMPLATE_PATH}")

    # 验证：报告每个占位符出现次数
    prs2 = Presentation(str(TEMPLATE_PATH))
    print("\n注入验证（每个字段应只出现 1 次）：")
    for field_name in M3_FIELD_PLACEMENTS:
        placeholder = f"{{{{{field_name}}}}}"
        count = sum(1 for slide in prs2.slides for shape in slide.shapes
                    if hasattr(shape, "text") and placeholder in shape.text)
        status = "OK" if count == 1 else f"FAIL ({count}次)"
        print(f"  {placeholder}: {status}")


if __name__ == "__main__":
    inject_placeholders()
